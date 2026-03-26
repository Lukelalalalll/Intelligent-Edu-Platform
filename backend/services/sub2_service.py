import os
import re
import json
import base64
import time
import requests
from zhipuai import ZhipuAI
import PyPDF2
from docx import Document
from pptx import Presentation
import tempfile
import fitz
from backend.config import Config
import opendataloader_pdf

def get_proxies():
    """如果在香港调 coze.com 报错，请取消下面 return 的注释"""
    # return {"http": "http://127.0.0.1:7890", "https": "http://127.0.0.1:7890"}
    return None


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in Config.ALLOWED_EXTENSIONS_SUB2


def extract_pdf_pages(pdf_path, page_numbers):
    """裁剪选定的 PDF 页面"""
    with open(pdf_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        pdf_writer = PyPDF2.PdfWriter()
        for page_num in page_numbers:
            if 0 <= page_num < len(pdf_reader.pages):
                pdf_writer.add_page(pdf_reader.pages[page_num])

        fd, temp_path = tempfile.mkstemp(suffix='.pdf')
        os.close(fd)
        with open(temp_path, 'wb') as output_file:
            pdf_writer.write(output_file)
        return temp_path


def _page_numbers_to_spec(page_numbers):
    """Convert 0-based selected pages to opendataloader page spec (1-based)."""
    if not page_numbers:
        return None

    pages = sorted({int(p) + 1 for p in page_numbers if int(p) >= 0})
    if not pages:
        return None

    ranges = []
    start = prev = pages[0]
    for page in pages[1:]:
        if page == prev + 1:
            prev = page
            continue
        ranges.append(f"{start}-{prev}" if start != prev else f"{start}")
        start = prev = page
    ranges.append(f"{start}-{prev}" if start != prev else f"{start}")
    return ",".join(ranges)


def extract_pdf_text_with_loader(pdf_path, page_numbers):
    """Use OpenDataLoader to quickly extract selected PDF pages as markdown text."""
    page_spec = _page_numbers_to_spec(page_numbers)

    with tempfile.TemporaryDirectory(prefix='sub2_odl_') as tmp_dir:
        opendataloader_pdf.convert(
            input_path=pdf_path,
            output_dir=tmp_dir,
            format="markdown",
            quiet=True,
            image_output="off",
            pages=page_spec,
        )

        stem = os.path.splitext(os.path.basename(pdf_path))[0]
        md_candidates = [
            os.path.join(tmp_dir, f"{stem}.md"),
            os.path.join(tmp_dir, f"{stem}_markdown.md"),
        ]
        md_path = next((p for p in md_candidates if os.path.exists(p)), None)

        if not md_path:
            md_files = [f for f in os.listdir(tmp_dir) if f.lower().endswith('.md')]
            if not md_files:
                raise Exception("pdf_loader did not produce markdown output")
            md_path = os.path.join(tmp_dir, md_files[0])

        with open(md_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()


# backend/services/sub2_service.py
import fitz
import base64
import os
import re
import json
from zhipuai import ZhipuAI


# backend/services/sub2_service.py

# backend/services/sub2_service.py

def call_zhipu_ocr(file_path):
    """智谱 OCR：修复版，不再乱动 JSON 结构"""
    client = ZhipuAI(api_key=Config.ZHIPU_API_KEY)

    # 1. PDF/图片处理逻辑 (保持不变)
    img_base = ""
    if file_path.lower().endswith('.pdf'):
        import fitz
        doc = fitz.open(file_path)
        page = doc.load_page(0)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")
        img_base = base64.b64encode(img_data).decode('utf-8')
        doc.close()
    else:
        with open(file_path, "rb") as f:
            img_base = base64.b64encode(f.read()).decode('utf-8')

    prompt = """请识别并提取图中题目。
    要求：
    1. 数学公式用 LaTeX（用 $ 包裹）。
    2. 必须输出标准 JSON：{"exercises": [{"chapter_number":"","sub_chapter_number":"","question_number":"","text":"题目内容","title":""}]}
    3. 特别注意：如果题目包含 Java/Python 代码，请将代码中的反斜杠进行转义。"""

    try:
        response = client.chat.completions.create(
            model="glm-4v",
            messages=[{"role": "user", "content": [{"type": "text", "text": prompt},
                                                   {"type": "image_url", "image_url": {"url": img_base}}]}],
            temperature=0.1
        )

        raw_text = response.choices[0].message.content
        print(f"智谱原始输出: {raw_text}")

        # --- 强力清洗逻辑 (修正版) ---

        # 1. 定位 JSON 块（寻找最外层的 {}）
        match = re.search(r'(\{[\s\S]*\})', raw_text)
        if not match:
            raise Exception("未能识别到 JSON 结构")

        clean_json = match.group(1)

        # 2. 【核心修复】：不要用 replace('\n')。
        # 我们只修复那些非法的反斜杠（LaTeX 或代码里的），但不动换行符
        # json.loads(strict=False) 能够自动处理字符串内部的物理换行！

        # 修复非法的反斜杠：将 \ 变成 \\，但避开已经是合法转义的（如 \n, \", \\）
        clean_json = re.sub(r'\\(?![\\"/bfnrtu])', r'\\\\', clean_json)

        try:
            # 使用 strict=False，它允许字符串中包含原始换行符
            return json.loads(clean_json, strict=False)
        except json.JSONDecodeError as e:
            # 如果还是失败，可能是因为 AI 在引号内输出了物理换行且 strict=False 没扛住
            # 我们针对性地处理：只有在引号内部的物理换行才转义
            print(f"初次解析失败: {e}，尝试深度清洗...")

            # 这种方法比较暴力，将所有不属于 JSON 结构的换行符干掉
            # 我们尝试将所有的物理换行符替换为 \\n，但要排除掉 JSON 结构间的换行
            # 既然难判断，我们就采用一种折中方案：先去掉所有物理换行，再解析
            # 但这对代码题不友好。所以我们最后试一次把控制字符过滤掉：
            processed = "".join(ch for ch in clean_json if ord(ch) >= 32 or ch in '\n\r\t')
            return json.loads(processed, strict=False)

    except Exception as e:
        print(f"提取最终失败: {str(e)}")
        raise e


def call_zhipu_layout_from_text(markdown_text):
    """Use Zhipu text model to format extracted markdown into exercise JSON."""
    if not markdown_text or not markdown_text.strip():
        raise Exception("No extracted markdown text from pdf_loader")

    client = ZhipuAI(api_key=Config.ZHIPU_API_KEY)
    prompt = f"""你是教育内容排版助手。下面是从 PDF 中高精度抽取出的 Markdown 文本。

任务目标：
1) 识别其中的题目内容并按题目切分；
2) 重点是“结构化排版”，不是改写题意；
3) 公式保持 LaTeX（用 $ 包裹）；
4) 输出必须是严格 JSON，且只输出 JSON 本体，不要解释。

严格输出 schema：
{{
  "exercises": [
    {{
      "chapter_number": "",
      "sub_chapter_number": "",
      "question_number": "",
      "page_number": "",
      "title": "",
      "text": ""
    }}
  ]
}}

排版规则：
- chapter_number / sub_chapter_number / question_number 若能从标题编号或题号中推断则填写；否则留空字符串。
- text 字段保留题干、选项、小问、条件、单位等必要信息；保持原始顺序。
- 如果文档中存在目录/章节标题，请用于辅助层级归类，但不要把目录本身当题目。

待处理 Markdown：
{markdown_text}
"""

    response = client.chat.completions.create(
        model="glm-4-plus",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.1,
    )

    raw_text = response.choices[0].message.content
    print(f"智谱排版原始输出: {raw_text}")

    match = re.search(r'(\{[\s\S]*\})', raw_text)
    if not match:
        raise Exception("Zhipu layout output does not contain JSON")

    clean_json = match.group(1)
    clean_json = re.sub(r'\\(?![\\"/bfnrtu])', r'\\\\', clean_json)
    parsed = json.loads(clean_json, strict=False)

    if not isinstance(parsed, dict) or 'exercises' not in parsed:
        raise Exception("Zhipu layout output missing exercises field")
    if not isinstance(parsed['exercises'], list):
        raise Exception("Zhipu layout exercises is not a list")

    return parsed

# backend/services/sub2_service.py

def call_coze_generate(base_content, user_requirements, output_language="Chinese", question_basis=None, knowledge_points="", saved_screenshots=None):
    """调用 Coze 生成新题目 (标准 V3 轮询逻辑)"""
    chat_url = f"{Config.COZE_API_ROOT}/v3/chat"
    headers = {
        "Authorization": f"Bearer {Config.COZE_TOKEN}",
        "Content-Type": "application/json"
    }

    saved_screenshots = saved_screenshots or []
    basis_hint = ""
    if question_basis == "knowledge_points" and knowledge_points.strip():
        basis_hint = f"\n【知识点约束】\n{knowledge_points.strip()}\n请严格围绕这些知识点出题。"
    elif question_basis == "example_images" and saved_screenshots:
        basis_hint = (
            "\n【截图参考】\n"
            f"共提供 {len(saved_screenshots)} 张截图作为题型参考：{', '.join(saved_screenshots[:12])}\n"
            "请基于这些题型风格裂变，不要照抄原题。"
        )

    language_rule = "请使用中文输出全部题干、选项、答案与解析。"
    if str(output_language).strip().lower().startswith("english"):
        language_rule = "Please output the full question set in English, including stem, options, answers, and explanations."

    prompt = f"""你是出题专家。请基于以下原始题目内容进行裂变，生成全新的题目：
    【原始题目内容】：{base_content}
    【生成要求】：{user_requirements}
    {basis_hint}
    【强制要求】：
    1) 包含详细的选项、答案和解析。
    2) 所有数学公式必须使用 LaTeX (用 $ 包裹)。
    3) 严禁与原题文字重复；应保持同知识点但换叙述与数据。
    4) {language_rule}"""

    payload = {
        "bot_id": Config.COZE_BOT_ID,
        "user_id": "sub2_user",
        "stream": False,
        "additional_messages": [{"role": "user", "content": prompt, "content_type": "text"}]
    }

    response = requests.post(chat_url, headers=headers, json=payload, proxies=get_proxies())
    res_data = response.json()

    print(f"Coze 发起对话返回: {res_data}")

    if response.status_code != 200 or 'data' not in res_data:
        raise Exception(f"Coze 发起对话失败: {res_data}")

    # 【核心修复】：同时拿到 chat_id 和 conversation_id
    chat_id = res_data['data']['id']
    conversation_id = res_data['data']['conversation_id']

    # 2. 轮询查询状态 (注意 URL 里拼上了 conversation_id)
    status_url = f"{Config.COZE_API_ROOT}/v3/chat/retrieve?chat_id={chat_id}&conversation_id={conversation_id}"

    for _ in range(40):  # 给它 60 秒的时间（40 * 1.5s）
        status_res = requests.get(status_url, headers=headers, proxies=get_proxies())
        status_data = status_res.json()

        status = status_data.get('data', {}).get('status')
        print(f"Coze 任务状态: {status}")

        if status == 'completed':
            # 3. 状态完成后，获取消息列表 (注意：获取列表是 GET 请求，也要带上 conversation_id)
            msg_url = f"{Config.COZE_API_ROOT}/v3/chat/message/list?chat_id={chat_id}&conversation_id={conversation_id}"
            msg_res = requests.get(msg_url, headers=headers, proxies=get_proxies())
            msg_data = msg_res.json()

            # 在消息列表中寻找 type 为 'answer' 的 AI 回复
            messages = msg_data.get('data', [])
            for msg in messages:
                if msg.get('type') == 'answer':
                    return msg.get('content')
            return "Coze 运行完成，但未能找到答案内容。"

        elif status in ['failed', 'canceled', 'requires_action']:
            raise Exception(f"Coze 任务异常终止: {status}")

        # 还没完成，等 1.5 秒再查
        time.sleep(1.5)

    return "Coze 生成题目超时，请稍后再试"


def create_word_document(content, filename):
    doc = Document()
    doc.add_heading('Generated Practice Questions', 0)
    # 处理生成的文本，按行写入
    for line in content.split('\n'):
        doc.add_paragraph(line)
    output_path = os.path.join(Config['GENERATED_FOLDER_SUB2'], f'{filename}.docx')
    doc.save(output_path)
    return output_path


def create_powerpoint(content, filename):
    prs = Presentation()
    # 简单的按题目切分逻辑
    questions = re.split(r'【Question \d+】|Question \d+:', content)
    for idx, q_text in enumerate(questions):
        if not q_text.strip(): continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title: slide.shapes.title.text = f"Exercise {idx}"
        slide.placeholders[1].text = q_text.strip()

    output_path = os.path.join(Config['GENERATED_FOLDER_SUB2'], f'{filename}.pptx')
    prs.save(output_path)
    return output_path