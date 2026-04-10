from pptx import Presentation
import os
from backend.services.slides.theme_catalog import build_theme_catalog, resolve_base_theme

class PPTTemplateManager:
    def __init__(self, templates_dir='static/ppt_templates'):
        self.templates_dir = templates_dir
        self.templates = {}
        self._load_templates()

    def _load_templates(self):
        """加载所有模板文件"""
        for file in os.listdir(self.templates_dir):
            if file.endswith('.pptx'):
                template_path = os.path.join(self.templates_dir, file)
                self.templates[file] = template_path

    def get_available_themes(self):
        """获取所有可用的主题"""
        base_themes = [os.path.splitext(template_name)[0] for template_name in self.templates.keys()]
        return build_theme_catalog(base_themes)

    def get_placeholders(self, theme_name):
        """获取指定主题的占位符信息"""
        base_themes = [os.path.splitext(template_name)[0] for template_name in self.templates.keys()]
        resolved_theme = resolve_base_theme(theme_name, base_themes)
        template_path = self.templates.get(f"{resolved_theme}.pptx")
        if not template_path:
            raise ValueError(f"Theme: {theme_name} does not exist")

        prs = Presentation(template_path)
        placeholders = []
        
        for slide_layout in prs.slide_masters[0].slide_layouts:
            layout_info = {
                'name': slide_layout.name,
                'placeholders': []
            }
            
            for shape in slide_layout.shapes:
                if shape.is_placeholder:
                    placeholder_info = {
                        'idx': shape.placeholder_format.idx,
                        'name': shape.name,
                        'type': shape.placeholder_format.type
                    }
                    layout_info['placeholders'].append(placeholder_info)
            
            placeholders.append(layout_info)
        
        return placeholders

    def create_presentation(self, theme_name, ppt_schema):
        """根据schema创建演示文稿"""
        base_themes = [os.path.splitext(template_name)[0] for template_name in self.templates.keys()]
        resolved_theme = resolve_base_theme(theme_name, base_themes)
        template_path = self.templates.get(f"{resolved_theme}.pptx")
        if not template_path:
            raise ValueError(f"Theme {theme_name} does not exist")

        prs = Presentation(template_path)
        
        # 设置演示文稿标题
        prs.core_properties.title = ppt_schema['presentation_title']
        
        # 为每个slide创建对应的页面
        for slide_data in ppt_schema['slides']:
            # 选择合适的布局
            slide_layout = self._select_layout(prs, slide_data)
            slide = prs.slides.add_slide(slide_layout)
            
            # 填充内容
            self._fill_slide_content(slide, slide_data)
        
        return prs

    def _select_layout(self, prs, slide_data):
        """根据内容选择合适的布局"""
        # 这里可以根据slide_data的内容选择合适的布局
        # 简单起见，这里使用第一个布局
        return prs.slide_layouts[0]

    def _fill_slide_content(self, slide, slide_data):
        """填充幻灯片内容"""
        for shape in slide.shapes:
            if shape.is_placeholder:
                if shape.placeholder_format.type == 1:  # 标题
                    shape.text = slide_data['title']
                elif shape.placeholder_format.type == 2:  # 正文
                    shape.text = '\n'.join(slide_data['content'])

# 打印占位符信息
def print_placeholders(prs):
    for slide_layout in prs.slide_masters[0].slide_layouts:
        print(f"Layout name: {slide_layout.name}")
        for shape in slide_layout.shapes:
            if shape.is_placeholder:
                idx = shape.placeholder_format.idx
                placeholder_type = shape.placeholder_format.type

                print(
                    f"  Place Holder index: {idx}, Name: {shape.name}, Type: {placeholder_type}")


