"""Dark 模板占位符处理器

Dark 模板占位符类型映射（来自模板实际扫描）：
  type 1  → TITLE         正文标题（所有内容布局、Section、Ending）
  type 2  → BODY          文字内容区域（所有内容布局、Section）
  type 3  → CENTER_TITLE  居中大标题（仅 Title 布局）
  type 4  → SUBTITLE      副标题（仅 Title 布局）
  type 7  → OBJECT        图表对象占位符（内容布局）
  type 13 → SLIDE_NUMBER  幻灯片编号 ← 跳过
  type 15 → FOOTER        页脚       ← 跳过
  type 16 → DATE          日期       ← 跳过
  type 18 → PICTURE       图片占位符（内容布局）← 仅在 chart_type 有值时由图片流程处理

填充规则：
  - 标题页：type 3 填大标题，type 4 填 metadata
  - 内容页：type 1 填幻灯片标题，type 2 填 bullet 列表
  - type 7/18 图片类：由 DarkPPTCreator._process_placeholders 父类处理
  - type 13/15/16：跳过，保持模板默认
"""

from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Pt
from ..output.text_layout_engine import fit_font_size, clean_bullets, shape_dimensions_pt, log_slide_layout_audit

_SKIP_TYPES = {13, 15, 16}


class DarkPlaceholderProcessor:
    """Dark 模板占位符处理器"""

    def __init__(self):
        pass

    # ------------------------------------------------------------------
    # 标题类占位符
    # ------------------------------------------------------------------

    def process_title_placeholders(self, slide, slide_data: dict, title_font_size):
        """填写标题类占位符

        - 标题页（Title 布局）：type 3 = 大标题，type 4 = metadata 副标题
        - 内容页：type 1 = 幻灯片标题
        """
        title = slide_data.get('title', '')
        metadata = slide_data.get('metadata', {})

        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            ph_type = shape.placeholder_format.type

            if ph_type in _SKIP_TYPES:
                continue

            if ph_type == 3:        # CENTER_TITLE（仅 Title 布局）
                shape.text = title
                if title_font_size:
                    self._set_uniform_font(shape, title_font_size)

            elif ph_type == 4:      # SUBTITLE（仅 Title 布局）
                shape.text = self._build_subtitle_text(metadata)

            elif ph_type == 1:      # TITLE（内容页 / Section / Ending）
                shape.text = title
                if title_font_size:
                    self._set_uniform_font(shape, title_font_size)

    # ------------------------------------------------------------------
    # 内容占位符（type 2 BODY）
    # ------------------------------------------------------------------

    def process_content(self, slide, slide_data: dict, content_font_size):
        """将 bullet 列表写入 type 2 BODY 占位符"""
        content_list = clean_bullets(slide_data.get('content', []))
        if not content_list:
            return

        target = self._find_body_shape(slide)
        if target is None:
            return

        self._fill_text(target, content_list, content_font_size, slide_data)

    # ------------------------------------------------------------------
    # 内部工具
    # ------------------------------------------------------------------

    def _find_body_shape(self, slide):
        """找到 type 2 (BODY) 占位符"""
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 2:
                return shape
        return None

    def _fill_text(self, shape, content_list: list, preferred_font_size, slide_data: dict):
        """向占位符写入 bullet 列表，自动缩放字体"""
        tf = shape.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.clear()

        w_pt, h_pt = shape_dimensions_pt(shape)
        preferred_pt = preferred_font_size.pt if hasattr(preferred_font_size, 'pt') else float(preferred_font_size)
        chosen_pt = fit_font_size(content_list, w_pt, h_pt, preferred_pt=preferred_pt)
        final_size = Pt(chosen_pt)

        log_slide_layout_audit(
            slide_idx=slide_data.get('slide_number', '?'),
            title=slide_data.get('title', ''),
            layout_name=getattr(getattr(shape, 'placeholder_format', None), 'type', 'N/A'),
            shape_w_pt=w_pt,
            shape_h_pt=h_pt,
            bullet_count=len(content_list),
            initial_pt=preferred_pt,
            final_pt=chosen_pt,
        )

        for i, text in enumerate(content_list):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            p.font.size = final_size

    @staticmethod
    def _set_uniform_font(shape, font_size):
        for para in shape.text_frame.paragraphs:
            para.font.size = font_size

    @staticmethod
    def _build_subtitle_text(metadata) -> str:
        if not isinstance(metadata, dict):
            return ''
        parts = []
        if metadata.get('author'):
            parts.append(metadata['author'])
        if metadata.get('date'):
            parts.append(metadata['date'])
        if metadata.get('description'):
            parts.append(metadata['description'])
        return '\n'.join(parts)
