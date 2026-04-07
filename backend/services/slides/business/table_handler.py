from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
import pandas as pd

class BusinessTableHandler:
    """Business模板的表格处理器"""
    
    def __init__(self):
        # Business模板的配色方案（更专业的商务风格）
        self.header_bg_color = RGBColor(47, 84, 150)  # #2f5496
        self.even_row_bg_color = RGBColor(242, 242, 242)  # #f2f2f2
        self.odd_row_bg_color = RGBColor(255, 255, 255)   # #ffffff
        self.header_text_color = RGBColor(255, 255, 255)
        self.data_text_color = RGBColor(68, 68, 68)  # #444444
    
    def create_business_table(self, slide, table_data, left, top, width, height):
        """Business模板专属的表格创建逻辑
        
        Args:
            slide: 幻灯片对象
            table_data (dict): 表格数据
            left: 左边距
            top: 上边距
            width: 表格宽度
            height: 表格高度
        """
        rows = len(table_data["rows"]) + 1  # +1 for header
        cols = len(table_data["header"])
        
        # 创建表格
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 设置表头
        for col_idx, header in enumerate(table_data["header"]):
            cell = table.cell(0, col_idx)
            self._set_cell_content_with_linebreaks(cell, str(header))
            
            # 设置表头背景色
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.header_bg_color
            
            # 设置表头样式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(12)  # Business模板使用较小的表头字体
                paragraph.font.color.rgb = self.header_text_color
        
        # 填充数据
        for row_idx, row_data in enumerate(table_data["rows"]):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                self._set_cell_content_with_linebreaks(cell, str(cell_data))
                
                # 设置交替行背景色
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.even_row_bg_color
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.odd_row_bg_color
                
                # 设置单元格样式
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)  # Business模板使用较小的数据字体
                    paragraph.font.color.rgb = self.data_text_color

    def _set_cell_content_with_linebreaks(self, cell, content):
        """设置单元格内容，正确处理<br>标签为换行
        
        Args:
            cell: 表格单元格对象
            content (str): 单元格内容
        """
        # 处理<br>标签
        if '<br>' in content:
            # 分割内容为多行
            lines = content.split('<br>')
            
            # 清空单元格内容
            cell.text_frame.clear()
            
            # 添加第一行
            if lines:
                first_paragraph = cell.text_frame.paragraphs[0]
                first_paragraph.text = lines[0].strip()
                
                # 添加后续行
                for line in lines[1:]:
                    new_paragraph = cell.text_frame.add_paragraph()
                    new_paragraph.text = line.strip()
        else:
            # 没有<br>标签，直接设置文本
            cell.text = content
    
    def read_table_csv(self, table_index, presentation_title):
        """从CSV文件读取表格数据
        
        Args:
            table_index (int): 表格索引
            presentation_title (str): 演示文稿标题
            
        Returns:
            dict: 表格数据，如果文件不存在则返回None
        """
        try:
            # 构建CSV文件路径
            csv_path = f"tables/{presentation_title}/table_{table_index}.csv"
            
            # 读取CSV文件
            df = pd.read_csv(csv_path)
            
            # 转换为表格数据格式
            table_data = {
                "header": df.columns.tolist(),
                "rows": df.values.tolist()
            }
            
            return table_data
        except Exception as e:
            print(f"Warning: Failed to read table CSV file: {e}")
            return None
    
    def process_tables_with_placeholders(self, slide, slide_data, presentation_title, placeholder_processor):
        """通用的表格处理逻辑，用于处理所有Business模板的表格
        
        Args:
            slide: 幻灯片对象
            slide_data (dict): 幻灯片数据
            presentation_title (str): 演示文稿标题
            placeholder_processor: 占位符处理器对象
        """
        if not slide_data.get('tables'):
            return
        
        # 收集其他类型的占位符信息
        other_placeholders = placeholder_processor.collect_other_placeholders(slide)
        
        # 按位置排序占位符（从上到下，从左到右）
        other_placeholders.sort(key=lambda x: (x['top'], x['left']))
        
        for idx, table_info in enumerate(slide_data['tables']):
            table_index = table_info.get('index')
            if table_index is not None:
                # 尝试从CSV文件读取表格数据
                table_data = self.read_table_csv(table_index, presentation_title)
                if table_data:
                    # 使用对应位置的占位符信息
                    if idx < len(other_placeholders):
                        placeholder = other_placeholders[idx]
                        self.create_business_table(slide, table_data, 
                                                  placeholder['left'], 
                                                  placeholder['top'], 
                                                  placeholder['width'], 
                                                  placeholder['height'])
                    else:
                        # 如果没有足够的占位符，使用默认值
                        left = Inches(1)
                        top = Inches(2 + idx * 3.5)  # Business模板使用较小的间距
                        width = Inches(8)
                        height = Inches(3.5)
                        self.create_business_table(slide, table_data, left, top, width, height)
                else:
                    # 如果CSV文件不存在，使用表格数据中的data
                    table_data = table_info.get('data')
                    if table_data:
                        if idx < len(other_placeholders):
                            placeholder = other_placeholders[idx]
                            self.create_business_table(slide, table_data, 
                                                      placeholder['left'], 
                                                      placeholder['top'], 
                                                      placeholder['width'], 
                                                      placeholder['height'])
                        else:
                            left = Inches(1)
                            top = Inches(2 + idx * 3.5)
                            width = Inches(8)
                            height = Inches(3.5)
                            self.create_business_table(slide, table_data, left, top, width, height)
    
    @staticmethod
    def process_tables_generic(slide, slide_data, presentation_title, table_reader_func, table_creator_func):
        """通用的表格处理逻辑，可以被任何PPT创建器使用
        
        Args:
            slide: 幻灯片对象
            slide_data (dict): 幻灯片数据
            presentation_title (str): 演示文稿标题
            table_reader_func: 表格数据读取函数 (table_index, presentation_title) -> table_data
            table_creator_func: 表格创建函数 (slide, table_data, left, top, width, height) -> None
        """
        if not slide_data.get('tables'):
            return
        
        # 收集其他类型的占位符信息（直接在这里实现，避免依赖placeholder_processor）
        other_placeholders = []
        excluded_types = {1, 2, 3, 4}  # 排除标题、内容、页码等常见占位符
        
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
                
            placeholder_type = shape.placeholder_format.type
            
            if placeholder_type not in excluded_types:
                other_placeholders.append({
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'type': placeholder_type,
                    'shape': shape
                })
        
        # 按位置排序占位符（从上到下，从左到右）
        other_placeholders.sort(key=lambda x: (x['top'], x['left']))
        
        for idx, table_info in enumerate(slide_data['tables']):
            table_index = table_info.get('index')
            if table_index is not None:
                # 尝试从CSV文件读取表格数据
                table_data = table_reader_func(table_index, presentation_title)
                if table_data:
                    # 使用对应位置的占位符信息
                    if idx < len(other_placeholders):
                        placeholder = other_placeholders[idx]
                        table_creator_func(slide, table_data, 
                                          placeholder['left'], 
                                          placeholder['top'], 
                                          placeholder['width'], 
                                          placeholder['height'])
                    else:
                        # 如果没有足够的占位符，使用默认值
                        left = Inches(1)
                        top = Inches(2 + idx * 4)  # 默认垂直堆叠
                        width = Inches(8)
                        height = Inches(4)
                        table_creator_func(slide, table_data, left, top, width, height)
                else:
                    # 如果CSV文件不存在，使用表格数据中的data
                    table_data = table_info.get('data')
                    if table_data:
                        if idx < len(other_placeholders):
                            placeholder = other_placeholders[idx]
                            table_creator_func(slide, table_data, 
                                              placeholder['left'], 
                                              placeholder['top'], 
                                              placeholder['width'], 
                                              placeholder['height'])
                        else:
                            left = Inches(1)
                            top = Inches(2 + idx * 4)
                            width = Inches(8)
                            height = Inches(4)
                            table_creator_func(slide, table_data, left, top, width, height) 