"""Business模板页码处理器

专门处理Business模板中页码的处理逻辑，包括：
- 页码位置计算
- 页码样式设置
- Business模板特有的页码格式
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class BusinessSlideNumberHandler:
    """Business模板页码处理器"""
    
    # Business模板页码默认配置
    DEFAULT_LEFT = Inches(9.20)
    DEFAULT_TOP = Inches(6.90)
    DEFAULT_WIDTH = Inches(3)
    DEFAULT_HEIGHT = Inches(0.5)
    DEFAULT_FONT_SIZE = Pt(10)
    
    def __init__(self, 
                 left=None, 
                 top=None, 
                 width=None, 
                 height=None, 
                 font_size=None):
        """初始化页码处理器
        
        Args:
            left: 页码左边距
            top: 页码上边距
            width: 页码宽度
            height: 页码高度
            font_size: 页码字体大小
        """
        self.left = left or self.DEFAULT_LEFT
        self.top = top or self.DEFAULT_TOP
        self.width = width or self.DEFAULT_WIDTH
        self.height = height or self.DEFAULT_HEIGHT
        self.font_size = font_size or self.DEFAULT_FONT_SIZE
    
    def add_slide_number(self, slide, slide_data, is_title_slide=False):
        """添加页码到幻灯片
        
        Args:
            slide: 幻灯片对象
            slide_data (dict): 幻灯片数据
            is_title_slide (bool): 是否为标题页
        """
        # 标题页通常不显示页码
        if is_title_slide:
            return
        
        slide_number = slide_data.get('slide_number')
        if not slide_number:
            return
        
        # 添加幻灯片编号文本框
        slide_number_box = slide.shapes.add_textbox(
            self.left, 
            self.top, 
            self.width, 
            self.height
        )
        
        # 设置文本内容和样式
        text_frame = slide_number_box.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.text = str(slide_number)
        paragraph.font.size = self.font_size
        paragraph.alignment = PP_ALIGN.RIGHT
        
        # Business模板特有的页码样式
        self._apply_business_style(paragraph)
    
    def _apply_business_style(self, paragraph):
        """应用Business模板特有的页码样式
        
        Args:
            paragraph: 段落对象
        """
        # Business模板页码样式：右对齐、小字体
        paragraph.alignment = PP_ALIGN.RIGHT
        
        # 可以在这里添加更多Business模板特有的样式
        # 例如：字体颜色、字体族等
        # paragraph.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)  # 灰色
        # paragraph.font.name = 'Calibri'
    
    def set_position(self, left, top, width=None, height=None):
        """设置页码位置
        
        Args:
            left: 左边距
            top: 上边距
            width: 宽度（可选）
            height: 高度（可选）
        """
        self.left = left
        self.top = top
        if width is not None:
            self.width = width
        if height is not None:
            self.height = height
    
    def set_font_size(self, font_size):
        """设置页码字体大小
        
        Args:
            font_size: 字体大小
        """
        self.font_size = font_size
    
    def create_custom_slide_number(self, slide, slide_number, prefix="", suffix=""):
        """创建自定义格式的页码
        
        Args:
            slide: 幻灯片对象
            slide_number: 页码数字
            prefix (str): 前缀
            suffix (str): 后缀
        """
        slide_number_text = f"{prefix}{slide_number}{suffix}"
        
        slide_number_box = slide.shapes.add_textbox(
            self.left, 
            self.top, 
            self.width, 
            self.height
        )
        
        text_frame = slide_number_box.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.text = slide_number_text
        paragraph.font.size = self.font_size
        
        self._apply_business_style(paragraph) 