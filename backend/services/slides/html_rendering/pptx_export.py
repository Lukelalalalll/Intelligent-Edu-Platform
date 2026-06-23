from __future__ import annotations

import os
import re
from html import unescape
from io import BytesIO


def images_to_pptx(images: list[bytes], output_path: str, *, logger) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for image_bytes in images:
        slide = prs.slides.add_slide(blank_layout)
        image_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        image_stream.close()

    prs.save(output_path)
    logger.info("PPTX saved to %s with %d slides", output_path, len(images))
    return output_path


def extract_plain_text_lines(slide: dict) -> list[str]:
    source = slide.get("raw") or ""
    if not source:
        heading = str(slide.get("heading") or "").strip()
        body = str(slide.get("body") or "").strip()
        bullets = [str(item).strip() for item in (slide.get("bullets") or []) if str(item).strip()]
        lines = ([heading] if heading else []) + ([body] if body else []) + [f"- {item}" for item in bullets]
        if lines:
            return lines
    if source:
        source = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", source)
        source = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", source)
        source = re.sub(r"`{1,3}", "", source)
        source = re.sub(r"^\s{0,3}#{1,6}\s*", "", source, flags=re.MULTILINE)
        source = re.sub(r"^\s*>\s?", "", source, flags=re.MULTILINE)
        source = re.sub(r"^\s*[-*+]\s+", "• ", source, flags=re.MULTILINE)
        source = re.sub(r"^\s*(\d+)\.\s+", r"\1. ", source, flags=re.MULTILINE)
        source = re.sub(r"^\s*[-=_]{3,}\s*$", "", source, flags=re.MULTILINE)
        lines = [unescape(re.sub(r"\s+", " ", line).strip()) for line in source.splitlines()]
        cleaned = [line for line in lines if line and line != "•"]
        if cleaned:
            return cleaned

    html_source = slide.get("body_html") or ""
    html_source = re.sub(r"<br\s*/?>", "\n", html_source, flags=re.IGNORECASE)
    html_source = re.sub(r"</(p|div|li|tr|h[1-6])\s*>", "\n", html_source, flags=re.IGNORECASE)
    html_source = re.sub(r"<li[^>]*>", "• ", html_source, flags=re.IGNORECASE)
    html_source = re.sub(r"<[^>]+>", " ", html_source)
    html_source = re.sub(r"\n{3,}", "\n\n", html_source)
    lines = [unescape(re.sub(r"\s+", " ", line).strip()) for line in html_source.splitlines()]
    return [line for line in lines if line and line != "•"]


def slides_to_basic_pptx(slides: list[dict], output_path: str, deck_title: str, *, logger) -> str:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for index, slide_data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 250, 252)

        title = slide_data.get("heading") or (deck_title if index == 1 else f"Slide {index}")
        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12.0), Inches(0.95))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.NONE
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = title
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.size = Pt(28 if index == 1 else 24)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(15, 23, 42)

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.7), Inches(5.45))
        body_frame = body_box.text_frame
        body_frame.clear()
        body_frame.word_wrap = True
        body_frame.auto_size = MSO_AUTO_SIZE.NONE

        body_lines = extract_plain_text_lines(slide_data)
        if not body_lines:
            body_lines = ["No slide body content was available."]

        body_font_size = 20 if len(body_lines) <= 4 else 18 if len(body_lines) <= 8 else 16
        for paragraph_index, line in enumerate(body_lines):
            paragraph = body_frame.paragraphs[0] if paragraph_index == 0 else body_frame.add_paragraph()
            paragraph.text = line
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.size = Pt(body_font_size)
            paragraph.font.color.rgb = RGBColor(51, 65, 85)
            paragraph.space_after = Pt(8)

    prs.save(output_path)
    logger.info("Fallback PPTX saved to %s with %d slides", output_path, len(slides))
    return output_path


def build_export_result(paths: dict[str, str], *, page_count: int, draft_slides: list[dict] | None = None) -> dict:
    result = {
        "pptx_filename": paths["pptx_filename"],
        "pptx_download_url": paths["pptx_download_url"],
        "html_preview_url": paths["html_preview_url"],
        "pptx_path": paths["pptx_path"],
        "html_path": paths["html_path"],
        "page_count": page_count,
        "render_mode": "screenshot_pptx",
        "renderer": {"available": True, "mode": "browser"},
    }
    if draft_slides is not None:
        result["draft_slides"] = draft_slides
    return result


async def render_and_export_impl(
    *,
    md_content: str,
    css_content: str,
    output_dir: str,
    title: str,
    os_module,
    logger,
    markdown_to_slides_fn,
    slides_to_theme_draft_fn,
    render_html_fn,
    build_output_paths_fn,
    save_single_slide_screenshots_fn,
    images_to_pptx_fn,
) -> dict:
    os_module.makedirs(output_dir, exist_ok=True)
    slides = slides_to_theme_draft_fn(markdown_to_slides_fn(md_content))
    logger.info("Parsed %d slides from markdown", len(slides))

    html = render_html_fn(slides=slides, css_content=css_content, title=title)
    paths = build_output_paths_fn(output_dir, title)
    with open(paths["html_path"], "w", encoding="utf-8") as file_obj:
        file_obj.write(html)
    logger.info("HTML preview saved to %s", paths["html_path"])

    screenshots = await save_single_slide_screenshots_fn(html, len(slides))
    images_to_pptx_fn(screenshots, paths["pptx_path"])
    return build_export_result(paths, page_count=len(slides), draft_slides=slides)


async def export_theme_draft_impl(
    *,
    slides: list[dict],
    css_content: str,
    output_dir: str,
    title: str,
    os_module,
    render_html_fn,
    build_output_paths_fn,
    save_single_slide_screenshots_fn,
    images_to_pptx_fn,
) -> dict:
    os_module.makedirs(output_dir, exist_ok=True)
    html = render_html_fn(slides=slides, css_content=css_content, title=title)
    paths = build_output_paths_fn(output_dir, title)
    with open(paths["html_path"], "w", encoding="utf-8") as file_obj:
        file_obj.write(html)

    screenshots = await save_single_slide_screenshots_fn(html, len(slides))
    images_to_pptx_fn(screenshots, paths["pptx_path"])
    return build_export_result(paths, page_count=len(slides))
