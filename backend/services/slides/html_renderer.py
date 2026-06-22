"""HTML-based slide renderer — converts Markdown + custom CSS into HTML and exports to PPTX.

Uses Playwright to screenshot individual slides (pixel-perfect rendering)
and python-pptx to package them into a downloadable PPTX file.
"""

from __future__ import annotations

import asyncio
import logging
import os
import re
import time
import uuid
from datetime import datetime
from html import unescape
from io import BytesIO
from typing import Optional, Literal

from jinja2 import Environment, FileSystemLoader, select_autoescape
from PIL import Image

logger = logging.getLogger(__name__)

RENDERER_ERROR_CODE = "browser_renderer_unavailable"
_RENDERER_CACHE_TTL_SECONDS = 30.0
_renderer_health_cache: dict[str, tuple[float, dict]] = {}
_PLAYWRIGHT_LAUNCH_ARGS = [
    "--disable-gpu",
    "--disable-dev-shm-usage",
    "--no-sandbox",
    "--disable-setuid-sandbox",
]

# ── Template paths ──
_TEMPLATE_DIR = os.path.join(
    os.path.dirname(__file__), "..", "..", "static", "slides_themes"
)
_SLIDE_TEMPLATE = "slide_template.html"

# ── Default slide dimensions (16:9 HD) ──
SLIDE_WIDTH = 1280
SLIDE_HEIGHT = 720

ThemeDraftLayout = Literal["cover", "content", "split", "quote"]


class BrowserRendererUnavailableError(RuntimeError):
    def __init__(
        self,
        *,
        stage: str,
        summary: str,
        renderer: dict | None = None,
        suggestion: str | None = None,
    ) -> None:
        super().__init__(summary)
        self.stage = stage
        self.summary = summary
        self.renderer = renderer or renderer_status_payload(
            {
                "available": False,
                "mode": "unavailable",
                "message": summary,
                "stage": stage,
            }
        )
        self.suggestion = suggestion or default_renderer_suggestion()

    def to_payload(self) -> dict:
        return build_renderer_error_payload(
            stage=self.stage,
            summary=self.summary,
            renderer=self.renderer,
            suggestion=self.suggestion,
        )


def default_renderer_suggestion() -> str:
    return (
        "Verify that Playwright is installed and Chromium can launch on the backend host, "
        "then retry the export. If needed, run `python -m playwright install chromium`."
    )


def renderer_status_payload(status: dict | None) -> dict:
    payload = {
        "available": bool(status and status.get("available")),
        "mode": "browser" if status and status.get("available") else "unavailable",
    }
    message = status.get("message") if status else None
    if message:
        payload["message"] = str(message)
    return payload


def build_renderer_error_payload(
    *,
    stage: str,
    summary: str,
    renderer: dict | None = None,
    suggestion: str | None = None,
) -> dict:
    return {
        "status": "error",
        "error_code": RENDERER_ERROR_CODE,
        "failed_stage": stage,
        "message": f"Browser rendering failed during {stage}.",
        "details": summary,
        "suggestion": suggestion or default_renderer_suggestion(),
        "renderer": renderer
        or renderer_status_payload(
            {
                "available": False,
                "mode": "unavailable",
                "message": summary,
                "stage": stage,
            }
        ),
    }


def _format_exception(exc: Exception) -> str:
    message = str(exc).strip()
    return message or exc.__class__.__name__


def _renderer_unavailable(stage: str, message: str) -> dict:
    return {
        "available": False,
        "mode": "unavailable",
        "message": message,
        "stage": stage,
    }


def _should_use_threaded_playwright() -> bool:
    if os.name != "nt":
        return False
    try:
        loop = asyncio.get_running_loop()
    except RuntimeError:
        return False
    selector_loop_type = getattr(asyncio, "SelectorEventLoop", None)
    return bool(selector_loop_type and isinstance(loop, selector_loop_type))


def _should_retry_playwright_in_thread(exc: Exception) -> bool:
    return os.name == "nt" and isinstance(exc, NotImplementedError)


def _check_browser_renderer_sync(smoke_test: bool) -> dict:
    try:
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        return _renderer_unavailable(
            "playwright_import",
            f"Playwright import failed: {_format_exception(exc)}",
        )

    try:
        with sync_playwright() as playwright:
            try:
                browser = playwright.chromium.launch(args=_PLAYWRIGHT_LAUNCH_ARGS)
            except Exception as exc:
                return _renderer_unavailable(
                    "chromium_launch",
                    f"Chromium launch failed: {_format_exception(exc)}",
                )

            try:
                if smoke_test:
                    page = browser.new_page(viewport={"width": 320, "height": 180})
                    page.set_content(
                        "<!DOCTYPE html><html><body><div style='width:100px;height:60px;background:#0f766e;color:white'>ok</div></body></html>",
                        wait_until="load",
                        timeout=30000,
                    )
                    page.screenshot(type="png")
            except Exception as exc:
                return _renderer_unavailable(
                    "screenshot_smoke_test",
                    f"Chromium screenshot smoke test failed: {_format_exception(exc)}",
                )
            finally:
                browser.close()
    except Exception as exc:
        return _renderer_unavailable(
            "renderer_check",
            f"Browser renderer health check failed: {_format_exception(exc)}",
        )

    return {
        "available": True,
        "mode": "browser",
        "stage": "ready",
    }


def _renderer_cache_key(smoke_test: bool) -> str:
    return "smoke_test" if smoke_test else "launch"


def _read_renderer_cache(smoke_test: bool) -> dict | None:
    now = time.monotonic()
    keys = [_renderer_cache_key(smoke_test)]
    if not smoke_test:
        keys.append(_renderer_cache_key(True))
    for key in keys:
        cached = _renderer_health_cache.get(key)
        if not cached:
            continue
        cached_at, status = cached
        if now - cached_at <= _RENDERER_CACHE_TTL_SECONDS:
            return dict(status)
    return None


def _write_renderer_cache(smoke_test: bool, status: dict) -> dict:
    stored = dict(status)
    _renderer_health_cache[_renderer_cache_key(smoke_test)] = (time.monotonic(), stored)
    return stored


async def check_browser_renderer(*, smoke_test: bool = True, use_cache: bool = True) -> dict:
    if use_cache:
        cached = _read_renderer_cache(smoke_test)
        if cached is not None:
            return cached

    if _should_use_threaded_playwright():
        logger.info("Using threaded Playwright renderer check because the current event loop cannot spawn subprocesses.")
        status = await asyncio.to_thread(_check_browser_renderer_sync, smoke_test)
        return _write_renderer_cache(smoke_test, status)

    try:
        from playwright.async_api import async_playwright
    except Exception as exc:
        return _write_renderer_cache(
            smoke_test,
            _renderer_unavailable(
                "playwright_import",
                f"Playwright import failed: {_format_exception(exc)}",
            ),
        )

    try:
        async with async_playwright() as playwright:
            try:
                browser = await playwright.chromium.launch(args=_PLAYWRIGHT_LAUNCH_ARGS)
            except Exception as exc:
                return _write_renderer_cache(
                    smoke_test,
                    _renderer_unavailable(
                        "chromium_launch",
                        f"Chromium launch failed: {_format_exception(exc)}",
                    ),
                )

            try:
                if smoke_test:
                    page = await browser.new_page(viewport={"width": 320, "height": 180})
                    await page.set_content(
                        "<!DOCTYPE html><html><body><div style='width:100px;height:60px;background:#0f766e;color:white'>ok</div></body></html>",
                        wait_until="load",
                        timeout=30000,
                    )
                    await page.screenshot(type="png")
            except Exception as exc:
                return _write_renderer_cache(
                    smoke_test,
                    _renderer_unavailable(
                        "screenshot_smoke_test",
                        f"Chromium screenshot smoke test failed: {_format_exception(exc)}",
                    ),
                )
            finally:
                await browser.close()
    except Exception as exc:
        if _should_retry_playwright_in_thread(exc):
            logger.warning("Async Playwright renderer check hit an unsupported subprocess loop; retrying in a worker thread.")
            status = await asyncio.to_thread(_check_browser_renderer_sync, smoke_test)
            return _write_renderer_cache(smoke_test, status)
        return _write_renderer_cache(
            smoke_test,
            _renderer_unavailable(
                "renderer_check",
                f"Browser renderer health check failed: {_format_exception(exc)}",
            ),
        )

    return _write_renderer_cache(
        smoke_test,
        {
            "available": True,
            "mode": "browser",
            "stage": "ready",
        },
    )


async def ensure_browser_renderer(*, smoke_test: bool = True, use_cache: bool = True) -> dict:
    status = await check_browser_renderer(smoke_test=smoke_test, use_cache=use_cache)
    if status.get("available"):
        return renderer_status_payload(status)
    raise BrowserRendererUnavailableError(
        stage=str(status.get("stage") or "browser_renderer"),
        summary=str(status.get("message") or "Browser renderer is unavailable."),
        renderer=renderer_status_payload(status),
    )


def build_theme_draft_preview(
    *,
    slides: list[dict],
    css_content: str,
    title: str = "Presentation",
    selected_slide_id: str | None = None,
    selected_index: int | None = None,
) -> dict:
    renderable_slides = theme_draft_to_render_slides(slides)
    if not renderable_slides:
        renderable_slides = [
            {
                "id": "slide-1",
                "heading": "Presentation Title",
                "body": "",
                "bullets": [],
                "accent_text": "",
                "layout": "cover",
                "align": "left",
            }
        ]

    selected_zero_index = 0
    if selected_slide_id:
        for index, slide in enumerate(renderable_slides):
            if slide.get("id") == selected_slide_id:
                selected_zero_index = index
                break
    elif selected_index is not None:
        selected_zero_index = max(0, min(int(selected_index), len(renderable_slides) - 1))

    html = render_html(slides=renderable_slides, css_content=css_content, title=title)
    return {
        "html": html,
        "selected_index": selected_zero_index,
        "selected_slide_id": renderable_slides[selected_zero_index].get("id") or "",
        "page_count": len(renderable_slides),
    }


def _build_output_paths(output_dir: str, title: str) -> dict[str, str]:
    safe_title = re.sub(r"[^a-zA-Z0-9_\-\u4e00-\u9fff]+", "_", title).strip("_")[:50] or "Presentation"
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    html_filename = f"{safe_title}_{timestamp}.html"
    pptx_filename = f"{safe_title}_{timestamp}.pptx"
    html_path = os.path.join(output_dir, html_filename)
    pptx_path = os.path.join(output_dir, pptx_filename)
    return {
        "html_filename": html_filename,
        "pptx_filename": pptx_filename,
        "html_path": html_path,
        "pptx_path": pptx_path,
        "html_preview_url": f"/api/slides/download_html/{html_filename}",
        "pptx_download_url": f"/api/slides/download_ppt/{pptx_filename}",
    }


# ── Markdown to slides ──

def markdown_to_slides(md_content: str) -> list[dict]:
    """Split Markdown content into an array of slide dictionaries.

    Each slide is a dict with:
        - heading: str
        - body_html: str (HTML rendered from markdown)
        - raw: str

    Splitting rules:
        1. Split on ``===SECTION_BREAK===`` markers first (explicit breaks).
        2. Then split each section on ``# `` and ``## `` headings.
        3. Content before the first heading becomes a title slide.
        4. If content is very short (<1 paragraph) it stays on the title slide.
    """
    if not md_content or not md_content.strip():
        return [{"heading": "Untitled", "body_html": "", "raw": ""}]

    import markdown2

    # Step 1: Split on explicit SECTION_BREAK markers first
    raw_sections = re.split(r"\n?===SECTION_BREAK===", md_content.strip())

    slides: list[dict] = []
    _heading_pattern = re.compile(r"^(#{1,2})\s+(.+)$", re.MULTILINE)

    for raw_section in raw_sections:
        raw_section = raw_section.strip()
        if not raw_section:
            continue

        # Step 2: Further split on # or ## headings at the start of a line
        sub_sections = re.split(r"\n(?=#{1,2}\s)", raw_section)

        for section in sub_sections:
            section = section.strip()
            if not section:
                continue

            # Extract heading (first line starting with # or ##)
            heading = ""
            body = section
            heading_match = _heading_pattern.match(section)
            if heading_match:
                heading = heading_match.group(2).strip()
                # Remove the heading line from body
                body = re.sub(r"^#{1,2}\s+.+\n", "", section, count=1).strip()

            # Convert body markdown to HTML
            body_html = markdown2.markdown(
                body,
                extras=[
                    "fenced-code-blocks",
                    "tables",
                    "break-on-newline",
                    "header-ids",
                    "cuddled-lists",
                ],
            )

            slides.append(
                {
                    "heading": heading,
                    "body_html": body_html,
                    "raw": body,
                }
            )

    # If we have NO slides at all (edge case)
    if not slides:
        slides.append({"heading": "Untitled", "body_html": "<p>No content</p>", "raw": ""})

    # If the first slide has no heading, make it a title slide
    if slides and not slides[0]["heading"]:
        first_line = slides[0]["raw"].split("\n")[0] if slides[0]["raw"] else ""
        slides[0]["heading"] = first_line[:100] or "Presentation Title"
        slides[0]["body_html"] = f'<div class="title-slide-content">{slides[0]["body_html"]}</div>'

    return slides


def _split_body_blocks(raw_body: str) -> tuple[str, list[str]]:
    lines = [line.strip() for line in (raw_body or "").splitlines()]
    bullets: list[str] = []
    paragraphs: list[str] = []
    for line in lines:
        if not line:
            continue
        if re.match(r"^[-*+]\s+", line):
            bullets.append(re.sub(r"^[-*+]\s+", "", line).strip())
            continue
        if re.match(r"^\d+\.\s+", line):
            bullets.append(re.sub(r"^\d+\.\s+", "", line).strip())
            continue
        paragraphs.append(line)
    body = "\n".join(paragraphs).strip()
    return body, bullets


def slides_to_theme_draft(slides: list[dict]) -> list[dict]:
    draft_slides: list[dict] = []
    total = len(slides)
    for index, slide in enumerate(slides):
        body, bullets = _split_body_blocks(slide.get("raw", ""))
        heading = str(slide.get("heading") or "").strip()
        if index == 0:
            layout: ThemeDraftLayout = "cover"
            align = "left"
        elif len(bullets) >= 4:
            layout = "split"
            align = "left"
        elif len(body) < 120 and not bullets:
            layout = "quote"
            align = "center"
        else:
            layout = "content"
            align = "left"
        accent_text = ""
        if index == 0 and total > 1:
            accent_text = f"{total} slides"
        elif bullets:
            accent_text = bullets[0]
        draft_slides.append(
            {
                "id": f"slide-{index + 1}",
                "heading": heading or (f"Slide {index + 1}" if index else "Presentation Title"),
                "body": body,
                "bullets": bullets[:6],
                "accent_text": accent_text,
                "layout": layout,
                "align": align,
            }
        )
    return draft_slides


def theme_draft_to_render_slides(slides: list[dict]) -> list[dict]:
    renderable: list[dict] = []
    for slide in slides:
        heading = str(slide.get("heading") or "").strip()
        body = str(slide.get("body") or "").strip()
        bullets = [str(item).strip() for item in (slide.get("bullets") or []) if str(item).strip()]
        accent_text = str(slide.get("accent_text") or "").strip()
        layout = str(slide.get("layout") or "content").strip().lower()
        align = str(slide.get("align") or "left").strip().lower()
        renderable.append(
            {
                "id": str(slide.get("id") or ""),
                "heading": heading,
                "body": body,
                "bullets": bullets,
                "accent_text": accent_text,
                "layout": layout if layout in {"cover", "content", "split", "quote"} else "content",
                "align": align if align in {"left", "center"} else "left",
            }
        )
    return renderable


# ── HTML rendering via Jinja2 ──

def _build_jinja_env() -> Environment:
    """Create a Jinja2 environment pointing to the slides_themes templates directory."""
    return Environment(
        loader=FileSystemLoader(_TEMPLATE_DIR),
        autoescape=select_autoescape(["html", "xml"]),
    )


def render_html(
    slides: list[dict],
    css_content: str,
    title: str = "Presentation",
) -> str:
    """Render a list of slide dicts into a complete HTML document.

    Args:
        slides: List of slide dicts from ``markdown_to_slides()``.
        css_content: Complete CSS string (including :root variables).
        title: HTML page title.

    Returns:
        Complete HTML string ready for browser rendering.
    """
    env = _build_jinja_env()
    template = env.get_template(_SLIDE_TEMPLATE)

    # Build slide content with richer editable layout hints
    enriched_slides = []
    for i, slide in enumerate(theme_draft_to_render_slides(slides)):
        tag = "h1" if i == 0 or slide["layout"] == "cover" else "h2"
        heading_html = f'<{tag} class="slide-heading">{heading_html_escape(slide["heading"])}</{tag}>' if slide["heading"] else ""

        body_html_parts: list[str] = []
        if slide["body"]:
            body_html_parts.append(
                "".join(f"<p>{heading_html_escape(line)}</p>" for line in slide["body"].split("\n") if line.strip())
            )
        bullets_html = ""
        if slide["bullets"]:
            bullets_html = "<ul>" + "".join(f"<li>{heading_html_escape(item)}</li>" for item in slide["bullets"]) + "</ul>"

        accent_html = ""
        if slide["accent_text"]:
            accent_html = f'<div class="theme-draft-accent">{heading_html_escape(slide["accent_text"])}</div>'

        if slide["layout"] == "split":
            side_html = bullets_html or '<div class="theme-draft-placeholder"></div>'
            content_html = (
                f'<div class="theme-draft split align-{slide["align"]}">'
                f'<div class="theme-draft-main">{heading_html}{accent_html}{"".join(body_html_parts) or "<p></p>"}</div>'
                f'<div class="theme-draft-side">{side_html}</div>'
                f'</div>'
            )
        elif slide["layout"] == "quote":
            quote_text = slide["body"] or "Add your key message here."
            content_html = (
                f'<div class="theme-draft quote align-{slide["align"]}">'
                f'{heading_html}'
                f'<blockquote class="theme-draft-quote">{heading_html_escape(quote_text)}</blockquote>'
                f'{accent_html}'
                f'</div>'
            )
        else:
            content_html = (
                f'<div class="theme-draft {slide["layout"]} align-{slide["align"]}">'
                f'{heading_html}'
                f'{accent_html}'
                f'{"".join(body_html_parts)}'
                f'{bullets_html}'
                f'</div>'
            )

        enriched_slides.append(
            {
                "content_html": content_html,
                "heading": slide["heading"],
                "layout": slide["layout"],
            }
        )

    html = template.render(
        css_content=css_content,
        slides=enriched_slides,
        title=title,
    )

    return html


def heading_html_escape(value: str) -> str:
    return (
        str(value or "")
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&#39;")
    )


# ── Playwright screenshotting ──

def _screenshot_slides_sync(html: str, slide_count: int) -> list[bytes]:
    from playwright.sync_api import sync_playwright

    screenshots: list[bytes] = []

    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(args=_PLAYWRIGHT_LAUNCH_ARGS)
        try:
            page = browser.new_page(
                viewport={"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT}
            )

            page.set_content(html, wait_until="load", timeout=30000)
            time.sleep(0.5)

            for slide_index in range(slide_count):
                page.evaluate(
                    """
                    (currentSlideIndex) => {
                        document.querySelectorAll('.viewport .slide').forEach((s, i) => {
                            if (i === currentSlideIndex) {
                                s.style.display = 'flex';
                                s.scrollIntoView({ behavior: 'instant', block: 'start' });
                            } else {
                                s.style.display = 'none';
                            }
                        });
                        var nav = document.querySelector('.slide-nav');
                        if (nav) nav.style.display = 'none';
                    }
                    """,
                    slide_index,
                )

                time.sleep(0.4)
                screenshots.append(page.screenshot(type="png", full_page=False))

            return screenshots
        finally:
            browser.close()


async def _screenshot_slides(html: str, slide_count: int) -> list[bytes]:
    """Render each slide individually and capture PNG screenshots.

    Returns a list of PNG image bytes (one per slide).
    """
    if _should_use_threaded_playwright():
        logger.info("Using threaded Playwright slide capture because the current event loop cannot spawn subprocesses.")
        return await asyncio.to_thread(_screenshot_slides_sync, html, slide_count)

    from playwright.async_api import async_playwright

    screenshots: list[bytes] = []

    try:
        async with async_playwright() as p:
            browser = await p.chromium.launch(args=_PLAYWRIGHT_LAUNCH_ARGS)
            try:
                page = await browser.new_page(
                    viewport={"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT}
                )

                # Set content once with a reasonable timeout, not networkidle
                # (networkidle hangs on Google Fonts imports)
                await page.set_content(html, wait_until="load", timeout=30000)
                await asyncio.sleep(0.5)

                for slide_index in range(slide_count):
                    # Show only the current slide, hide nav UI
                    await page.evaluate(
                        f"""
                        () => {{
                            document.querySelectorAll('.viewport .slide').forEach((s, i) => {{
                                if (i === {slide_index}) {{
                                    s.style.display = 'flex';
                                    s.scrollIntoView({{ behavior: 'instant', block: 'start' }});
                                }} else {{
                                    s.style.display = 'none';
                                }}
                            }});
                            // Hide navigation bar during screenshot
                            var nav = document.querySelector('.slide-nav');
                            if (nav) nav.style.display = 'none';
                        }}
                        """
                    )

                    await asyncio.sleep(0.4)

                    screenshot = await page.screenshot(type="png", full_page=False)
                    screenshots.append(screenshot)

                return screenshots
            finally:
                await browser.close()
    except Exception as exc:
        if _should_retry_playwright_in_thread(exc):
            logger.warning("Async Playwright slide capture hit an unsupported subprocess loop; retrying in a worker thread.")
            return await asyncio.to_thread(_screenshot_slides_sync, html, slide_count)
        raise


# ── PPTX packaging ──

def _images_to_pptx(images: list[bytes], output_path: str) -> str:
    """Create a PPTX file where each slide is a full-screen image.

    Args:
        images: List of PNG image byte arrays.
        output_path: Where to save the PPTX file.

    Returns:
        The *output_path* (for chaining).
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Set 16:9 slide size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Usually 'Blank'

    for img_bytes in images:
        slide = prs.slides.add_slide(blank_layout)
        img_stream = BytesIO(img_bytes)
        # Add picture filling the entire slide
        slide.shapes.add_picture(
            img_stream,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        img_stream.close()

    prs.save(output_path)
    logger.info("PPTX saved to %s with %d slides", output_path, len(images))
    return output_path


def _extract_plain_text_lines(slide: dict) -> list[str]:
    """Best-effort plain-text extraction for fallback PPTX export."""
    source = slide.get("raw") or ""
    if not source:
        heading = str(slide.get("heading") or "").strip()
        body = str(slide.get("body") or "").strip()
        bullets = [str(item).strip() for item in (slide.get("bullets") or []) if str(item).strip()]
        lines = ([heading] if heading else []) + ([body] if body else []) + [f"- {item}" for item in bullets]
        if lines:
            return lines
    if source:
        source = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", source)
        source = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", source)
        source = re.sub(r"`{1,3}", "", source)
        source = re.sub(r"^\s{0,3}#{1,6}\s*", "", source, flags=re.MULTILINE)
        source = re.sub(r"^\s*>\s?", "", source, flags=re.MULTILINE)
        source = re.sub(r"^\s*[-*+]\s+", "• ", source, flags=re.MULTILINE)
        source = re.sub(r"^\s*(\d+)\.\s+", r"\1. ", source, flags=re.MULTILINE)
        source = re.sub(r"^\s*[-=_]{3,}\s*$", "", source, flags=re.MULTILINE)
        lines = [unescape(re.sub(r"\s+", " ", line).strip()) for line in source.splitlines()]
        cleaned = [line for line in lines if line and line != "•"]
        if cleaned:
            return cleaned

    html_source = slide.get("body_html") or ""
    html_source = re.sub(r"<br\s*/?>", "\n", html_source, flags=re.IGNORECASE)
    html_source = re.sub(r"</(p|div|li|tr|h[1-6])\s*>", "\n", html_source, flags=re.IGNORECASE)
    html_source = re.sub(r"<li[^>]*>", "• ", html_source, flags=re.IGNORECASE)
    html_source = re.sub(r"<[^>]+>", " ", html_source)
    html_source = re.sub(r"\n{3,}", "\n\n", html_source)
    lines = [unescape(re.sub(r"\s+", " ", line).strip()) for line in html_source.splitlines()]
    return [line for line in lines if line and line != "•"]


def _slides_to_basic_pptx(slides: list[dict], output_path: str, deck_title: str) -> str:
    """Fallback PPTX export when browser-based rendering is unavailable."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for index, slide_data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 250, 252)

        title = slide_data.get("heading") or (deck_title if index == 1 else f"Slide {index}")
        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12.0), Inches(0.95))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.NONE
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = title
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.size = Pt(28 if index == 1 else 24)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(15, 23, 42)

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.7), Inches(5.45))
        body_frame = body_box.text_frame
        body_frame.clear()
        body_frame.word_wrap = True
        body_frame.auto_size = MSO_AUTO_SIZE.NONE

        body_lines = _extract_plain_text_lines(slide_data)
        if not body_lines:
            body_lines = ["No slide body content was available."]

        body_font_size = 20 if len(body_lines) <= 4 else 18 if len(body_lines) <= 8 else 16

        for paragraph_index, line in enumerate(body_lines):
            paragraph = body_frame.paragraphs[0] if paragraph_index == 0 else body_frame.add_paragraph()
            paragraph.text = line
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.size = Pt(body_font_size)
            paragraph.font.color.rgb = RGBColor(51, 65, 85)
            paragraph.space_after = Pt(8)

    prs.save(output_path)
    logger.info("Fallback PPTX saved to %s with %d slides", output_path, len(slides))
    return output_path


async def _save_single_slide_screenshots(html: str, slide_count: int) -> list[bytes]:
    """Use Playwright to screenshot individual slides."""
    status = await ensure_browser_renderer(smoke_test=True)
    try:
        return await _screenshot_slides(html, slide_count)
    except Exception as exc:
        logger.error("Playwright screenshot failed: %s", _format_exception(exc))
        raise BrowserRendererUnavailableError(
            stage="slide_screenshot",
            summary=f"Playwright screenshot failed: {_format_exception(exc)}",
            renderer=status,
        ) from exc


# ── Service class wrapper ──

class SlidesHtmlRenderer:
    """Service that renders Markdown + CSS into HTML slides and exports them as PPTX.

    Usage::

        renderer = SlidesHtmlRenderer()
        result = await renderer.render_and_export(md_content, css_content, output_dir, title="My Deck")
    """

    async def render_and_export(
        self,
        md_content: str,
        css_content: str,
        output_dir: str,
        title: str = "Presentation",
    ) -> dict:
        """Full pipeline: Markdown → HTML → Screenshots → PPTX.

        Args:
            md_content: Raw Markdown content for the presentation.
            css_content: Complete CSS (variables + rules) for styling.
            output_dir: Directory where output files will be written.
            title: Presentation title (used as filename prefix).

        Returns:
            dict with keys:
                - pptx_filename: str (just the filename, not full path)
                - pptx_download_url: str (relative URL for download)
                - html_preview_url: str (relative URL for preview)
                - pptx_path: str (full path)
                - html_path: str (full path to preview HTML)
                - page_count: int
        """
        return await render_and_export(md_content, css_content, output_dir, title)


# ── Main render-and-export entry point (module-level) ──

async def render_and_export(
    md_content: str,
    css_content: str,
    output_dir: str,
    title: str = "Presentation",
) -> dict:
    """Full pipeline: Markdown → HTML → Screenshots → PPTX.

    Args:
        md_content: Raw Markdown content for the presentation.
        css_content: Complete CSS (variables + rules) for styling.
        output_dir: Directory where output files will be written.
        title: Presentation title (used as filename prefix).

    Returns:
        dict with keys:
            - pptx_filename: str (just the filename, not full path)
            - pptx_download_url: str (relative URL for download)
            - html_preview_url: str (relative URL for preview)
            - pptx_path: str (full path)
            - html_path: str (full path to preview HTML)
            - page_count: int
    """
    os.makedirs(output_dir, exist_ok=True)

    slides = slides_to_theme_draft(markdown_to_slides(md_content))
    logger.info("Parsed %d slides from markdown", len(slides))

    html = render_html(slides=slides, css_content=css_content, title=title)
    paths = _build_output_paths(output_dir, title)

    with open(paths["html_path"], "w", encoding="utf-8") as file_obj:
        file_obj.write(html)
    logger.info("HTML preview saved to %s", paths["html_path"])

    screenshots = await _save_single_slide_screenshots(html, len(slides))
    _images_to_pptx(screenshots, paths["pptx_path"])

    return {
        "pptx_filename": paths["pptx_filename"],
        "pptx_download_url": paths["pptx_download_url"],
        "html_preview_url": paths["html_preview_url"],
        "pptx_path": paths["pptx_path"],
        "html_path": paths["html_path"],
        "page_count": len(slides),
        "draft_slides": slides,
        "render_mode": "screenshot_pptx",
        "renderer": {"available": True, "mode": "browser"},
    }


async def export_theme_draft(
    *,
    slides: list[dict],
    css_content: str,
    output_dir: str,
    title: str = "Presentation",
) -> dict:
    os.makedirs(output_dir, exist_ok=True)
    html = render_html(slides=slides, css_content=css_content, title=title)
    paths = _build_output_paths(output_dir, title)

    with open(paths["html_path"], "w", encoding="utf-8") as file_obj:
        file_obj.write(html)

    screenshots = await _save_single_slide_screenshots(html, len(slides))
    _images_to_pptx(screenshots, paths["pptx_path"])
    return {
        "pptx_filename": paths["pptx_filename"],
        "pptx_download_url": paths["pptx_download_url"],
        "html_preview_url": paths["html_preview_url"],
        "pptx_path": paths["pptx_path"],
        "html_path": paths["html_path"],
        "page_count": len(slides),
        "render_mode": "screenshot_pptx",
        "renderer": {"available": True, "mode": "browser"},
    }
