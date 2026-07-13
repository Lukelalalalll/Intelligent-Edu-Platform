"""Light 模板占位符处理器

Light 模板占位符类型映射（来自模板实际扫描）：
  type 1  → TITLE          正文标题
  type 2  → BODY           文字内容（Chart layout 系列才有）
  type 3  → CENTER_TITLE   居中大标题（Title/Section title/End Slide）
  type 4  → SUBTITLE       副标题
  type 7  → OBJECT         通用内容（Light 主要的正文占位符）
  type 8  → CHART          图表（1_Chart layout 系列）
  type 13 → SLIDE_NUMBER   幻灯片编号
  type 15 → FOOTER         页脚
  type 16 → DATE           日期

核心处理规则：
  - 内容文字写入 type 2 或 type 7（优先 type 2，其次 type 7）
  - 双栏时：左列写入第一个 type 7，右列写入第二个 type 7
  - type 13/15/16 留空（不填充，避免样式破坏）
  - type 8 图表占位符跳过（由 AI 图表流程单独处理）
"""

from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Pt
from ..output.text_layout_engine import fit_font_size, clean_bullets, shape_dimensions_pt, log_slide_layout_audit

# 不需要处理的辅助占位符类型
_SKIP_TYPES = {8, 13, 15, 16}


class LightPlaceholderProcessor:
    """Light 模板占位符处理器"""

    def __init__(self):
        pass

    # ------------------------------------------------------------------
    # 标题占位符
    # ------------------------------------------------------------------

    def process_title_placeholders(self, slide, slide_data: dict, title_font_size):
        """填写标题类占位符（type 1, 3, 4）"""
        title = slide_data.get('title', '')
        metadata = slide_data.get('metadata', {})

        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            ph_type = shape.placeholder_format.type

            if ph_type == 1:          # TITLE
                shape.text = title
                if title_font_size:
                    self._set_uniform_font(shape, title_font_size)

            elif ph_type == 3:        # CENTER_TITLE
                shape.text = title
                if title_font_size:
                    self._set_uniform_font(shape, title_font_size)

            elif ph_type == 4:        # SUBTITLE
                # 标题页时写 metadata，普通页留空
                text = self._build_subtitle_text(metadata)
                shape.text = text

    # ------------------------------------------------------------------
    # 内容占位符（单栏）
    # ------------------------------------------------------------------

    def process_content_single_column(self, slide, slide_data: dict, content_font_size):
        """将 content_list 写入单个内容占位符（type 2 优先，type 7 次之）"""
        content_list = clean_bullets(slide_data.get('content', []))
        if not content_list:
            return

        target = self._find_primary_content_shape(slide)
        if target is None:
            return

        self._fill_text(target, content_list, content_font_size, slide_data)

    # ------------------------------------------------------------------
    # 内容占位符（双栏）
    # ------------------------------------------------------------------

    def process_content_two_columns(self, slide, slide_data: dict,
                                    left_col: list, right_col: list,
                                    content_font_size):
        """将左右两列内容分别写入两个 type 7 占位符"""
        type7_shapes = self._collect_type7_shapes(slide)

        if len(type7_shapes) >= 2:
            self._fill_text(type7_shapes[0], left_col, content_font_size, slide_data)
            self._fill_text(type7_shapes[1], right_col, content_font_size, slide_data)
        elif len(type7_shapes) == 1:
            # 降级：把全部内容塞进一个占位符
            merged = left_col + right_col
            self._fill_text(type7_shapes[0], merged, content_font_size, slide_data)
        else:
            print("⚠️ [Light] No OBJECT placeholders found for two-column layout")

    # ------------------------------------------------------------------
    # 内部工具
    # ------------------------------------------------------------------

    def _find_primary_content_shape(self, slide):
        """找到第一个可写入文字内容的占位符（type 2 优先，type 7 次之）"""
        type2 = None
        type7 = None
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            ph_type = shape.placeholder_format.type
            if ph_type == 2 and type2 is None:
                type2 = shape
            elif ph_type == 7 and type7 is None:
                type7 = shape
        return type2 or type7

    def _collect_type7_shapes(self, slide) -> list:
        """收集所有 type 7 (OBJECT) 占位符，按左→右顺序排列"""
        shapes = [
            s for s in slide.shapes
            if s.is_placeholder and s.placeholder_format.type == 7
        ]
        shapes.sort(key=lambda s: (s.top, s.left))
        return shapes

    def _fill_text(self, shape, content_list: list, preferred_font_size, slide_data: dict):
        """向 shape 写入 bullet 列表，自动缩放字体"""
        tf = shape.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.clear()

        w_pt, h_pt = shape_dimensions_pt(shape)
        preferred_pt = preferred_font_size.pt if hasattr(preferred_font_size, 'pt') else float(preferred_font_size)
        chosen_pt = fit_font_size(content_list, w_pt, h_pt, preferred_pt=preferred_pt)
        final_size = Pt(chosen_pt)

        log_slide_layout_audit(
            slide_idx=slide_data.get('slide_number', '?'),
            title=slide_data.get('title', ''),
            layout_name=getattr(getattr(shape, 'placeholder_format', None), 'type', 'N/A'),
            shape_w_pt=w_pt,
            shape_h_pt=h_pt,
            bullet_count=len(content_list),
            initial_pt=preferred_pt,
            final_pt=chosen_pt,
        )

        for i, text in enumerate(content_list):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            p.font.size = final_size

    @staticmethod
    def _set_uniform_font(shape, font_size):
        """对 shape 所有段落统一设置字体大小"""
        for para in shape.text_frame.paragraphs:
            para.font.size = font_size

    @staticmethod
    def _build_subtitle_text(metadata) -> str:
        """从 metadata dict 构建副标题字符串"""
        if not isinstance(metadata, dict):
            return ''
        parts = []
        if metadata.get('author'):
            parts.append(metadata['author'])
        if metadata.get('date'):
            parts.append(metadata['date'])
        if metadata.get('description'):
            parts.append(metadata['description'])
        return '\n'.join(parts)
