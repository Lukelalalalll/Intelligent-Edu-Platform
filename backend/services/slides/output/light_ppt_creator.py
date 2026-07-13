"""Light 模板专属 PPT 创建器

继承 PPTCreator，为 Light.pptx 提供：
1. 章节分隔页（Section title 布局）
2. 结尾页（End Slide 布局）
3. 双栏内容（content bullets >= 5 时自动切双栏）
4. 正确处理 DATE / FOOTER / SLIDE_NUMBER 辅助占位符
5. 标题页使用 '标题幻灯片' 布局

Light 模板布局（来自模板实际扫描）：
  '标题幻灯片'              → 标题页  (type 3)
  '标题和内容'              → 单栏内容 (type 1 + type 7)
  'Title and 2 Column Content' / 'Title and 2 content'
                           → 双栏内容 (type 1 + type 7 × 2)
  'Section title'          → 章节分隔页 (type 3 + type 4)
  'End Slide'              → 结尾页 (type 3 + type 4)
  'Chart layout 1/2'       → 有图表的内容页 (type 1 + type 2 + type 7)
"""

import asyncio
import os

from pptx import Presentation

from .ppt_creator import PPTCreator
from .text_layout_engine import clean_bullets, log_slide_layout_audit
from ..light import (
    LightContentProcessor,
    LightSectionHandler,
    LightPlaceholderProcessor,
)

# 双栏布局候选名，按优先级排列
_TWO_COL_LAYOUT_NAMES = [
    'Title and 2 Column Content',
    'Title and 2 content',
    '1_Title and 2 Column Content',
    '1_Title and 2 content',
]

# 单栏布局候选名（标准内容页）
_SINGLE_COL_LAYOUT_NAMES = [
    '标题和内容',
    'Title and content 2',
    '1_Title and content 2',
]

# 结尾页大标题默认文字
_END_SLIDE_TITLE = 'Thank You'


class LightPPTCreator(PPTCreator):
    """Light 模板专属 PPT 创建器

    继承自 PPTCreator，覆盖 create_presentation() 以实现
    章节页、结尾页与双栏布局等 Light-specific 能力。
    """

    def __init__(self, template_base_path=None):
        super().__init__(template_base_path)
        self.template_name = 'Light'
        self.content_processor = LightContentProcessor()
        self.section_handler = LightSectionHandler()
        self.placeholder_processor = LightPlaceholderProcessor()

    # ------------------------------------------------------------------
    # 模板路径
    # ------------------------------------------------------------------

    def _get_template_path(self, theme):
        if theme.lower() == 'light':
            return os.path.join(self.template_base_path, 'Light.pptx')
        return super()._get_template_path(theme)

    # ------------------------------------------------------------------
    # 主入口
    # ------------------------------------------------------------------

    def create_presentation(self, ppt_schema: dict, output_path: str) -> str:
        """创建 Light 主题演示文稿

        完整流程：
          Title → [Section页 → 内容页...] × N → End Slide

        Args:
            ppt_schema:  PPT 结构 dict（含 'slides', 'theme', 'presentation_title'）
            output_path: 输出文件路径

        Returns:
            output_path（写入成功后返回）
        """
        if not ppt_schema or 'slides' not in ppt_schema:
            raise ValueError('Invalid PPT schema')

        theme = ppt_schema.get('theme', 'Light')
        template_path = self._get_template_path(theme)

        if not os.path.exists(template_path):
            raise FileNotFoundError(f'Light template not found: {template_path}')

        prs = Presentation(template_path)
        self._clear_existing_slides(prs)

        presentation_title = ppt_schema.get('presentation_title', '')
        slides_data: list[dict] = ppt_schema['slides']

        # ── 1. 标题页 ──────────────────────────────────────────────────
        self._add_title_slide(prs, presentation_title, ppt_schema.get('metadata', {}))

        # ── 2. 选出主章节，用于插入 Section 分隔页 ────────────────────
        main_headers = self.section_handler.select_main_headers(slides_data)
        print(f'[Light] {len(slides_data)} content slides, '
              f'{len(main_headers)} main sections: {main_headers}')

        # ── 3. 收集模式：延迟批量处理 AI 图片 ─────────────────────────
        print('🔄 [Light Batch] Starting image collection mode...')
        self.start_collecting()

        slide_number = 2  # Title=1，从 2 开始编页
        failed_slides: list[dict] = []

        for slide_index, slide_data in enumerate(slides_data):
            try:
                # ── 3a. 需要时插入章节分隔页 ──────────────────────────
                if slide_data['title'] in self.section_handler.main_headers_with_numbers:
                    subtitle = self.section_handler.get_section_subtitle(
                        slide_data['title'], slides_data
                    )
                    section_slide = self.section_handler.create_section_slide(
                        prs,
                        {'title': slide_data['title'], 'subtitle': subtitle},
                        self._find_layout_by_name,
                    )
                    if section_slide:
                        slide_number += 1
                        print(f'  [Light] Section slide → {slide_data["title"]}')

                # ── 3b. 内容页 ─────────────────────────────────────────
                slide_data['slide_number'] = str(slide_number)
                self._add_content_slide(prs, slide_data, presentation_title)
                slide_number += 1

            except Exception as exc:
                title = slide_data.get('title', f'Slide {slide_index + 1}')
                print(f'❌ [Light] Slide {slide_number} "{title}" failed: {exc}')
                failed_slides.append({'slide_number': slide_number, 'title': title, 'error': str(exc)})

        # ── 4. 结尾页 ──────────────────────────────────────────────────
        end_slide = self.section_handler.create_end_slide(
            prs, self._find_layout_by_name,
            title_text=_END_SLIDE_TITLE,
            subtitle_text=presentation_title,
        )
        if end_slide:
            print('[Light] End Slide added.')

        # ── 5. 批量处理 AI 图片 ────────────────────────────────────────
        self.stop_collecting()
        print('⚡ [Light Batch] Processing collected image placeholders...')
        try:
            asyncio.run(self.process_all_collected_tasks())
        except Exception as exc:
            print(f'⚠️ [Light Batch] Image batch processing failed: {exc}')

        if failed_slides:
            print(f'⚠️ [Light] {len(failed_slides)} slides had errors: '
                  f'{[s["title"] for s in failed_slides]}')

        prs.save(output_path)
        return output_path

    # ------------------------------------------------------------------
    # 标题页
    # ------------------------------------------------------------------

    def _add_title_slide(self, prs, presentation_title: str, metadata: dict):
        """添加 Light 标题页（布局：'标题幻灯片'）"""
        layout = self._find_layout_by_name(prs, '标题幻灯片')
        if layout is None:
            # 通用回退：找 CENTER_TITLE 无 BODY 的布局
            for lay in prs.slide_layouts:
                ph_types = {sh.placeholder_format.type for sh in lay.placeholders}
                if 3 in ph_types and 2 not in ph_types and 7 not in ph_types:
                    layout = lay
                    break
        if layout is None and prs.slide_layouts:
            layout = prs.slide_layouts[0]
        if layout is None:
            return

        title_slide = prs.slides.add_slide(layout)
        title_data = {'title': presentation_title, 'metadata': metadata, 'content': []}
        _, title_font_size = self.content_processor.get_font_sizes([], presentation_title)
        self.placeholder_processor.process_title_placeholders(title_slide, title_data, title_font_size)

    # ------------------------------------------------------------------
    # 内容页
    # ------------------------------------------------------------------

    def _add_content_slide(self, prs, slide_data: dict, presentation_title: str):
        """添加单张内容幻灯片"""
        layout_raw = slide_data.get('layout')
        if isinstance(layout_raw, dict):
            layout_name = (layout_raw.get('name') or '').strip()
        elif layout_raw is None:
            layout_name = ''
        else:
            layout_name = str(layout_raw).strip()

        content_list = clean_bullets(slide_data.get('content', []))
        slide_data = {**slide_data, 'content': content_list}

        use_two_cols = self.content_processor.should_use_two_columns(content_list)

        # ── 布局解析 ──────────────────────────────────────────────────
        layout = self._find_layout_by_name(prs, layout_name) if layout_name else None

        if use_two_cols:
            layout = self._find_two_col_layout(prs) or layout

        if layout is None:
            layout = self._find_single_col_layout(prs)

        if layout is None:
            # 终极回退
            layout = self._find_content_layout(prs)

        if layout is None:
            print(f'⚠️ [Light] No suitable layout for slide "{slide_data.get("title", "")}", skipping')
            return

        # 如 AI 指定了有内容的布局但模板里找不到体，则换去单栏
        if content_list and not self._layout_has_body(layout):
            fallback = self._find_single_col_layout(prs) or self._find_content_layout(prs)
            if fallback:
                print(f'⚠️ [Light] Layout "{layout.name}" has no body — '
                      f'falling back to "{fallback.name}"')
                layout = fallback

        # ── 创建幻灯片 ────────────────────────────────────────────────
        slide = prs.slides.add_slide(layout)

        content_font_size, title_font_size = self.content_processor.get_font_sizes(
            content_list, slide_data.get('title', '')
        )

        # 填标题
        self.placeholder_processor.process_title_placeholders(slide, slide_data, title_font_size)

        # 填内容
        if use_two_cols and len(content_list) >= 2:
            left_col, right_col = self.content_processor.split_two_columns(content_list)
            self.placeholder_processor.process_content_two_columns(
                slide, slide_data, left_col, right_col, content_font_size
            )
        else:
            self.placeholder_processor.process_content_single_column(
                slide, slide_data, content_font_size
            )

        # 处理 AI 图片占位符（进入收集队列）
        # Bug1/2 fix: 直接调用 _collect_visual_tasks，不触碰已写好的文字占位符
        self._collect_visual_tasks(slide, slide_data)

        # 讲者备注
        self._apply_speaker_notes(slide, slide_data)

    # ------------------------------------------------------------------
    # 图片/对象占位符（复用父类 _process_placeholders，仅图片部分）
    # ------------------------------------------------------------------

    def _process_placeholders_light(self, slide, slide_data: dict, presentation_title: str):
        """已弃用：请直接使用 _collect_visual_tasks()。保留仅供兼容性参考。"""
        # Bug 1/2 root-cause was here: fake_data with 'content': [] caused base class
        # to write the slide title into BODY (type 2) placeholders and broke the
        # type-7 chart guard. Replaced by _collect_visual_tasks().
        self._collect_visual_tasks(slide, slide_data)

    # ------------------------------------------------------------------
    # 布局查找辅助
    # ------------------------------------------------------------------

    def _find_two_col_layout(self, prs):
        """按优先级找双栏布局"""
        for name in _TWO_COL_LAYOUT_NAMES:
            layout = self._find_layout_by_name(prs, name)
            if layout:
                return layout
        return None

    def _find_single_col_layout(self, prs):
        """按优先级找单栏内容布局"""
        for name in _SINGLE_COL_LAYOUT_NAMES:
            layout = self._find_layout_by_name(prs, name)
            if layout:
                return layout
        return None
