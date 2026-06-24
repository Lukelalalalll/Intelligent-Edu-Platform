from __future__ import annotations

import asyncio
import os

from pptx import Presentation


def create_presentation(creator, ppt_schema, output_path):
    if not ppt_schema or "slides" not in ppt_schema:
        raise ValueError("Invalid PPT schema")

    theme = ppt_schema.get("theme", "Business")
    template_path = creator._get_template_path(theme)
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Business template not found: {template_path}")

    prs = Presentation(template_path)
    presentation_title = ppt_schema.get("presentation_title", "")
    slides_data = ppt_schema["slides"]

    _create_title_slide(creator, prs, ppt_schema, presentation_title)
    _create_catalogue_slide(creator, prs, slides_data)

    slide_number = 3
    print("🔄 [Batch Processing] Starting batch collection mode for image placeholders...")
    creator.image_processor.start_collecting()

    for slide_index, slide_data in enumerate(slides_data):
        slide_number = _maybe_insert_section_slide(
            creator,
            prs,
            slides_data,
            slide_index,
            slide_data,
            slide_number,
        )
        if "layout" not in slide_data:
            continue

        layout = _resolve_slide_layout(creator, prs, slide_data)
        if not layout:
            continue

        if hasattr(layout, "is_dynamic") and layout.is_dynamic:
            content_list = slide_data.get("content", [])
            target_indices = creator.layout_manager.get_target_indices(
                len(content_list),
                len(layout.group_templates),
            )
            creator.layout_manager.modify_layout_hide_groups(layout, target_indices)

        slide = prs.slides.add_slide(layout)
        slide.custom_layout = layout
        slide_data["slide_number"] = str(slide_number)
        slide_number += 1
        creator._process_business_placeholders(slide, slide_data, presentation_title, prs)
        creator._apply_speaker_notes(slide, slide_data)

    creator.image_processor.stop_collecting()
    print("⚡ [Batch Processing] Executing batch processing for all collected image placeholders...")
    asyncio.run(creator.image_processor.process_all_collected_tasks())
    _append_ending_slide(creator, prs)
    prs.save(output_path)
    return output_path


def _create_title_slide(creator, prs, ppt_schema, presentation_title):
    title_layout = creator._find_layout_by_name(prs, "Title")
    if not title_layout:
        return
    title_slide = prs.slides.add_slide(title_layout)
    creator._process_business_placeholders(
        title_slide,
        {"title": presentation_title, "metadata": ppt_schema.get("metadata", {})},
        presentation_title,
        prs,
        True,
    )


def _create_catalogue_slide(creator, prs, slides_data):
    main_headers = creator.section_handler.select_main_headers(slides_data)
    print("Business Template Processing:")
    print(f"  - Total slides: {len(slides_data)}")
    print(f"  - Selected main headers: {main_headers}")
    print(f"  - All slide titles: {[slide['title'] for slide in slides_data]}")
    catalogue_slide = creator.section_handler.create_catalogue_slide(
        prs,
        main_headers,
        creator._find_layout_by_name,
    )
    if catalogue_slide:
        print("  - Created catalogue slide")
    else:
        print("  - Failed to create catalogue slide")


def _maybe_insert_section_slide(
    creator,
    prs,
    slides_data,
    slide_index,
    slide_data,
    slide_number,
):
    if slide_data["title"] not in creator.section_handler.main_headers_with_numbers:
        return slide_number

    next_main_header = None
    for index in range(slide_index + 1, len(slides_data)):
        if slides_data[index]["title"] in creator.section_handler.main_headers_with_numbers:
            next_main_header = slides_data[index]["title"]
            break

    section_data = creator.section_handler.get_section_content(
        slide_data["title"],
        next_main_header,
        slides_data,
    )
    if not section_data:
        return slide_number

    section_slide = creator.section_handler.create_section_slide(
        prs,
        section_data,
        creator._find_layout_by_name,
    )
    if section_slide:
        print(f"  - Created section slide for: {slide_data['title']}")
        return slide_number + 1

    print(f"  - Failed to create section slide for: {slide_data['title']}")
    return slide_number


def _resolve_slide_layout(creator, prs, slide_data):
    layout_raw = slide_data.get("layout")
    if isinstance(layout_raw, dict):
        layout_name = (layout_raw.get("name") or "").strip()
    elif layout_raw is None:
        layout_name = ""
    else:
        layout_name = str(layout_raw).strip()
    if not layout_name:
        return None

    layout = creator._find_layout_by_name(prs, layout_name)
    if not layout:
        print(f"Warning: Layout '{layout_name}' not found in Business template")
        return None

    content_list = slide_data.get("content", [])
    if content_list and not creator._layout_has_body(layout):
        fallback = creator._find_content_layout(prs)
        if fallback:
            print(
                f"⚠️ Layout '{layout_name}' has no body placeholder "
                f"but slide has {len(content_list)} bullets — "
                f"falling back to '{fallback.name}'"
            )
            layout = fallback
    return layout


def _append_ending_slide(creator, prs):
    ending_layout = creator._find_layout_by_name(prs, "Ending")
    if ending_layout:
        prs.slides.add_slide(ending_layout)
    else:
        print("Warning: Ending layout not found in Business template")
