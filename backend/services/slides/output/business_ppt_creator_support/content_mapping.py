from __future__ import annotations

import asyncio

from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


def process_dynamic_layout_content(
    creator,
    slide,
    slide_data,
    layout,
    title_font_size,
    content_font_size,
    presentation_title=None,
):
    if not hasattr(layout, "is_dynamic") or not layout.is_dynamic:
        return

    title = slide_data.get("title", "")
    content_list = slide_data.get("content", [])
    target_indices = creator.layout_manager.get_target_indices(
        len(content_list),
        len(layout.group_templates),
    )
    asyncio.run(creator.image_processor.process_image_placeholders(slide, slide_data))

    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.type == 1:
            shape.text = title
            if title_font_size:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = title_font_size

    for index, content in enumerate(content_list):
        if index >= len(target_indices):
            continue
        template_index = target_indices[index]
        if template_index >= len(layout.group_templates):
            continue
        template = layout.group_templates[template_index]
        subtitle_content = creator.subtitle_generator.generate_subtitle_content(
            index + 1,
            content,
        )
        _add_subtitle_shape(slide, template, subtitle_content, layout.name)
        _add_body_shape(slide, template, content, content_font_size)

    if presentation_title is None:
        presentation_title = slide_data.get(
            "presentation_title",
            title.split(" ")[0] if title else "Default",
        )

    creator.table_handler.process_tables_with_placeholders(
        slide,
        slide_data,
        presentation_title,
        creator.placeholder_processor,
    )
    creator.latex_processor.process_latex_formulas(
        slide,
        slide_data,
        creator.placeholder_processor,
    )


def _add_subtitle_shape(slide, template, subtitle_content, layout_name):
    if len(template["text_boxes"]) <= 0:
        return
    text_box = template["text_boxes"][0]
    subtitle_shape = slide.shapes.add_textbox(
        text_box["left"],
        text_box["top"],
        text_box["width"],
        text_box["height"],
    )
    subtitle_shape.text = subtitle_content
    if not subtitle_shape.text_frame.paragraphs:
        return

    paragraph = subtitle_shape.text_frame.paragraphs[0]
    paragraph.font.size = Pt(18)
    paragraph.font.bold = True
    if layout_name == "Rectangular Style_dynamic":
        paragraph.alignment = PP_ALIGN.CENTER


def _add_body_shape(slide, template, content, content_font_size):
    if len(template["text_boxes"]) <= 1:
        return
    text_box = template["text_boxes"][1]
    body_shape = slide.shapes.add_textbox(
        text_box["left"],
        text_box["top"],
        text_box["width"],
        text_box["height"],
    )
    body_shape.text = content
    if body_shape.text_frame.paragraphs:
        paragraph = body_shape.text_frame.paragraphs[0]
        paragraph.font.size = content_font_size
        paragraph.font.bold = False
    body_shape.text_frame.word_wrap = True
