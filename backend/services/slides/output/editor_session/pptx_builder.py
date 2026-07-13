from __future__ import annotations

from io import BytesIO
from typing import Any


def build_pptx_from_json(payload: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    slides_data = payload.get("slides", [])
    if not slides_data:
        presentation.slides.add_slide(presentation.slide_layouts[6])
    else:
        for slide_obj in slides_data:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            background = slide.background
            fill = background.fill
            fill.solid()
            try:
                fill.fore_color.rgb = RGBColor.from_string(str(slide_obj.get("background_color", "#FFFFFF")).lstrip("#"))
            except Exception:
                fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            for element in slide_obj.get("elements", []):
                element_type = element.get("type")
                x = Inches(element.get("x", 0))
                y = Inches(element.get("y", 0))
                w = Inches(element.get("w", 1))
                h = Inches(element.get("h", 1))
                if element_type == "textbox":
                    textbox = slide.shapes.add_textbox(x, y, w, h)
                    text_frame = textbox.text_frame
                    text_frame.word_wrap = True
                    text = element.get("text", "")
                    if text:
                        paragraph = text_frame.paragraphs[0]
                        paragraph.text = text
                        paragraph.font.size = Pt(int(element.get("font_size", 24)))
                        paragraph.font.bold = element.get("bold", False)
                        paragraph.font.italic = element.get("italic", False)
                        try:
                            paragraph.font.color.rgb = RGBColor.from_string(
                                str(element.get("font_color", "000000")).lstrip("#")
                            )
                        except Exception:
                            paragraph.font.color.rgb = RGBColor(0, 0, 0)
                        paragraph.alignment = {
                            "left": PP_ALIGN.LEFT,
                            "center": PP_ALIGN.CENTER,
                            "right": PP_ALIGN.RIGHT,
                        }.get(element.get("align", "left"), PP_ALIGN.LEFT)
                elif element_type == "shape" and element.get("shape", "rectangle") == "rectangle":
                    shape = slide.shapes.add_shape(1, x, y, w, h)
                    try:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor.from_string(
                            str(element.get("fill_color", "FFFFFF")).lstrip("#")
                        )
                    except Exception:
                        pass
                elif element_type == "image":
                    src = element.get("src", "")
                    import os

                    if src and os.path.isfile(src):
                        slide.shapes.add_picture(src, x, y, w, h)

    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.read()
