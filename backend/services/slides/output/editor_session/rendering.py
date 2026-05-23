"""Fallback slide rendering via python-pptx + Pillow.

Extracted from EditorSession._render_via_pillow so the core module stays focused.
"""

import os
import sys
from io import BytesIO as Bio
from typing import List


def render_slides_via_pillow(pptx_bytes: bytes) -> List[bytes]:
    """Render slides from in-memory PPTX using Pillow.

    Returns a list of PNG bytes, one per slide.
    """
    try:
        from pptx import Presentation as Prs
    except ImportError:
        raise RuntimeError(
            "Cannot render slides: LibreOffice is not available and "
            "python-pptx is not installed."
        )

    try:
        from PIL import Image, ImageDraw, ImageFont
        _HAS_PIL = True
    except ImportError:
        _HAS_PIL = False

    SLIDE_W_PX = 2000
    SLIDE_H_PX = 1125

    prs = Prs(Bio(pptx_bytes))
    slide_count = len(prs.slides)

    if not _HAS_PIL:
        slide_pngs = []
        for _ in range(slide_count):
            img = Image.new("RGB", (SLIDE_W_PX, SLIDE_H_PX), color=(240, 245, 250))
            draw = ImageDraw.Draw(img)
            draw.text(
                (SLIDE_W_PX // 2 - 200, SLIDE_H_PX // 2 - 12),
                "Preview unavailable (install LibreOffice for full preview)",
                fill=(100, 100, 100),
            )
            buf = Bio()
            img.save(buf, format="PNG")
            slide_pngs.append(buf.getvalue())
        return slide_pngs

    title_font = None
    text_font = None
    try:
        if sys.platform == "darwin":
            font_paths = [
                "/System/Library/Fonts/Helvetica.ttc",
                "/System/Library/Fonts/PingFang.ttc",
                "/Library/Fonts/Arial.ttf",
            ]
        else:
            font_paths = [
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
            ]
        for fp in font_paths:
            if os.path.isfile(fp):
                title_font = ImageFont.truetype(fp, 36)
                text_font = ImageFont.truetype(fp, 24)
                break
    except Exception:
        pass

    slide_pngs = []
    for slide_idx, slide in enumerate(prs.slides):
        img = Image.new("RGB", (SLIDE_W_PX, SLIDE_H_PX), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        slide_num_text = f"Slide {slide_idx + 1}"
        try:
            draw.text(
                (SLIDE_W_PX - 180, 20),
                slide_num_text,
                fill=(180, 180, 180),
                font=text_font,
            )
        except Exception:
            draw.text((SLIDE_W_PX - 180, 20), slide_num_text, fill=(180, 180, 180))

        y_offset = 80
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    left = int(shape.left / 914400 * 72 * (150 / 72)) if shape.left else 60
                    top = int(shape.top / 914400 * 72 * (150 / 72)) if shape.top else y_offset
                    left = max(10, min(left, SLIDE_W_PX - 400))
                    top = max(10, min(top, SLIDE_H_PX - 40))
                    is_title = shape == slide.shapes[0] if slide.shapes else False
                    font = title_font if is_title else text_font
                    fill_color = (30, 30, 30) if is_title else (80, 80, 80)
                    display_text = text[:120] + ("..." if len(text) > 120 else "")
                    try:
                        draw.text((left, top), display_text, fill=fill_color, font=font)
                    except Exception:
                        draw.text((left, top), display_text, fill=fill_color)
                    y_offset = top + 40 if not is_title else y_offset
            elif shape.has_table:
                left = int(shape.left / 914400 * 72 * (150 / 72)) if shape.left else 100
                top = int(shape.top / 914400 * 72 * (150 / 72)) if shape.top else y_offset
                w = int(shape.width / 914400 * 72 * (150 / 72)) if shape.width else 1600
                h = int(shape.height / 914400 * 72 * (150 / 72)) if shape.height else 300
                draw.rectangle(
                    [left, top, left + w, top + h],
                    outline=(200, 200, 200),
                    width=2,
                )
                try:
                    draw.text(
                        (left + 10, top + 10),
                        "[Table]",
                        fill=(150, 150, 150),
                        font=text_font,
                    )
                except Exception:
                    draw.text((left + 10, top + 10), "[Table]", fill=(150, 150, 150))

        if y_offset <= 80 and len(slide.shapes) == 0:
            try:
                draw.text(
                    (SLIDE_W_PX // 2 - 200, SLIDE_H_PX // 2 - 12),
                    "(Empty slide)",
                    fill=(180, 180, 180),
                    font=title_font,
                )
            except Exception:
                draw.text(
                    (SLIDE_W_PX // 2 - 200, SLIDE_H_PX // 2 - 12),
                    "(Empty slide)",
                    fill=(180, 180, 180),
                )

        buf = Bio()
        img.save(buf, format="PNG")
        slide_pngs.append(buf.getvalue())

    return slide_pngs
