"""Editor session management — create PPTX, render per-slide PNGs, extract element tree."""
from __future__ import annotations

import io
import json
import logging
import os
import shutil
import subprocess
import tempfile
import time
import uuid
import zipfile
import hashlib
from typing import Any

from pptx import Presentation
from pptx.util import Emu

from backend.config import Config

logger = logging.getLogger(__name__)

SESSION_DIR = os.path.join(Config.BASE_DIR, "generated", "editor_sessions")
SOFFICE_BIN = shutil.which("soffice") or shutil.which("libreoffice")

# ── Helpers ───────────────────────────────────────────────────────────────────


def _dedup_zip(src_bytes: bytes) -> bytes:
    """Remove duplicate ZIP entries produced by python-pptx."""
    seen: set[str] = set()
    out = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(src_bytes), "r") as zin:
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename not in seen:
                    seen.add(item.filename)
                    zout.writestr(item, zin.read(item.filename))
    return out.getvalue()


def _dedup_zip_file(path: str) -> None:
    """Rewrite a PPTX file in-place with duplicate ZIP entries removed."""
    with open(path, "rb") as f:
        original = f.read()
    deduped = _dedup_zip(original)
    if deduped != original:
        with open(path, "wb") as f:
            f.write(deduped)


def _emu_to_pt(emu_value: int) -> float:
    return round(Emu(emu_value).pt, 2)


def _get_align(text_frame) -> str:
    """Return CSS-compatible alignment string from a text frame."""
    from pptx.enum.text import PP_ALIGN
    mapping = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
    }
    if text_frame.paragraphs:
        return mapping.get(text_frame.paragraphs[0].alignment, "left")
    return "left"


def _extract_font_color(run) -> str | None:
    """Extract CSS hex color from a python-pptx Run, resolving theme/scheme colors."""
    try:
        color = run.font.color
        if color.rgb:
            return f"#{color.rgb}"
        if color.theme_color is not None:
            # Resolve theme color via the XML tree
            rPr = run._r.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
            if rPr is None:
                rPr = run._r.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            # If we can get the RGB from the element directly
            if color.rgb:
                return f"#{color.rgb}"
    except Exception:
        pass
    return None


# ── Element extraction ────────────────────────────────────────────────────────


def _extract_elements(slide, slide_idx: int) -> list[dict[str, Any]]:
    """Extract editable element tree from a python-pptx Slide."""
    elements: list[dict[str, Any]] = []
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        ph = shape.placeholder_format
        elem: dict[str, Any] = {
            "id": f"s{slide_idx}_ph{ph.idx}",
            "placeholder_idx": ph.idx,
            "bbox": {
                "x": _emu_to_pt(shape.left),
                "y": _emu_to_pt(shape.top),
                "w": _emu_to_pt(shape.width),
                "h": _emu_to_pt(shape.height),
            },
            "editable": True,
        }
        # Determine type: picture placeholder (idx >= 10) or text
        if hasattr(shape, "image"):
            elem["type"] = "image"
            elem["content"] = None
        elif hasattr(shape, "text_frame") and shape.text_frame:
            elem["type"] = "text"
            tf = shape.text_frame
            elem["content"] = shape.text or ""
            # Font info from first run
            if tf.paragraphs and tf.paragraphs[0].runs:
                run = tf.paragraphs[0].runs[0]
                elem["font_size"] = round(run.font.size.pt, 1) if run.font.size else None
                elem["bold"] = bool(run.font.bold)
                elem["font_color"] = _extract_font_color(run)
            else:
                elem["font_size"] = None
                elem["bold"] = False
                elem["font_color"] = None
            elem["align"] = _get_align(tf)
        else:
            elem["type"] = "text"
            elem["content"] = ""
            elem["font_size"] = None
            elem["bold"] = False
            elem["font_color"] = None
            elem["align"] = "left"
        elements.append(elem)
    return elements


# ── LibreOffice rendering ─────────────────────────────────────────────────────


def _render_pptx_to_pngs(pptx_path: str, out_dir: str) -> list[str]:
    """Convert a PPTX file into per-slide PNGs using LibreOffice.

    Returns a sorted list of PNG file paths (one per slide).
    """
    if not SOFFICE_BIN:
        raise RuntimeError("LibreOffice (soffice) not found on PATH.")

    # Prefer a stable two-step pipeline:
    # 1) PPTX -> PDF via LibreOffice
    # 2) PDF pages -> PNG via PyMuPDF (fitz)
    # Some LibreOffice builds export only the first PPT slide when converting
    # directly to PNG; PDF conversion is consistently multi-page.
    t0 = time.time()
    result = subprocess.run(
        [SOFFICE_BIN, "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
        capture_output=True, text=True, timeout=120,
    )
    elapsed = time.time() - t0
    logger.info("[editor-session] soffice PDF rendered in %.1fs (rc=%d)", elapsed, result.returncode)

    import glob
    pdfs = sorted(glob.glob(os.path.join(out_dir, "*.pdf")))
    if pdfs:
        try:
            import fitz

            pdf_path = pdfs[0]
            doc = fitz.open(pdf_path)
            pngs: list[str] = []
            for i in range(doc.page_count):
                page = doc.load_page(i)
                # Render at ~2x then downstream resize keeps dimensions stable.
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                png_path = os.path.join(out_dir, f"page_{i}.png")
                pix.save(png_path)
                pngs.append(png_path)
            doc.close()
            if pngs:
                return pngs
        except Exception as e:
            logger.warning("[editor-session] PDF->PNG fallback failed: %s", e)

    # Fallback: direct PPTX->PNG (kept for compatibility)
    t1 = time.time()
    result_png = subprocess.run(
        [SOFFICE_BIN, "--headless", "--convert-to", "png", "--outdir", out_dir, pptx_path],
        capture_output=True, text=True, timeout=120,
    )
    logger.info(
        "[editor-session] soffice direct PNG rendered in %.1fs (rc=%d)",
        time.time() - t1,
        result_png.returncode,
    )

    pngs = sorted(glob.glob(os.path.join(out_dir, "*.png")))
    if not pngs:
        stderr = result.stderr or result_png.stderr
        raise RuntimeError(f"LibreOffice produced no PNGs. stderr: {stderr[:500]}")
    return pngs


def _resize_png(src: str, dst: str, size: tuple[int, int] = (1280, 720)) -> None:
    """Resize a PNG to the target size using Pillow."""
    try:
        from PIL import Image
        img = Image.open(src).convert("RGB")
        img = img.resize(size, Image.LANCZOS)
        img.save(dst, "PNG", optimize=True)
    except ImportError:
        shutil.copy2(src, dst)


# ── Session CRUD ──────────────────────────────────────────────────────────────


def _session_path(session_id: str) -> str:
    return os.path.join(SESSION_DIR, session_id)


def _file_md5(path: str) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _themes_match(session_theme: str | None, request_theme: str) -> bool:
    """Return True if both themes resolve to the same base template.

    Handles None session_theme (always matches), exact string equality,
    and variant names that share a common base (e.g. "Classic Formal" → "Classic").
    """
    if session_theme is None:
        return True
    if (session_theme or "").strip().lower() == (request_theme or "").strip().lower():
        return True
    # Resolve both to their base template name and compare.
    try:
        from .theme_catalog import resolve_base_theme

        available = [
            os.path.splitext(name)[0]
            for name in os.listdir(Config.PPT_TEMPLATES_FOLDER)
            if name.endswith(".pptx")
        ]
        base_session = resolve_base_theme((session_theme or "").strip(), available)
        base_request = resolve_base_theme((request_theme or "").strip(), available)
        return base_session == base_request
    except Exception:
        return False


def create_session(theme: str, ppt_schema: dict) -> dict[str, Any]:
    """Create an editor session: generate PPTX → render PNGs → extract elements.

    Returns the full session payload for the frontend.
    """
    from .ppt_creator import PPTCreator

    session_id = uuid.uuid4().hex[:12]
    sess_dir = _session_path(session_id)
    slides_dir = os.path.join(sess_dir, "slides")
    os.makedirs(slides_dir, exist_ok=True)

    # 1. Generate PPTX using the exact same backend pipeline as normal export.
    # This keeps all themes (except specialized ones like Business) consistent
    # with the proven dark-theme placeholder filling behavior.
    # Inject 'theme' so create_presentation uses the correct template even when
    # ppt_schema came from localStorage without a theme field.
    creator = PPTCreator(template_base_path=Config.PPT_TEMPLATES_FOLDER)
    pptx_path = os.path.join(sess_dir, "original.pptx")
    creator.create_presentation({**ppt_schema, 'theme': theme}, pptx_path)
    _dedup_zip_file(pptx_path)

    # 2. Render all slides to PNG
    rendered_count = 0
    with tempfile.TemporaryDirectory() as tmp:
        raw_pngs = _render_pptx_to_pngs(pptx_path, tmp)
        rendered_count = len(raw_pngs)
        for i, raw_png in enumerate(raw_pngs):
            _resize_png(raw_png, os.path.join(slides_dir, f"{i}.png"))

    # 3. Extract element tree from rendered PPTX
    prs_rendered = Presentation(pptx_path)
    slide_width_pt = _emu_to_pt(prs_rendered.slide_width)
    slide_height_pt = _emu_to_pt(prs_rendered.slide_height)

    total_slides = len(prs_rendered.slides)
    payload_count = min(total_slides, rendered_count)
    if payload_count != total_slides:
        logger.warning(
            "[editor-session] rendered slide count mismatch: ppt=%d png=%d; using %d previews",
            total_slides,
            rendered_count,
            payload_count,
        )

    slides_payload: list[dict] = []
    for idx in range(payload_count):
        slide = prs_rendered.slides[idx]
        slides_payload.append({
            "index": idx,
            "preview_url": f"/api/slides/editor-img/{session_id}/{idx}.png",
            "elements": _extract_elements(slide, idx),
        })

    # 4. Write meta
    meta = {
        "session_id": session_id,
        "theme": theme,
        "slide_width_pt": slide_width_pt,
        "slide_height_pt": slide_height_pt,
        "slides": slides_payload,
    }
    with open(os.path.join(sess_dir, "meta.json"), "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

    return meta


def get_session_meta(session_id: str) -> dict[str, Any]:
    """Load session metadata from disk."""
    meta_path = os.path.join(_session_path(session_id), "meta.json")
    if not os.path.exists(meta_path):
        raise FileNotFoundError(f"Session not found: {session_id}")
    with open(meta_path, encoding="utf-8") as f:
        return json.load(f)


def get_slide_png_path(session_id: str, page: int) -> str:
    """Return the filesystem path of a session slide PNG."""
    png_path = os.path.join(_session_path(session_id), "slides", f"{page}.png")
    if not os.path.exists(png_path):
        raise FileNotFoundError(f"Slide PNG not found: {session_id}/{page}")
    return png_path


def export_pptx(
    session_id: str,
    theme: str,
    ppt_schema: dict,
    edits: list[dict],
    slide_images: list[dict] | None = None,
) -> str:
    """Apply edits + inserted images to PPTX and return the file path for download."""
    sess_dir = _session_path(session_id)
    original_path = os.path.join(sess_dir, "original.pptx")
    if not os.path.exists(original_path):
        raise FileNotFoundError(f"Session not found: {session_id}")

    # Read session metadata to detect theme mismatch.
    meta_path = os.path.join(sess_dir, "meta.json")
    session_theme: str | None = None
    if os.path.exists(meta_path):
        try:
            with open(meta_path, encoding="utf-8") as f:
                meta = json.load(f)
            session_theme = meta.get("theme")
            if session_theme and session_theme != theme:
                logger.warning(
                    "[editor-session] export theme mismatch: session_id=%s session_theme=%s req_theme=%s",
                    session_id,
                    session_theme,
                    theme,
                )
        except Exception:
            pass

    # When user didn't modify anything in the editor, always return original.pptx
    # used by preview rendering. This guarantees download == preview and avoids
    # accidental regeneration caused by theme variant/name mismatch.
    if not edits and not slide_images:
        try:
            logger.info(
                "[editor-session] export no-edits: returning original.pptx | session_id=%s req_theme=%s session_theme=%s size=%d md5=%s slides=%d",
                session_id,
                theme,
                session_theme,
                os.path.getsize(original_path),
                _file_md5(original_path),
                len(Presentation(original_path).slides),
            )
        except Exception:
            pass
        return original_path

    prs = Presentation(original_path)

    # Build edit lookup: { "s0_ph1": { "content": "...", "image_asset_id": "..." } }
    edit_map: dict[str, dict] = {}
    for edit in edits:
        edit_map[edit["element_id"]] = edit

    slides_list = list(prs.slides)
    has_effective_changes = False

    # Apply element edits
    for slide_idx, slide in enumerate(slides_list):
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            ph = shape.placeholder_format
            elem_id = f"s{slide_idx}_ph{ph.idx}"
            edit = edit_map.get(elem_id)
            if not edit:
                continue
            if "content" in edit and hasattr(shape, "text_frame"):
                # Preserve original placeholder formatting when content is unchanged.
                new_text = (edit.get("content") or "")
                old_text = (shape.text or "")
                if new_text != old_text:
                    shape.text = new_text
                    has_effective_changes = True
            if "image_asset_id" in edit:
                asset_path = os.path.join(SESSION_DIR, "assets", f"{edit['image_asset_id']}.png")
                if os.path.exists(asset_path):
                    slide.shapes.add_picture(
                        asset_path, shape.left, shape.top, shape.width, shape.height,
                    )
                    has_effective_changes = True

    # Insert per-slide images (user-uploaded, centered on each slide)
    if slide_images:
        _insert_slide_images(prs, slides_list, slide_images)
        has_effective_changes = True

    # If edits are all no-op and no images were inserted, return original file
    # to keep binary output aligned with preview source.
    if not has_effective_changes:
        try:
            logger.info(
                "[editor-session] export no-op edits: returning original.pptx | session_id=%s req_theme=%s edits=%d",
                session_id,
                theme,
                len(edits or []),
            )
        except Exception:
            pass
        return original_path

    # Save final PPTX
    out_path = os.path.join(sess_dir, "exported.pptx")
    prs.save(out_path)
    _dedup_zip_file(out_path)

    try:
        logger.info(
            "[editor-session] export rebuilt: session_id=%s req_theme=%s file=%s size=%d md5=%s slides=%d edits=%d images=%d",
            session_id,
            theme,
            os.path.basename(out_path),
            os.path.getsize(out_path),
            _file_md5(out_path),
            len(Presentation(out_path).slides),
            len(edits or []),
            len(slide_images or []),
        )
    except Exception:
        pass

    return out_path


def _insert_slide_images(
    prs: Presentation,
    slides_list: list,
    slide_images: list[dict],
) -> None:
    """Insert user-uploaded images on the specified slides using position data."""
    for img_info in slide_images:
        slide_idx = img_info.get("slide_index", -1)
        asset_id = img_info.get("asset_id", "")
        ext = img_info.get("ext", ".png")
        if slide_idx < 0 or slide_idx >= len(slides_list):
            continue
        asset_path = os.path.join(SESSION_DIR, "assets", f"{asset_id}{ext}")
        if not os.path.exists(asset_path):
            continue

        slide = slides_list[slide_idx]
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Use position/size from frontend if provided (values are 0-1 fractions)
        x_pct = img_info.get("x_pct")
        y_pct = img_info.get("y_pct")
        w_pct = img_info.get("w_pct")

        if x_pct is not None and y_pct is not None and w_pct is not None:
            # Frontend provides exact position and width as fractions of slide
            target_w = int(slide_w * w_pct)
            # Preserve original aspect ratio for height
            try:
                from PIL import Image as PILImage
                with PILImage.open(asset_path) as pil_img:
                    orig_w, orig_h = pil_img.size
                aspect = orig_h / orig_w if orig_w else 1.0
            except Exception:
                aspect = 0.75
            target_h = int(target_w * aspect)
            left = int(slide_w * x_pct)
            top = int(slide_h * y_pct)
        else:
            # Fallback: center at 40% width
            try:
                from PIL import Image as PILImage
                with PILImage.open(asset_path) as pil_img:
                    orig_w, orig_h = pil_img.size
                aspect = orig_h / orig_w if orig_w else 1.0
            except Exception:
                aspect = 0.75
            target_w = int(slide_w * 0.4)
            target_h = int(target_w * aspect)
            if target_h > int(slide_h * 0.6):
                target_h = int(slide_h * 0.6)
                target_w = int(target_h / aspect) if aspect else target_w
            left = (slide_w - target_w) // 2
            top = (slide_h - target_h) // 2

        slide.shapes.add_picture(asset_path, left, top, target_w, target_h)
        logger.info("[export] Inserted image %s on slide %d at (%.2f, %.2f) w=%.2f",
                    asset_id, slide_idx, x_pct or 0, y_pct or 0, w_pct or 0.4)


# ── Re-render session with edits ──────────────────────────────────────────────


def re_render_session(
    session_id: str,
    edits: list[dict],
    slide_images: list[dict] | None = None,
) -> dict[str, Any]:
    """Apply text edits and images to the session PPTX, re-render PNGs, return updated meta.

    This modifies the session's ``original.pptx`` in-place so all subsequent
    previews and exports reflect the saved edits.
    """
    sess_dir = _session_path(session_id)
    original_path = os.path.join(sess_dir, "original.pptx")
    if not os.path.exists(original_path):
        raise FileNotFoundError(f"Session not found: {session_id}")

    # Load meta for theme info
    meta_path = os.path.join(sess_dir, "meta.json")
    with open(meta_path, encoding="utf-8") as f:
        old_meta = json.load(f)

    prs = Presentation(original_path)

    # Build edit lookup: { "s0_ph1": "new text" }
    edit_map: dict[str, str] = {}
    for edit in edits:
        if "content" in edit:
            edit_map[edit["element_id"]] = edit["content"]

    # Apply edits to shapes
    slides_list = list(prs.slides)
    for slide_idx, slide in enumerate(slides_list):
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            ph = shape.placeholder_format
            elem_id = f"s{slide_idx}_ph{ph.idx}"
            new_text = edit_map.get(elem_id)
            if new_text is None:
                continue
            if hasattr(shape, "text_frame") and (new_text != (shape.text or "")):
                shape.text = new_text

    # Insert per-slide images (user-uploaded)
    if slide_images:
        slides_list = list(prs.slides)
        _insert_slide_images(prs, slides_list, slide_images)

    # Save edited PPTX back as original so future exports use it
    prs.save(original_path)
    _dedup_zip_file(original_path)

    # Re-render PNGs
    slides_dir = os.path.join(sess_dir, "slides")
    with tempfile.TemporaryDirectory() as tmp:
        raw_pngs = _render_pptx_to_pngs(original_path, tmp)
        for i, raw_png in enumerate(raw_pngs):
            _resize_png(raw_png, os.path.join(slides_dir, f"{i}.png"))

    # Re-extract elements from updated PPTX
    prs_updated = Presentation(original_path)
    slide_width_pt = _emu_to_pt(prs_updated.slide_width)
    slide_height_pt = _emu_to_pt(prs_updated.slide_height)

    total_slides = len(prs_updated.slides)
    slides_payload: list[dict] = []
    for idx in range(total_slides):
        slide = prs_updated.slides[idx]
        slides_payload.append({
            "index": idx,
            "preview_url": f"/api/slides/editor-img/{session_id}/{idx}.png",
            "elements": _extract_elements(slide, idx),
        })

    # Update meta
    meta = {
        "session_id": session_id,
        "theme": old_meta.get("theme", ""),
        "slide_width_pt": slide_width_pt,
        "slide_height_pt": slide_height_pt,
        "slides": slides_payload,
    }
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

    return meta
