from __future__ import annotations

import asyncio

from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from ...business.table_handler import BusinessTableHandler
from .styles import TableBuilder


class PlaceholderProcessingMixin:
    def _read_table_csv(self, table_index, presentation_title):
        return TableBuilder.read_table_csv(table_index, presentation_title)

    def _create_table(self, slide, table_data, left, top, width, height):
        TableBuilder.create_table(slide, table_data, left, top, width, height)

    def _set_cell_content_with_linebreaks(self, cell, content):
        TableBuilder._set_cell_content_with_linebreaks(cell, content)

    def _determine_content_font_size(self, bullet_count, avg_words_per_bullet):
        return self.ppt_utils.determine_content_font_size(bullet_count, avg_words_per_bullet)

    def _insert_picture_into_placeholder(self, slide, placeholder, image_path):
        self.ppt_utils.insert_picture_into_placeholder(slide, placeholder, image_path)

    def _process_placeholders(self, slide, slide_data, presentation_title, prs=None, is_title_slide=False):
        title = slide_data.get("title", "")
        title_word_count = len(title.split()) if title else 0
        content_list = self.clean_bullets(slide_data.get("content", []))
        bullet_count = len(content_list)
        total_words = sum(len(content.split()) for content in content_list)
        avg_words_per_bullet = total_words / bullet_count if bullet_count > 0 else 0

        title_font_size = Pt(24) if title_word_count > 4 else None
        content_font_size = self._determine_content_font_size(bullet_count, avg_words_per_bullet)
        placeholder_infos = []
        image_data_list = []
        chart_type = slide_data.get("chart_type", "")
        chart_reasoning = slide_data.get("chart_reasoning", [])
        original_text = slide_data.get("original_text", "")
        content_written = False

        def fill_content_text(target_shape):
            nonlocal content_written
            text_frame = target_shape.text_frame
            text_frame.auto_size = MSO_AUTO_SIZE.NONE
            text_frame.word_wrap = True
            text_frame.clear()

            shape_w_pt, shape_h_pt = self.shape_dimensions_pt(target_shape)
            chosen_pt = self.fit_font_size(
                content_list,
                shape_w_pt,
                shape_h_pt,
                preferred_pt=content_font_size.pt,
            )
            final_font_size = Pt(chosen_pt)
            self.log_slide_layout_audit(
                slide_idx=slide_data.get("slide_number", "?"),
                title=title,
                layout_name=getattr(getattr(target_shape, "placeholder_format", None), "type", "N/A"),
                shape_w_pt=shape_w_pt,
                shape_h_pt=shape_h_pt,
                bullet_count=len(content_list),
                initial_pt=content_font_size.pt,
                final_pt=chosen_pt,
            )

            if content_list:
                for index, content in enumerate(content_list):
                    paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                    paragraph.text = content
                    paragraph.level = 0
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.font.size = final_font_size
            elif title:
                paragraph = text_frame.paragraphs[0]
                paragraph.text = title
                paragraph.font.size = final_font_size
            content_written = True

        def is_fallback_text_placeholder(target_shape):
            if not getattr(target_shape, "is_placeholder", False):
                return False
            if not getattr(target_shape, "has_text_frame", False):
                return False
            placeholder_type = target_shape.placeholder_format.type
            return placeholder_type not in {1, 3, 4, 5, 6, 7, 8, 10, 11, 12, 13, 15, 16, 18}

        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            placeholder_type = shape.placeholder_format.type

            if placeholder_type == 3:
                shape.text = slide_data.get("title", "")
            elif placeholder_type == 4:
                if isinstance(slide_data.get("metadata"), dict):
                    metadata_text = []
                    if slide_data["metadata"].get("author"):
                        metadata_text.append(f"Author: {slide_data['metadata']['author']}")
                    if slide_data["metadata"].get("date"):
                        metadata_text.append(f"Date: {slide_data['metadata']['date']}")
                    if slide_data["metadata"].get("description"):
                        metadata_text.append(f"Description: {slide_data['metadata']['description']}")
                    shape.text = "\n".join(metadata_text)
                else:
                    shape.text = slide_data.get("title", "")
            elif placeholder_type == 1:
                shape.text = slide_data.get("title", "")
                if title_font_size:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = title_font_size
            elif placeholder_type == 2:
                fill_content_text(shape)
            elif placeholder_type == 18:
                aspect_ratio = shape.width / shape.height
                if abs(aspect_ratio - 1.778) < 0.1:
                    ratio = 1
                elif abs(aspect_ratio - 1.333) < 0.1:
                    ratio = 0
                else:
                    continue
                placeholder_info = {
                    "shape": shape,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "placeholder_type": placeholder_type,
                    "aspect_ratio": aspect_ratio,
                    "ratio": ratio,
                    "image_type": "image",
                }
                placeholder_infos.append(placeholder_info)
                image_data_list.append(
                    {
                        "title": title,
                        "content_list": content_list,
                        "ratio": ratio,
                        "type": "image",
                        "chart_type": chart_type,
                        "chart_reasoning": chart_reasoning,
                        "original_text": original_text,
                        "placeholder_index": len(image_data_list),
                    }
                )
            elif placeholder_type == 7:
                if (not content_written) and getattr(shape, "has_text_frame", False) and content_list:
                    fill_content_text(shape)
                    continue
                if content_written and not self._is_meaningful_chart_type(chart_type):
                    continue

                aspect_ratio = shape.width / shape.height
                if abs(aspect_ratio - 1.778) < 0.1:
                    ratio = 1
                elif abs(aspect_ratio - 1.333) < 0.1:
                    ratio = 0
                else:
                    continue

                placeholder_info = {
                    "shape": shape,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "placeholder_type": placeholder_type,
                    "aspect_ratio": aspect_ratio,
                    "ratio": ratio,
                    "image_type": "diagram",
                }
                placeholder_infos.append(placeholder_info)
                image_data_list.append(
                    {
                        "title": title,
                        "content_list": content_list,
                        "ratio": ratio,
                        "type": "diagram",
                        "chart_type": chart_type,
                        "chart_reasoning": chart_reasoning,
                        "original_text": original_text,
                        "placeholder_index": len(image_data_list),
                    }
                )
            elif not content_written and is_fallback_text_placeholder(shape):
                fill_content_text(shape)

        if content_list and not content_written:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if getattr(shape, "is_placeholder", False):
                    placeholder_type = shape.placeholder_format.type
                    if placeholder_type in {1, 3, 4, 5, 6, 7, 8, 10, 11, 12, 13, 15, 16, 18}:
                        continue
                fill_content_text(shape)
                break

        if placeholder_infos:
            slide_title = slide_data.get("title", "Unknown")
            print(f"🔍 [Image Processing] Found {len(placeholder_infos)} placeholders in slide: {slide_title}")
            if self.is_collecting:
                self.collected_tasks.append(
                    {
                        "slide": slide,
                        "slide_data": slide_data,
                        "placeholder_infos": placeholder_infos,
                        "slide_title": slide_title,
                    }
                )
                print(f"📦 [Batch Processing] Collected task for slide: {slide_title} ({len(placeholder_infos)} placeholders)")
            else:
                print(f"🚀 [Image Processing] Processing {len(placeholder_infos)} placeholders in parallel for slide: {slide_title}")
                image_paths = asyncio.run(self.image_processor.process_multiple_images_async(image_data_list))
                for index, (placeholder_info, image_path) in enumerate(zip(placeholder_infos, image_paths)):
                    try:
                        if image_path:
                            self._insert_picture_into_placeholder(slide, placeholder_info["shape"], image_path)
                            print(f"✅ [Image Processing] Applied image {index + 1}/{len(image_paths)} to placeholder successfully")
                        else:
                            print(f"⚠️ [Image Processing] No image generated for placeholder {index + 1}")
                    except Exception as exc:
                        print(f"❌ [Image Processing] Error applying image {index + 1} to placeholder: {exc}")
                print(f"🎉 [Image Processing] Completed processing {len(placeholder_infos)} placeholders for slide: {slide_title}")

        if not is_title_slide and slide_data.get("slide_number"):
            left = Inches(9.40)
            top = Inches(6.95)
            width = Inches(3)
            height = Inches(0.5)
            slide_number_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = slide_number_box.text_frame
            paragraph = text_frame.paragraphs[0]
            paragraph.text = str(slide_data.get("slide_number", ""))
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.RIGHT

        BusinessTableHandler.process_tables_generic(
            slide,
            slide_data,
            presentation_title,
            self._read_table_csv,
            self._create_table,
        )
        self._process_latex_formulas(slide, slide_data)
