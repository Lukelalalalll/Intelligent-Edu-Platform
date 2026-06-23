from __future__ import annotations

import asyncio
import os

from pptx import Presentation

from .. import ppt_utils
from ..theme_catalog import resolve_base_theme


class PresentationBuilderMixin:
    def _get_template_path(self, theme):
        return self.ppt_utils.get_template_path(self.template_base_path, theme)

    def _find_layout_by_name(self, prs, layout_name):
        return self.ppt_utils.find_layout_by_name(prs, layout_name)

    @staticmethod
    def _is_meaningful_chart_type(chart_type) -> bool:
        return ppt_utils.is_meaningful_chart_type(chart_type)

    @staticmethod
    def _layout_has_body(layout) -> bool:
        return ppt_utils.layout_has_body(layout)

    def _find_content_layout(self, prs):
        return self.ppt_utils.find_content_layout(prs)

    @staticmethod
    def _clear_existing_slides(prs):
        ppt_utils.clear_existing_slides(prs)

    def _get_template_creator_mapping(self):
        return dict(self.ppt_utils.THEME_CREATOR_MAPPING)

    def _should_use_specialized_creator(self, theme):
        return self.ppt_utils.should_use_specialized_creator(self.template_base_path, theme)

    def _get_specialized_creator(self, theme):
        creator_mapping = self._get_template_creator_mapping()
        try:
            available_themes = [
                os.path.splitext(name)[0]
                for name in os.listdir(self.template_base_path)
                if name.endswith(".pptx")
            ]
            resolved_theme = resolve_base_theme(theme, available_themes)
        except Exception:
            resolved_theme = theme

        creator_class_name = creator_mapping.get(resolved_theme.lower())
        if not creator_class_name:
            return None

        creator_import_map: dict[str, type] = {}
        try:
            from ..business_ppt_creator import BusinessPPTCreator

            creator_import_map["BusinessPPTCreator"] = BusinessPPTCreator
        except ImportError as import_error:
            print(f"⚠️ [CreatorImport] Could not import BusinessPPTCreator: {import_error}")
        try:
            from ..light_ppt_creator import LightPPTCreator

            creator_import_map["LightPPTCreator"] = LightPPTCreator
        except ImportError as import_error:
            print(f"⚠️ [CreatorImport] Could not import LightPPTCreator: {import_error}")
        try:
            from ..dark_ppt_creator import DarkPPTCreator

            creator_import_map["DarkPPTCreator"] = DarkPPTCreator
        except ImportError as import_error:
            print(f"⚠️ [CreatorImport] Could not import DarkPPTCreator: {import_error}")

        creator_class = creator_import_map.get(creator_class_name)
        if creator_class is None:
            print(f"⚠️ [CreatorImport] No class found for '{creator_class_name}' (theme='{theme}'). Falling back to default creator.")
            return None

        creator_instance = creator_class(self.template_base_path)
        print(f"✅ [CreatorImport] Using {creator_class_name} for theme: {resolved_theme}")
        return creator_instance

    def create_presentation(self, ppt_schema, output_path):
        if not ppt_schema or "slides" not in ppt_schema:
            raise ValueError("Invalid PPT schema")

        theme = ppt_schema.get("theme", "Dark")
        if self._should_use_specialized_creator(theme):
            specialized_creator = self._get_specialized_creator(theme)
            if specialized_creator:
                return specialized_creator.create_presentation(ppt_schema, output_path)
            print(f"Specialized creator not available, falling back to default creator for theme: {theme}")

        template_path = self._get_template_path(theme)
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")

        prs = Presentation(template_path)
        self._clear_existing_slides(prs)
        presentation_title = ppt_schema.get("presentation_title", "")

        title_layout = self._find_layout_by_name(prs, "Title")
        if not title_layout:
            for layout in prs.slide_layouts:
                placeholder_types = {shape.placeholder_format.type for shape in layout.placeholders}
                if placeholder_types & {1, 3} and not (placeholder_types & {2, 7}):
                    title_layout = layout
                    break
            if not title_layout and prs.slide_layouts:
                title_layout = prs.slide_layouts[0]

        if title_layout:
            title_slide = prs.slides.add_slide(title_layout)
            title_data = {
                "title": presentation_title,
                "metadata": ppt_schema.get("metadata", {}),
            }
            self._process_placeholders(title_slide, title_data, presentation_title, prs, True)

        print("🔄 [Batch Processing] Starting batch collection mode for image placeholders...")
        self.start_collecting()
        failed_slides = []

        for slide_index, slide_data in enumerate(ppt_schema["slides"]):
            try:
                if "layout" not in slide_data:
                    continue

                layout_raw = slide_data.get("layout")
                if isinstance(layout_raw, dict):
                    layout_name = (layout_raw.get("name") or "").strip()
                elif layout_raw is None:
                    layout_name = ""
                else:
                    layout_name = str(layout_raw).strip()

                if not layout_name:
                    continue

                layout = self._find_layout_by_name(prs, layout_name)
                if not layout:
                    print(f"Warning: Layout '{layout_name}' not found in template — trying content-capable fallback")
                    layout = self._find_content_layout(prs)
                    if not layout:
                        print("Warning: No content-capable fallback layout found — skipping slide")
                        continue

                content_list = slide_data.get("content", [])
                if content_list and not self._layout_has_body(layout):
                    fallback = self._find_content_layout(prs)
                    if fallback:
                        print(
                            f"⚠️ Layout '{layout_name}' has no body placeholder but slide has {len(content_list)} bullets — "
                            f"falling back to '{fallback.name}'"
                        )
                        layout = fallback

                slide = prs.slides.add_slide(layout)
                slide_data["slide_number"] = str(slide_index + 1)
                self._process_placeholders(slide, slide_data, presentation_title, prs)
            except Exception as exc:
                slide_title = slide_data.get("title", f"Slide {slide_index + 1}")
                print(f"❌ Slide {slide_index + 1} '{slide_title}' failed: {exc}")
                failed_slides.append(
                    {
                        "slide_index": slide_index + 1,
                        "title": slide_title,
                        "error": str(exc),
                    }
                )

        self.stop_collecting()
        print("⚡ [Batch Processing] Executing batch processing for all collected image placeholders...")
        try:
            asyncio.run(self.process_all_collected_tasks())
        except Exception as exc:
            print(f"⚠️ Batch image processing failed: {exc} — PPT saved without some images")

        if failed_slides:
            print(f"⚠️ {len(failed_slides)}/{len(ppt_schema['slides'])} slides had errors: {[slide['title'] for slide in failed_slides]}")

        prs.save(output_path)
        return output_path
