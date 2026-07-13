"""Table creation and styling helpers extracted from PPTCreator."""

import os
import pandas as pd
from pptx.dml.color import RGBColor
from pptx.util import Pt


class TableBuilder:
    """Builds styled tables for PPT slides using a seaborn-inspired palette."""

    @staticmethod
    def read_table_csv(table_index: int, presentation_title: str) -> dict | None:
        csv_path = f"md/csv/{presentation_title}_tables_{table_index}.csv"
        if not os.path.exists(csv_path):
            return None

        try:
            df = pd.read_csv(csv_path)
            headers = [
                '' if 'Unnamed:' in str(col) else str(col)
                for col in df.columns
            ]
            rows = []
            for _, row in df.iterrows():
                cleaned_row = [
                    '' if (pd.isna(value) or str(value).lower() == 'nan')
                    else str(value)
                    for value in row
                ]
                rows.append(cleaned_row)
            return {"header": headers, "rows": rows}
        except Exception as e:
            print(f"Error reading table CSV: {e}")
            return None

    @staticmethod
    def create_table(slide, table_data: dict, left, top, width, height):
        rows = len(table_data["rows"]) + 1
        cols = len(table_data["header"])
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        header_bg_color = RGBColor(31, 119, 180)
        even_row_bg_color = RGBColor(248, 248, 248)
        odd_row_bg_color = RGBColor(255, 255, 255)
        header_text_color = RGBColor(255, 255, 255)
        data_text_color = RGBColor(51, 51, 51)

        for col_idx, header in enumerate(table_data["header"]):
            cell = table.cell(0, col_idx)
            TableBuilder._set_cell_content_with_linebreaks(cell, str(header))
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg_color
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = header_text_color

        for row_idx, row_data in enumerate(table_data["rows"]):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                TableBuilder._set_cell_content_with_linebreaks(cell, str(cell_data))
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = even_row_bg_color
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = odd_row_bg_color
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = data_text_color

    @staticmethod
    def _set_cell_content_with_linebreaks(cell, content: str):
        if '<br>' in content:
            lines = content.split('<br>')
            cell.text_frame.clear()
            if lines:
                first_paragraph = cell.text_frame.paragraphs[0]
                first_paragraph.text = lines[0].strip()
                for line in lines[1:]:
                    new_paragraph = cell.text_frame.add_paragraph()
                    new_paragraph.text = line.strip()
        else:
            cell.text = content

