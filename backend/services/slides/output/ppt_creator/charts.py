"""LaTeX formula rendering and image placement helpers extracted from PPTCreator."""

from pptx.util import Pt


class LatexRenderer:
    """Renders LaTeX formulas as images and inserts them into slides."""

    @staticmethod
    def process_slide_latex(slide_data: dict, slide_id: str) -> dict:
        """Convert LaTeX formulas in slide_data to images.

        Returns a dict mapping formula strings to generated image paths.
        """
        from backend.services.slides.generation.latex_generator import process_slide_latex as _process

        latex_formulas = slide_data.get('latex', [])
        if not latex_formulas or not any(formula.strip() for formula in latex_formulas):
            return {}

        formula_images = _process(slide_data, slide_id)
        if not formula_images:
            print(f"⚠️ Slide {slide_id} LaTeX formulas processing failed")
            return {}

        print(f"✅ Slide {slide_id} successfully processed {len(formula_images)} formulas")
        return formula_images

    @staticmethod
    def insert_latex_images(slide, formula_images: dict, insert_picture_fn):
        """Insert LaTeX formula images into the slide using available placeholders."""
        other_placeholders = LatexRenderer._collect_other_placeholders(slide)

        for i, (formula, image_path) in enumerate(formula_images.items()):
            if i < len(other_placeholders):
                placeholder_info = other_placeholders[i]
                try:
                    insert_picture_fn(
                        slide,
                        placeholder_info['shape'],
                        image_path,
                        placeholder_info['left'],
                        placeholder_info['top'],
                        placeholder_info['width'],
                        placeholder_info['height'],
                    )
                    print(f"✅ Formula image inserted into the slide: {formula[:50]}...")
                except Exception as e:
                    print(f"❌ Formula image failed to insert into the slide: {e}")
            else:
                print(f"⚠️ No enough placeholders to insert formula: {formula[:50]}...")

    @staticmethod
    def _collect_other_placeholders(slide) -> list:
        """Collect placeholder info excluding common types (title, body, etc.)."""
        other_placeholders = []
        excluded_types = {1, 2, 3, 4, 13}

        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            placeholder_type = shape.placeholder_format.type
            if placeholder_type not in excluded_types:
                other_placeholders.append({
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'type': placeholder_type,
                    'shape': shape,
                })

        other_placeholders.sort(key=lambda x: (x['top'], x['left']))
        return other_placeholders
