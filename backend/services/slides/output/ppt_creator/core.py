import os
import asyncio
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from ...generation.img_chart_processor import ImageChartProcessor
from ...business.table_handler import BusinessTableHandler
from ..theme_catalog import resolve_base_theme
from ..text_layout_engine import (
    fit_font_size,
    clean_bullets,
    shape_dimensions_pt,
    log_slide_layout_audit,
)
from ..batch_context import BatchContext
from .. import ppt_utils
from .styles import TableBuilder
from .charts import LatexRenderer


class PPTCreator:
    def __init__(self, template_base_path=None):
        self.template_base_path = template_base_path or "static/ppt_templates"
        self.image_processor = ImageChartProcessor()
        self._ctx = BatchContext(
            template_base_path=self.template_base_path,
            image_processor=self.image_processor,
        )

        # Keep legacy attributes as aliases to BatchContext fields
        # so subclass code that reads/writes them still works.
        self.collected_tasks = self._ctx.collected_tasks
        self.is_collecting = self._ctx.is_collecting
        self.batch_results = self._ctx.batch_results
        
    def start_collecting(self):
        """开始收集模式"""
        self._ctx.start_collecting()
        self.collected_tasks = self._ctx.collected_tasks
        self.is_collecting = self._ctx.is_collecting
        self.batch_results = self._ctx.batch_results
        
    def stop_collecting(self):
        """停止收集模式"""
        self._ctx.stop_collecting()
        self.is_collecting = self._ctx.is_collecting
        
    async def process_all_collected_tasks(self):
        """批量处理所有收集的任务"""
        if not self.collected_tasks:
            print("ℹ️ [Batch Processing] No collected tasks to process")
            return
            
        print(f"🚀 [Batch Processing] Starting batch processing of {len(self.collected_tasks)} collected tasks...")
        
        # 准备所有图片数据
        all_image_data = []
        task_mappings = []  # 用于映射结果回对应的任务和占位符
        
        for task_idx, task_info in enumerate(self.collected_tasks):
            slide_data = task_info['slide_data']
            placeholder_infos = task_info['placeholder_infos']
            slide_title = task_info['slide_title']
            
            # 为每个任务的每个占位符创建图片数据
            for placeholder_idx, placeholder_info in enumerate(placeholder_infos):
                image_data = self._prepare_image_data(slide_data, placeholder_info, placeholder_idx)
                all_image_data.append(image_data)
                
                # 记录映射关系
                task_mappings.append({
                    'task_idx': task_idx,
                    'placeholder_idx': placeholder_idx,
                    'slide_title': slide_title
                })
        
        # 批量并行处理所有图片数据
        total_count = len(all_image_data)
        print(f"⚡ [Batch Processing] Processing {total_count} image placeholders in parallel...")
        
        self.batch_results = await self.image_processor.process_multiple_images_async(all_image_data)
        
        print(f"✅ [Batch Processing] Batch processing completed! Generated {len(self.batch_results)} images")
        
        # 应用结果到对应的占位符
        await self._apply_batch_results(task_mappings)
        
    async def _apply_batch_results(self, task_mappings):
        """应用批量处理结果到对应的占位符"""
        print(f"📌 [Batch Processing] Applying {len(self.batch_results)} results to placeholders...")
        
        applied_count = 0
        error_count = 0
        
        for result_idx, (image_path, mapping) in enumerate(zip(self.batch_results, task_mappings)):
            try:
                task_info = self.collected_tasks[mapping['task_idx']]
                slide = task_info['slide']
                placeholder_info = task_info['placeholder_infos'][mapping['placeholder_idx']]
                slide_title = mapping['slide_title']
                
                if image_path:
                    self._insert_picture_into_placeholder(slide, placeholder_info['shape'], image_path)
                    applied_count += 1
                    print(f"✅ [Batch Processing] Applied result {result_idx+1}/{len(self.batch_results)} to slide: {slide_title}")
                else:
                    print(f"⚠️ [Batch Processing] No image generated for result {result_idx+1} in slide: {slide_title}")
                    
            except Exception as e:
                error_count += 1
                print(f"❌ [Batch Processing] Error applying result {result_idx+1}: {e}")
        
        print(f"🎉 [Batch Processing] Batch application completed! Applied: {applied_count}, Errors: {error_count}")
        
    def _prepare_image_data(self, slide_data, placeholder_info, placeholder_index):
        """准备单个占位符的图片数据"""
        return ppt_utils.prepare_image_data(slide_data, placeholder_info, placeholder_index)

    def _get_template_path(self, theme):
        """获取主题模板路径"""
        return ppt_utils.get_template_path(self.template_base_path, theme)
    
    def _find_layout_by_name(self, prs, layout_name):
        """根据布局名称查找对应的布局"""
        return ppt_utils.find_layout_by_name(prs, layout_name)

    @staticmethod
    def _is_meaningful_chart_type(chart_type) -> bool:
        """Return True only when chart_type represents a real chart/diagram request."""
        return ppt_utils.is_meaningful_chart_type(chart_type)

    @staticmethod
    def _layout_has_body(layout) -> bool:
        """Check if a slide layout has a BODY (2) or OBJECT (7) placeholder."""
        return ppt_utils.layout_has_body(layout)

    def _find_content_layout(self, prs):
        """Find the first slide layout that contains a body/object placeholder."""
        return ppt_utils.find_content_layout(prs)

    @staticmethod
    def _clear_existing_slides(prs):
        """Remove all slides from a template while keeping masters/layouts intact."""
        ppt_utils.clear_existing_slides(prs)

    def _collect_visual_tasks(
        self, slide, slide_data: dict, *, content_was_handled: bool = True
    ):
        """Collect PICTURE (type 18) and OBJECT/Diagram (type 7) placeholders for
        batch or immediate AI image/diagram generation.

        Unlike _process_placeholders(), this method NEVER writes text to any
        placeholder — it is safe to call after the specialized creator has already
        filled title and body text (Bug 1 fix).

        Args:
            slide:               python-pptx Slide object.
            slide_data:          slide dict (with original content list intact).
            content_was_handled: when True, type-7 OBJECT placeholders are only
                                 queued when chart_type is a meaningful diagram
                                 type (Bug 2 fix / guard now always fires).
        """
        chart_type = slide_data.get('chart_type', '')
        slide_title = slide_data.get('title', 'Unknown')
        placeholder_infos: list = []

        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            ptype = shape.placeholder_format.type

            if ptype == 18:  # PICTURE
                aspect_ratio = shape.width / shape.height
                if abs(aspect_ratio - 1.778) < 0.1:
                    ratio = 1
                elif abs(aspect_ratio - 1.333) < 0.1:
                    ratio = 0
                else:
                    continue  # unsupported aspect ratio
                placeholder_infos.append({
                    'shape': shape,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'placeholder_type': ptype,
                    'aspect_ratio': aspect_ratio,
                    'ratio': ratio,
                    'image_type': 'image',
                })

            elif ptype == 7:  # OBJECT / Diagram
                # Bug 2 fix: guard now always fires correctly because
                # content_was_handled is an explicit flag, not inferred from
                # a mutable `content_written` variable that could be False.
                if content_was_handled and not self._is_meaningful_chart_type(chart_type):
                    continue
                aspect_ratio = shape.width / shape.height
                if abs(aspect_ratio - 1.778) < 0.1:
                    ratio = 1
                elif abs(aspect_ratio - 1.333) < 0.1:
                    ratio = 0
                else:
                    continue
                placeholder_infos.append({
                    'shape': shape,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'placeholder_type': ptype,
                    'aspect_ratio': aspect_ratio,
                    'ratio': ratio,
                    'image_type': 'diagram',
                })

        if not placeholder_infos:
            return

        print(f'🔍 [Visual] {len(placeholder_infos)} visual placeholder(s) in: "{slide_title}"')

        if self.is_collecting:
            self.collected_tasks.append({
                'slide': slide,
                'slide_data': slide_data,
                'placeholder_infos': placeholder_infos,
                'slide_title': slide_title,
            })
            print(f'📦 [Batch] Collected visual task for: "{slide_title}"')
        else:
            image_data_list = [
                self._prepare_image_data(slide_data, ph, i)
                for i, ph in enumerate(placeholder_infos)
            ]
            print(f'🚀 [Visual] Processing {len(placeholder_infos)} placeholder(s) immediately...')
            image_paths = asyncio.run(
                self.image_processor.process_multiple_images_async(image_data_list)
            )
            for ph_info, image_path in zip(placeholder_infos, image_paths):
                if image_path:
                    self._insert_picture_into_placeholder(slide, ph_info['shape'], image_path)

    def _get_template_creator_mapping(self):
        """获取模板创建器映射"""
        return dict(ppt_utils.THEME_CREATOR_MAPPING)
    
    def _should_use_specialized_creator(self, theme):
        """判断是否应该使用专门的模板创建器"""
        return ppt_utils.should_use_specialized_creator(self.template_base_path, theme)
    
    def _get_specialized_creator(self, theme):
        """获取专门的模板创建器实例
        
        Args:
            theme (str): 主题名称
            
        Returns:
            PPTCreator: 专门的模板创建器实例，如果不存在则返回None
        """
        creator_mapping = self._get_template_creator_mapping()
        try:
            available_themes = [
                os.path.splitext(name)[0]
                for name in os.listdir(self.template_base_path)
                if name.endswith('.pptx')
            ]
            resolved_theme = resolve_base_theme(theme, available_themes)
        except Exception:
            resolved_theme = theme

        creator_class_name = creator_mapping.get(resolved_theme.lower())

        if not creator_class_name:
            return None

        # Explicit import map — avoids unreliable string-based reflection (P0 fix)
        _CREATOR_IMPORT_MAP: dict[str, type] = {}
        try:
            from ..business_ppt_creator import BusinessPPTCreator  # lazy, avoids circular at module level
            _CREATOR_IMPORT_MAP['BusinessPPTCreator'] = BusinessPPTCreator
        except ImportError as ie:
            print(f"⚠️ [CreatorImport] Could not import BusinessPPTCreator: {ie}")
        try:
            from ..light_ppt_creator import LightPPTCreator  # lazy import
            _CREATOR_IMPORT_MAP['LightPPTCreator'] = LightPPTCreator
        except ImportError as ie:
            print(f"⚠️ [CreatorImport] Could not import LightPPTCreator: {ie}")
        try:
            from ..dark_ppt_creator import DarkPPTCreator  # lazy import
            _CREATOR_IMPORT_MAP['DarkPPTCreator'] = DarkPPTCreator
        except ImportError as ie:
            print(f"⚠️ [CreatorImport] Could not import DarkPPTCreator: {ie}")

        creator_class = _CREATOR_IMPORT_MAP.get(creator_class_name)

        if creator_class is None:
            print(f"⚠️ [CreatorImport] No class found for '{creator_class_name}' (theme='{theme}'). Falling back to default creator.")
            return None

        creator_instance = creator_class(self.template_base_path)
        print(f"✅ [CreatorImport] Using {creator_class_name} for theme: {resolved_theme}")
        return creator_instance
    
    def _read_table_csv(self, table_index, presentation_title):
        return TableBuilder.read_table_csv(table_index, presentation_title)

    def _create_table(self, slide, table_data, left, top, width, height):
        TableBuilder.create_table(slide, table_data, left, top, width, height)

    def _set_cell_content_with_linebreaks(self, cell, content):
        TableBuilder._set_cell_content_with_linebreaks(cell, content)
    
    def _determine_content_font_size(self, bullet_count, avg_words_per_bullet):
        """根据bullet points数量和平均词数确定字体大小"""
        return ppt_utils.determine_content_font_size(bullet_count, avg_words_per_bullet)
    
    def _insert_picture_into_placeholder(self, slide, placeholder, image_path):
        """将图片插入到占位符中"""
        ppt_utils.insert_picture_into_placeholder(slide, placeholder, image_path)
    
    def _process_placeholders(self, slide, slide_data, presentation_title, prs=None, is_title_slide=False):
        """处理幻灯片占位符"""
        
        # 计算标题词数
        title = slide_data.get('title', '')
        title_word_count = len(title.split()) if title else 0
        
        # 计算bullet points数量和平均词数
        content_list = slide_data.get('content', [])

        # ── Bullet cleaner (P0/P1 fix) ───────────────────────────────────────
        content_list = clean_bullets(content_list)

        bullet_count = len(content_list)
        total_words = sum(len(content.split()) for content in content_list)
        avg_words_per_bullet = total_words / bullet_count if bullet_count > 0 else 0
        
        # 确定字体大小
        title_font_size = Pt(24) if title_word_count > 4 else None
        content_font_size = self._determine_content_font_size(bullet_count, avg_words_per_bullet)
        
        # 收集图片和对象占位符信息
        placeholder_infos = []
        image_data_list = []
        
        # 准备图片处理数据
        chart_type = slide_data.get('chart_type', '')
        chart_reasoning = slide_data.get('chart_reasoning', [])
        original_text = slide_data.get('original_text', '')
        content_written = False

        def fill_content_text(target_shape):
            nonlocal content_written
            text_frame = target_shape.text_frame

            # ── Controlled scaling (P0 fix: replace TEXT_TO_FIT_SHAPE) ───────
            # Disable PowerPoint's built-in auto-size — it produces inconsistent
            # results across slides.  Instead, we estimate capacity ourselves and
            # step down the font size through a pre-defined ladder.
            text_frame.auto_size = MSO_AUTO_SIZE.NONE
            text_frame.word_wrap = True
            text_frame.clear()

            # Determine the best font size for this specific shape.
            shape_w_pt, shape_h_pt = shape_dimensions_pt(target_shape)
            chosen_pt = fit_font_size(
                content_list,
                shape_w_pt,
                shape_h_pt,
                preferred_pt=content_font_size.pt,
            )
            final_font_size = Pt(chosen_pt)

            # Emit layout audit line.
            log_slide_layout_audit(
                slide_idx=slide_data.get('slide_number', '?'),
                title=title,
                layout_name=getattr(getattr(target_shape, 'placeholder_format', None), 'type', 'N/A'),
                shape_w_pt=shape_w_pt,
                shape_h_pt=shape_h_pt,
                bullet_count=len(content_list),
                initial_pt=content_font_size.pt,
                final_pt=chosen_pt,
            )

            if content_list:
                for i, content in enumerate(content_list):
                    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                    p.text = content
                    p.level = 0
                    p.alignment = PP_ALIGN.LEFT
                    p.font.size = final_font_size
            elif title:
                p = text_frame.paragraphs[0]
                p.text = title
                p.font.size = final_font_size
            content_written = True

        def is_fallback_text_placeholder(target_shape):
            if not getattr(target_shape, 'is_placeholder', False):
                return False
            if not getattr(target_shape, 'has_text_frame', False):
                return False
            ptype = target_shape.placeholder_format.type
            return ptype not in {1, 3, 4, 5, 6, 7, 8, 10, 11, 12, 13, 15, 16, 18}
        
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
                
            placeholder_type = shape.placeholder_format.type
            
            # 处理大标题 (type=3)
            if placeholder_type == 3:
                shape.text = slide_data.get('title', '')
                
            # 处理副标题 (type=4)
            elif placeholder_type == 4:
                if isinstance(slide_data.get('metadata'), dict):
                    metadata_text = []
                    if slide_data['metadata'].get('author'):
                        metadata_text.append(f"Author: {slide_data['metadata']['author']}")
                    if slide_data['metadata'].get('date'):
                        metadata_text.append(f"Date: {slide_data['metadata']['date']}")
                    if slide_data['metadata'].get('description'):
                        metadata_text.append(f"Description: {slide_data['metadata']['description']}")
                    shape.text = "\n".join(metadata_text)
                else:
                    shape.text = slide_data.get('title', '')

            # 处理正文标题 (type=1)
            elif placeholder_type == 1:
                shape.text = slide_data.get('title', '')
                # 设置标题字体大小
                if title_font_size:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = title_font_size

            elif placeholder_type == 2:
                fill_content_text(shape)

                # text_frame.clear()
                #
                # if content_list:
                #     # 设置第一个段落
                #     p = text_frame.paragraphs[0]
                #     p.text = content_list[0]
                #     p.alignment = PP_ALIGN.LEFT
                #     p.level = 1
                #     p.space_after = Pt(12)  # 设置段落间距
                #     p.font.size = content_font_size
                #
                #     # 添加后续段落
                #     for content in content_list[1:]:
                #         p = text_frame.add_paragraph()
                #         p.text = content
                #         p.alignment = PP_ALIGN.LEFT
                #         p.level = 1
                #         p.space_after = Pt(12)  # 设置段落间距
                #         p.font.size = content_font_size
            
            # 收集图片占位符 (type=18)
            elif placeholder_type == 18:
                aspect_ratio = shape.width / shape.height
                if abs(aspect_ratio - 1.778) < 0.1:  # 16:9
                    ratio = 1
                elif abs(aspect_ratio - 1.333) < 0.1:  # 4:3
                    ratio = 0
                else:
                    continue  # 不支持的比例
                
                placeholder_info = {
                    'shape': shape,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'placeholder_type': placeholder_type,
                    'aspect_ratio': aspect_ratio,
                    'ratio': ratio,
                    'image_type': 'image'
                }
                placeholder_infos.append(placeholder_info)
                
                # 准备图片数据
                image_data = {
                    'title': title,
                    'content_list': content_list,
                    'ratio': ratio,
                    'type': 'image',
                    'chart_type': chart_type,
                    'chart_reasoning': chart_reasoning,
                    'original_text': original_text,
                    'placeholder_index': len(image_data_list)
                }
                image_data_list.append(image_data)

            # 收集对象占位符 (type=7)
            elif placeholder_type == 7:
                # Some templates (e.g. Light) use OBJECT as the only text body slot.
                # If content has not been written yet and this shape supports text,
                # prioritize writing bullets here instead of treating it as media.
                if (not content_written) and getattr(shape, 'has_text_frame', False) and content_list:
                    fill_content_text(shape)
                    continue

                # If text content was already written to a BODY (type=2) placeholder on
                # this slide, only proceed with AI chart generation when chart_type is
                # explicitly set to a real diagram type.  Without this guard, layouts
                # that contain BOTH a BODY and an OBJECT placeholder (e.g. Classic
                # "B1-D1-H", Light "Chart layout 1") would trigger spurious AI-generated
                # images on every slide, producing garbled output.
                if content_written and not self._is_meaningful_chart_type(chart_type):
                    continue

                aspect_ratio = shape.width / shape.height
                if abs(aspect_ratio - 1.778) < 0.1:  # 16:9
                    ratio = 1
                elif abs(aspect_ratio - 1.333) < 0.1:  # 4:3
                    ratio = 0
                else:
                    continue  # 不支持的比例

                placeholder_info = {
                    'shape': shape,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'placeholder_type': placeholder_type,
                    'aspect_ratio': aspect_ratio,
                    'ratio': ratio,
                    'image_type': 'diagram'
                }
                placeholder_infos.append(placeholder_info)

                # 准备图片数据
                image_data = {
                    'title': title,
                    'content_list': content_list,
                    'ratio': ratio,
                    'type': 'diagram',
                    'chart_type': chart_type,
                    'chart_reasoning': chart_reasoning,
                    'original_text': original_text,
                    'placeholder_index': len(image_data_list)
                }
                image_data_list.append(image_data)

            elif not content_written and is_fallback_text_placeholder(shape):
                fill_content_text(shape)

        # If layout has no body placeholder, still write content to first text-capable shape.
        if content_list and not content_written:
            for shape in slide.shapes:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                if getattr(shape, 'is_placeholder', False):
                    ptype = shape.placeholder_format.type
                    if ptype in {1, 3, 4, 5, 6, 7, 8, 10, 11, 12, 13, 15, 16, 18}:
                        continue
                fill_content_text(shape)
                break
        
        # 处理图片占位符
        if placeholder_infos:
            slide_title = slide_data.get('title', 'Unknown')
            print(f"🔍 [Image Processing] Found {len(placeholder_infos)} placeholders in slide: {slide_title}")
            
            # 根据是否在收集模式决定处理方式
            if self.is_collecting:
                # 收集模式：将任务加入队列
                task_info = {
                    'slide': slide,
                    'slide_data': slide_data,
                    'placeholder_infos': placeholder_infos,
                    'slide_title': slide_title
                }
                self.collected_tasks.append(task_info)
                print(f"📦 [Batch Processing] Collected task for slide: {slide_title} ({len(placeholder_infos)} placeholders)")
            else:
                # 正常模式：立即处理
                print(f"🚀 [Image Processing] Processing {len(placeholder_infos)} placeholders in parallel for slide: {slide_title}")
                image_paths = asyncio.run(self.image_processor.process_multiple_images_async(image_data_list))
                
                # 应用结果到对应占位符
                for i, (placeholder_info, image_path) in enumerate(zip(placeholder_infos, image_paths)):
                    try:
                        if image_path:
                            self._insert_picture_into_placeholder(slide, placeholder_info['shape'], image_path)
                            print(f"✅ [Image Processing] Applied image {i+1}/{len(image_paths)} to placeholder successfully")
                        else:
                            print(f"⚠️ [Image Processing] No image generated for placeholder {i+1}")
                    except Exception as e:
                        print(f"❌ [Image Processing] Error applying image {i+1} to placeholder: {e}")
                
                print(f"🎉 [Image Processing] Completed processing {len(placeholder_infos)} placeholders for slide: {slide_title}")
        
        # 如果不是标题页，且有slide_number数据，则在右下角添加幻灯片编号
        if not is_title_slide and slide_data.get('slide_number'):      
            # 在右下角添加页码文本框
            left = Inches(9.40)
            top = Inches(6.95)
            width = Inches(3)
            height = Inches(0.5)
            
            # 添加幻灯片编号文本框
            slide_number_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = slide_number_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = str(slide_data.get('slide_number', ''))
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.RIGHT  # 右对齐
        
        # 处理表格数据（使用通用表格处理方法）
        BusinessTableHandler.process_tables_generic(
            slide, slide_data, presentation_title, 
            self._read_table_csv, self._create_table
        )
        
        # 处理LaTeX公式（如果存在）
        self._process_latex_formulas(slide, slide_data)
    
    def _apply_speaker_notes(self, slide, slide_data):
        """Write speaker notes to a slide if provided in slide_data."""
        ppt_utils.apply_speaker_notes(slide, slide_data)

    def _process_latex_formulas(self, slide, slide_data):
        slide_id = f"slide_{slide_data.get('slide_number', 'unknown')}"
        formula_images = LatexRenderer.process_slide_latex(slide_data, slide_id)
        if formula_images:
            self._insert_latex_images(slide, formula_images)

    def _insert_latex_images(self, slide, formula_images):
        LatexRenderer.insert_latex_images(
            slide, formula_images,
            insert_picture_fn=self._insert_picture_with_aspect_ratio,
        )

    def _collect_other_placeholders(self, slide):
        return LatexRenderer._collect_other_placeholders(slide)
    
    def _insert_picture_with_aspect_ratio(self, slide, placeholder_shape, image_path, left, top, width, height):
        """将图片插入到占位符中，根据原始宽高比计算高度"""
        ppt_utils.insert_picture_with_aspect_ratio(slide, placeholder_shape, image_path, left, top, width, height)
    
    def create_presentation(self, ppt_schema, output_path):
        """创建演示文稿
        
        Args:
            ppt_schema (dict): PPT结构数据
            output_path (str): 输出文件路径
        """
        if not ppt_schema or 'slides' not in ppt_schema:
            raise ValueError("Invalid PPT schema")
            
        # 获取主题信息
        theme = ppt_schema.get('theme', 'Dark')
        
        # 检查是否应该使用专门的模板创建器
        if self._should_use_specialized_creator(theme):
            specialized_creator = self._get_specialized_creator(theme)
            if specialized_creator:
                return specialized_creator.create_presentation(ppt_schema, output_path)
            else:
                print(f"Specialized creator not available, falling back to default creator for theme: {theme}")
        
        # 使用默认的创建逻辑
        template_path = self._get_template_path(theme)
        
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
            
        # 创建演示文稿
        prs = Presentation(template_path)
        # Template files may contain demo/sample slides. Start from a clean deck
        # so exported content never includes original template pages.
        self._clear_existing_slides(prs)
        
        # 获取演示文稿标题
        presentation_title = ppt_schema.get('presentation_title', '')
        
        # 创建标题幻灯片
        title_layout = self._find_layout_by_name(prs, 'Title')
        if not title_layout:
            # Fuzzy fallback: some templates name the title layout in other languages
            # (e.g. Light.pptx uses "标题幻灯片").  Find the first layout that has a
            # CENTER_TITLE or TITLE placeholder but no BODY/OBJECT content area.
            for _layout in prs.slide_layouts:
                _ph_types = {sh.placeholder_format.type for sh in _layout.placeholders}
                if _ph_types & {1, 3} and not (_ph_types & {2, 7}):
                    title_layout = _layout
                    break
            if not title_layout and prs.slide_layouts:
                title_layout = prs.slide_layouts[0]
        if title_layout:
            title_slide = prs.slides.add_slide(title_layout)
            title_data = {
                'title': presentation_title,
                'metadata': ppt_schema.get('metadata', {})
            }
            self._process_placeholders(title_slide, title_data, presentation_title, prs, True)
        
        # 处理每个幻灯片
        print("🔄 [Batch Processing] Starting batch collection mode for image placeholders...")
        self.start_collecting()
        
        failed_slides = []
        for slide_index, slide_data in enumerate(ppt_schema['slides']):
            try:
                if 'layout' not in slide_data:
                    continue

                layout_raw = slide_data.get('layout')
                if isinstance(layout_raw, dict):
                    layout_name = (layout_raw.get('name') or '').strip()
                elif layout_raw is None:
                    layout_name = ''
                else:
                    layout_name = str(layout_raw).strip()

                if not layout_name:
                    continue
                    
                # 查找对应的布局
                layout = self._find_layout_by_name(prs, layout_name)
                if not layout:
                    print(f"Warning: Layout '{layout_name}' not found in template — trying content-capable fallback")
                    layout = self._find_content_layout(prs)
                    if not layout:
                        print(f"Warning: No content-capable fallback layout found — skipping slide")
                        continue

                # If this slide has content bullets but the assigned layout has no
                # body/object placeholder, swap to a content-capable layout so the
                # text is not silently dropped.
                content_list = slide_data.get('content', [])
                if content_list and not self._layout_has_body(layout):
                    fallback = self._find_content_layout(prs)
                    if fallback:
                        print(f"⚠️ Layout '{layout_name}' has no body placeholder "
                              f"but slide has {len(content_list)} bullets — "
                              f"falling back to '{fallback.name}'")
                        layout = fallback
                    
                # 创建幻灯片
                slide = prs.slides.add_slide(layout)
                
                # 添加页码信息到slide_data
                slide_data['slide_number'] = str(slide_index + 1)  # 从第1页开始
                
                # 处理占位符
                self._process_placeholders(slide, slide_data, presentation_title, prs)
            except Exception as e:
                slide_title = slide_data.get('title', f'Slide {slide_index + 1}')
                print(f"❌ Slide {slide_index + 1} '{slide_title}' failed: {e}")
                failed_slides.append({
                    "slide_index": slide_index + 1,
                    "title": slide_title,
                    "error": str(e)
                })
                # Continue with remaining slides — don't abort the whole presentation
        
        # 停止收集模式并执行批量处理
        self.stop_collecting()
        print("⚡ [Batch Processing] Executing batch processing for all collected image placeholders...")
        try:
            asyncio.run(self.process_all_collected_tasks())
        except Exception as e:
            print(f"⚠️ Batch image processing failed: {e} — PPT saved without some images")
        
        if failed_slides:
            print(f"⚠️ {len(failed_slides)}/{len(ppt_schema['slides'])} slides had errors: {[s['title'] for s in failed_slides]}")
        
        # 保存演示文稿
        prs.save(output_path)
        return output_path


