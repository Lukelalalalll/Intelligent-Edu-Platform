"""Pure utility functions extracted from PPTCreator.

These functions have no dependency on PPTCreator instance state and can be
used by any theme creator or tested independently.
"""
from __future__ import annotations

import os
from typing import Any, Dict, List, Optional

from pptx.util import Pt
from PIL import Image

from .theme_catalog import resolve_base_theme


# ── Template & layout utilities ────────────────────────────────────────

def get_template_path(template_base_path: str, theme: str) -> str:
    """Resolve *theme* to a ``.pptx`` file path under *template_base_path*."""
    available_themes = [
        os.path.splitext(name)[0]
        for name in os.listdir(template_base_path)
        if name.endswith(".pptx")
    ]
    resolved_theme = resolve_base_theme(theme, available_themes)
    return os.path.join(template_base_path, f"{resolved_theme}.pptx")


def find_layout_by_name(prs, layout_name: str):
    """Find a slide layout by *layout_name* (case-insensitive fallback)."""
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    lower = layout_name.lower()
    for layout in prs.slide_layouts:
        if layout.name.lower() == lower:
            return layout
    return None


def layout_has_body(layout) -> bool:
    """Check if *layout* has a BODY (2) or OBJECT (7) placeholder."""
    for shape in layout.placeholders:
        if shape.placeholder_format.type in (2, 7):
            return True
    return False


def find_content_layout(prs):
    """Return the first slide layout that contains a body/object placeholder."""
    for layout in prs.slide_layouts:
        if layout_has_body(layout):
            return layout
    return None


def clear_existing_slides(prs) -> None:
    """Remove all slides from *prs* while keeping masters/layouts intact."""
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(slide_id)


def is_meaningful_chart_type(chart_type) -> bool:
    """Return True only when *chart_type* represents a real chart/diagram request."""
    if not chart_type:
        return False
    lower = str(chart_type).strip().lower()
    return lower not in ("", "none", "null", "n/a", "no chart", "no_chart", "false", "0")


# ── Font sizing ────────────────────────────────────────────────────────

def determine_content_font_size(bullet_count: int, avg_words_per_bullet: float) -> Pt:
    """Simple heuristic: fewer/shorter bullets → larger font."""
    if bullet_count <= 3 and avg_words_per_bullet < 20:
        return Pt(18)
    elif bullet_count == 4 and avg_words_per_bullet < 12:
        return Pt(16)
    else:
        return Pt(14)


# ── Image insertion helpers ────────────────────────────────────────────

def insert_picture_into_placeholder(slide, placeholder, image_path: str) -> None:
    """Insert *image_path* into *placeholder* at its current geometry."""
    if not os.path.exists(image_path):
        print(f"Warning: Image file not found: {image_path}")
        return
    try:
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height
        slide.shapes.add_picture(image_path, left, top, width, height)
    except Exception as e:
        print(f"Warning: Failed to insert picture: {e}")


def insert_picture_with_aspect_ratio(slide, placeholder_shape, image_path: str,
                                      left, top, width, height) -> None:
    """Insert *image_path* preserving its native aspect ratio."""
    try:
        with Image.open(image_path) as img:
            original_width, original_height = img.size
        aspect_ratio = original_height / original_width
        calculated_height = int(width * aspect_ratio)
        print(
            f"📏 Image dimensions - Original: {original_width}x{original_height}, "
            f"Target: {width}x{calculated_height} (aspect ratio: {aspect_ratio:.3f})"
        )
        slide.shapes.add_picture(image_path, left, top, width, calculated_height)
    except Exception as e:
        print(f"❌ Failed to get image dimensions, using original height: {e}")
        slide.shapes.add_picture(image_path, left, top, width, height)


# ── Speaker notes ──────────────────────────────────────────────────────

def apply_speaker_notes(slide, slide_data: Dict[str, Any]) -> None:
    """Write speaker notes to *slide* if provided in *slide_data*."""
    notes_text = slide_data.get("speaker_notes") or slide_data.get("notes", "")
    if not notes_text:
        return
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = str(notes_text)
    except Exception as e:
        print(f"⚠️ [SpeakerNotes] Failed to write notes for slide '{slide_data.get('title', '')}': {e}")


# ── Visual task preparation ───────────────────────────────────────────

def prepare_image_data(slide_data: Dict[str, Any], placeholder_info: Dict[str, Any],
                       placeholder_index: int) -> Dict[str, Any]:
    """Build the image-generation request dict for *ImageChartProcessor*."""
    return {
        "title": slide_data.get("title", ""),
        "content_list": slide_data.get("content", []),
        "ratio": placeholder_info["ratio"],
        "type": placeholder_info["image_type"],
        "chart_type": slide_data.get("chart_type", ""),
        "chart_reasoning": slide_data.get("chart_reasoning", ""),
        "original_text": slide_data.get("original_text", ""),
        "placeholder_type": placeholder_info["placeholder_type"],
        "placeholder_index": placeholder_index,
        "aspect_ratio": placeholder_info["aspect_ratio"],
    }


# ── Theme dispatch helpers ─────────────────────────────────────────────

THEME_CREATOR_MAPPING: Dict[str, str] = {
    "business": "BusinessPPTCreator",
    "light": "LightPPTCreator",
    "dark": "DarkPPTCreator",
}


def should_use_specialized_creator(template_base_path: str, theme: str) -> bool:
    """Check whether *theme* has a dedicated PPTCreator subclass."""
    try:
        available_themes = [
            os.path.splitext(name)[0]
            for name in os.listdir(template_base_path)
            if name.endswith(".pptx")
        ]
        resolved_theme = resolve_base_theme(theme, available_themes)
    except Exception:
        resolved_theme = theme
    return resolved_theme.lower() in THEME_CREATOR_MAPPING
