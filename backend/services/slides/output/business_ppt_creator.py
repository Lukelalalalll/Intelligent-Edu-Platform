import os
import asyncio
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from ..business import (
    BusinessContentProcessor,
    BusinessSectionHandler,
    BusinessLayoutManager,
    BusinessTableHandler,
    BusinessPlaceholderProcessor,
    BusinessImageProcessor,
    BusinessSubtitleGenerator,
    BusinessSlideNumberHandler,
    BusinessLatexProcessor
)
from ..generation.img_chart_processor import ImageChartProcessor
from .ppt_creator import PPTCreator
import asyncio


class BusinessPPTCreator(PPTCreator):
    """Business模板专属的PPT创建器
    
    继承自PPTCreator，为Business模板提供特殊的内容处理逻辑
    """

    
    def __init__(self, template_base_path=None):
        super().__init__(template_base_path)
        self.template_name = "Business"
        # 初始化所有business模块
        self.content_processor = BusinessContentProcessor()
        self.section_handler = BusinessSectionHandler()
        self.layout_manager = BusinessLayoutManager()
        self.table_handler = BusinessTableHandler()
        self.placeholder_processor = BusinessPlaceholderProcessor()
        self.image_processor = BusinessImageProcessor()
        self.subtitle_generator = BusinessSubtitleGenerator()
        self.slide_number_handler = BusinessSlideNumberHandler()
        self.latex_processor = BusinessLatexProcessor()
        
    def _get_template_path(self, theme):
        """获取Business模板路径"""
        if theme.lower() == "business":
            return os.path.join(self.template_base_path, f"{self.template_name}.pptx")
        else:
            # 如果不是Business主题，回退到父类处理
            return super()._get_template_path(theme)
    
    def _find_layout_by_name(self, prs, layout_name):
        """根据布局名称查找对应的布局，支持动态布局"""
        # 检查是否为动态布局
        global base_layout
        if "dynamic" in layout_name.lower():
            for layout in prs.slide_layouts:
                if layout.name == layout_name:
                    base_layout = layout
                    break
            
            if base_layout:
                # 返回动态布局处理后的布局
                return self._create_dynamic_layout(base_layout, layout_name)
            else:
                print(f"Warning: Base layout '{layout_name}' not found for dynamic layout '{layout_name}'")
                return None
        else:
            # 非动态布局，使用父类方法
            return super()._find_layout_by_name(prs, layout_name)
    
    def _create_dynamic_layout(self, base_layout, dynamic_layout_name):
        """创建动态布局
        
        Args:
            base_layout: 基础布局
            dynamic_layout_name (str): 动态布局名称
            
        Returns:
            layout: 处理后的布局
        """
        # 分析基础布局中的组合形状
        group_templates = self.layout_manager.analyze_layout_groups(base_layout)
        
        # 将组合模板信息存储到布局对象中，供后续使用
        base_layout.group_templates = group_templates
        base_layout.is_dynamic = True
        
        # 注意：这里不进行隐藏操作，因为此时还不知道内容数量
        # 隐藏操作将在_process_dynamic_layout_content中根据实际内容数量进行
        
        return base_layout
    
    def _insert_picture_into_placeholder(self, slide, placeholder, image_path):
        """将图片插入到占位符中
        
        Args:
            slide: 幻灯片对象
            placeholder: 占位符对象
            image_path (str): 图片文件路径
        """
        try:
            # 🎯 直接在占位符中插入图片，保持层次顺序
            placeholder.insert_picture(image_path)
            print(f"✅ [Insert Picture] Successfully inserted image into placeholder: {image_path}")
        except Exception as e:
            print(f"⚠️ [Insert Picture] insert_picture failed, falling back to add_picture: {e}")
            # 如果插入失败，回退到原来的方法
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height
            slide.shapes.add_picture(image_path, left, top, width, height)

    def _process_dynamic_layout_content(self, slide, slide_data, layout, title_font_size, content_font_size, presentation_title=None):
        """处理动态布局的内容插入"""
        if not hasattr(layout, 'is_dynamic') or not layout.is_dynamic:
            return
        
        # 获取标题和内容
        title = slide_data.get('title', '')
        content_list = slide_data.get('content', [])
        
        content_count = len(content_list)
        group_templates = layout.group_templates
        template_count = len(group_templates)
        
        # 确定要插入的组合索引
        target_indices = self.layout_manager.get_target_indices(content_count, template_count)

        # 异步并行处理图片占位符（在收集模式下将被延迟执行）
        asyncio.run(self.image_processor.process_image_placeholders(slide, slide_data))
        
        # 处理标题（查找type=1的占位符）
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
                
            placeholder_type = shape.placeholder_format.type
            
            # 处理标题占位符
            if placeholder_type == 1:
                shape.text = title
                # 设置Business模板的标题字体大小
                if title_font_size:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = title_font_size
        
        # 根据策略插入文本框和内容
        for i, content in enumerate(content_list):
            if i < len(target_indices):
                template_index = target_indices[i]
                if template_index < len(group_templates):
                    template = group_templates[template_index]
                    
                    # 生成副标题内容
                    subtitle_content = self.subtitle_generator.generate_subtitle_content(i + 1, content)
                    
                    # 插入副标题文本框
                    if len(template['text_boxes']) > 0:
                        text_box = template['text_boxes'][0]  # 使用第一个文本框作为副标题
                        subtitle_shape = slide.shapes.add_textbox(
                            text_box['left'], text_box['top'], 
                            text_box['width'], text_box['height']
                        )
                        subtitle_shape.text = subtitle_content
                        
                        # 设置副标题样式
                        if subtitle_shape.text_frame.paragraphs:
                            paragraph = subtitle_shape.text_frame.paragraphs[0]
                            paragraph.font.size = Pt(18)
                            paragraph.font.bold = True
                            if layout.name == "Rectangular Style_dynamic":
                                paragraph.alignment = PP_ALIGN.CENTER
                    
                    # 插入正文文本框
                    if len(template['text_boxes']) > 1:
                        text_box = template['text_boxes'][1]  # 使用第二个文本框
                        body_shape = slide.shapes.add_textbox(
                            text_box['left'], text_box['top'], 
                            text_box['width'], text_box['height']
                        )
                        body_shape.text = content
                        
                        # 设置正文样式，使用传入的字体大小
                        if body_shape.text_frame.paragraphs:
                            paragraph = body_shape.text_frame.paragraphs[0]
                            paragraph.font.size = content_font_size
                            paragraph.font.bold = False
                        
                        # 启用自动换行
                        body_shape.text_frame.word_wrap = True
        
        # ===== 处理表格数据（使用通用表格处理方法）=====
        # 如果没有传入presentation_title，从slide_data中获取或使用默认值
        if presentation_title is None:
            presentation_title = slide_data.get('presentation_title', title.split(' ')[0] if title else 'Default')
        
        self.table_handler.process_tables_with_placeholders(slide, slide_data, presentation_title, self.placeholder_processor)
        # ===== 表格处理结束 =====
        
        # 处理LaTeX公式（如果存在）
        self.latex_processor.process_latex_formulas(slide, slide_data, self.placeholder_processor)
    
    def _process_business_placeholders(self, slide, slide_data, presentation_title, prs=None, is_title_slide=False):
        """Business模板专属的占位符处理逻辑"""
        # 获取标题和内容
        title = slide_data.get('title', '')
        content_list = slide_data.get('content', [])
        
        # 使用content_processor计算字体大小
        content_font_size, title_font_size = self.content_processor.get_font_sizes(content_list, title)
        
        # 检查是否为动态布局，如果是则使用动态布局处理
        layout = getattr(slide, 'custom_layout', None) or getattr(slide, 'slide_layout', None)
        if layout and hasattr(layout, 'is_dynamic') and layout.is_dynamic:
            self._process_dynamic_layout_content(slide, slide_data, layout, title_font_size, content_font_size, presentation_title)
            return
        
        
        # 使用专门的占位符处理器处理标题占位符
        self.placeholder_processor.process_title_placeholders(slide, slide_data, title_font_size)
        
        # 收集内容占位符 (type=2)
        type2_placeholders = self.placeholder_processor.collect_content_placeholders(slide)
        
        # 处理图片和对象占位符 (使用专门的图片处理器)
        asyncio.run(self.image_processor.process_image_placeholders(slide, slide_data))
        
        # 使用专门的占位符处理器进行智能内容分配
        self.placeholder_processor.apply_smart_content_distribution(type2_placeholders, content_list, content_font_size)
        
        # Business模板的页码处理
        self.slide_number_handler.add_slide_number(slide, slide_data, is_title_slide)
        
        # 处理表格数据（使用通用表格处理方法）
        self.table_handler.process_tables_with_placeholders(slide, slide_data, presentation_title, self.placeholder_processor)
        
        # 处理LaTeX公式（如果存在）
        self.latex_processor.process_latex_formulas(slide, slide_data, self.placeholder_processor)
    
    def create_presentation(self, ppt_schema, output_path):
        """创建Business模板演示文稿
        
        Args:
            ppt_schema (dict): PPT结构数据
            output_path (str): 输出文件路径
        """
        if not ppt_schema or 'slides' not in ppt_schema:
            raise ValueError("Invalid PPT schema")
            
        # 获取主题信息
        theme = ppt_schema.get('theme', 'Business')
        template_path = self._get_template_path(theme)
        
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Business template not found: {template_path}")
            
        # 创建演示文稿
        prs = Presentation(template_path)
        
        # 获取演示文稿标题和幻灯片数据
        presentation_title = ppt_schema.get('presentation_title', '')
        slides_data = ppt_schema['slides']
        
        # 1. 创建标题幻灯片
        title_layout = self._find_layout_by_name(prs, 'Title')
        if title_layout:
            title_slide = prs.slides.add_slide(title_layout)
            title_data = {
                'title': presentation_title,
                'metadata': ppt_schema.get('metadata', {})
            }
            self._process_business_placeholders(title_slide, title_data, presentation_title, prs, True)
        
        # 2. 选择主要章节标题
        main_headers = self.section_handler.select_main_headers(slides_data)
        print(f"Business Template Processing:")
        print(f"  - Total slides: {len(slides_data)}")
        print(f"  - Selected main headers: {main_headers}")
        print(f"  - All slide titles: {[slide['title'] for slide in slides_data]}")
        
        # 3. 创建目录页
        catalogue_slide = self.section_handler.create_catalogue_slide(prs, main_headers, self._find_layout_by_name)
        if catalogue_slide:
            print(f"  - Created catalogue slide")
        else:
            print(f"  - Failed to create catalogue slide")
        
        # 4. 处理每个幻灯片，在主要章节前插入Section页
        slide_number = 3  # 从第3页开始（Title + Catalogue + 内容页）
        
        # 启动图片占位符收集模式
        print("🔄 [Batch Processing] Starting batch collection mode for image placeholders...")
        self.image_processor.start_collecting()
        
        for slide_index, slide_data in enumerate(slides_data):
            # 检查是否需要在此幻灯片前插入Section页
            # 使用带编号的标题列表进行匹配，因为slide_data['title']是带编号的
            if slide_data['title'] in self.section_handler.main_headers_with_numbers:
                # 找到下一个主要章节标题
                next_main_header = None
                for i in range(slide_index + 1, len(slides_data)):
                    if slides_data[i]['title'] in self.section_handler.main_headers_with_numbers:
                        next_main_header = slides_data[i]['title']
                        break
                
                # 创建Section页
                section_data = self.section_handler.get_section_content(
                    slide_data['title'], 
                    next_main_header, 
                    slides_data
                )
                if section_data:
                    section_slide = self.section_handler.create_section_slide(prs, section_data, self._find_layout_by_name)
                    if section_slide:
                        print(f"  - Created section slide for: {slide_data['title']}")
                        slide_number += 1
                    else:
                        print(f"  - Failed to create section slide for: {slide_data['title']}")
            
            # 创建内容幻灯片
            if 'layout' not in slide_data:
                continue
                
            layout_name = slide_data['layout'].get('name')
            if not layout_name:
                continue
                
            # 查找对应的布局
            layout = self._find_layout_by_name(prs, layout_name)
            if not layout:
                print(f"Warning: Layout '{layout_name}' not found in Business template")
                continue
                
            # 如果是动态布局，在创建幻灯片之前先修改布局
            if hasattr(layout, 'is_dynamic') and layout.is_dynamic:
                content_list = slide_data.get('content', [])
                content_count = len(content_list)
                group_templates = layout.group_templates
                template_count = len(group_templates)
                
                # 确定要插入的组合索引
                target_indices = self.layout_manager.get_target_indices(content_count, template_count)
                
                # 先修改layout，隐藏不需要的组合
                self.layout_manager.modify_layout_hide_groups(layout, target_indices)
            
            # 创建幻灯片
            slide = prs.slides.add_slide(layout)
            
            # 确保幻灯片能够访问布局信息（特别是动态布局）
            slide.custom_layout = layout
            
            # 添加页码信息到slide_data
            slide_data['slide_number'] = str(slide_number)
            slide_number += 1
            
            # 处理占位符（使用Business模板专属的处理逻辑）
            self._process_business_placeholders(slide, slide_data, presentation_title, prs)

            # 写入讲者备注（如果有）
            self._apply_speaker_notes(slide, slide_data)
        
        # 停止收集模式并执行批量处理
        self.image_processor.stop_collecting()
        print("⚡ [Batch Processing] Executing batch processing for all collected image placeholders...")
        asyncio.run(self.image_processor.process_all_collected_tasks())
        
        # 5. 在所有内容幻灯片创建完成后，插入Ending幻灯片
        ending_layout = self._find_layout_by_name(prs, 'Ending')
        if ending_layout:
            prs.slides.add_slide(ending_layout)
        else:
            print(f"Warning: Ending layout not found in Business template")
        
        # 保存演示文稿
        prs.save(output_path)
        return output_path

    