"""Dark 模板专属 PPT 创建器

继承 PPTCreator，为 Dark.pptx 提供：
1. 章节分隔页（'Section' 布局 — type 1 标题 + type 2 副标题列表）
2. 结尾页（'Ending' 布局 — type 1 标题）
3. 正确填充 BODY(2) 文字内容，跳过 DATE/FOOTER/SLIDE_NUMBER
4. PICTURE(18) / OBJECT(7) 图片占位符：仅当 chart_type 有值时触发 AI 图片

Dark 模板布局（来自模板实际扫描）：
  [0] 'Title'                          → 标题页
  [1] 'Single Picture'                 → 单图+文字
  [2] 'Picture & Diagram_Vertical'     → 左图右文+图表
  [3] 'Picture & Diagram_Horizontal'   → 上图下文+图表
  [4] 'Picture & Diagram_Vertical_reverse' → 左文右图+图表
  [5] '2 Pictures & Diagram'           → 双图+图表
  [6] 'Section'                        → 章节分隔页 (新增)
  [7] 'Ending'                         → 结尾页 (新增)
"""

import asyncio
import os

from pptx import Presentation

from .ppt_creator import PPTCreator
from .text_layout_engine import clean_bullets
from ..dark import (
    DarkContentProcessor,
    DarkSectionHandler,
    DarkPlaceholderProcessor,
)

_END_SLIDE_TITLE = 'Thank You'


class DarkPPTCreator(PPTCreator):
    """Dark 模板专属 PPT 创建器"""

    def __init__(self, template_base_path=None):
        super().__init__(template_base_path)
        self.template_name = 'Dark'
        self.content_processor = DarkContentProcessor()
        self.section_handler = DarkSectionHandler()
        self.placeholder_processor = DarkPlaceholderProcessor()

    # ------------------------------------------------------------------
    # 模板路径
    # ------------------------------------------------------------------

    def _get_template_path(self, theme):
        if theme.lower() == 'dark':
            return os.path.join(self.template_base_path, 'Dark.pptx')
        return super()._get_template_path(theme)

    # ------------------------------------------------------------------
    # 主入口
    # ------------------------------------------------------------------

    def create_presentation(self, ppt_schema: dict, output_path: str) -> str:
        """创建 Dark 主题演示文稿

        生成结构：Title → [Section页 → 内容页...] × N → Ending
        """
        if not ppt_schema or 'slides' not in ppt_schema:
            raise ValueError('Invalid PPT schema')

        theme = ppt_schema.get('theme', 'Dark')
        template_path = self._get_template_path(theme)

        if not os.path.exists(template_path):
            raise FileNotFoundError(f'Dark template not found: {template_path}')

        prs = Presentation(template_path)
        self._clear_existing_slides(prs)

        presentation_title = ppt_schema.get('presentation_title', '')
        slides_data: list[dict] = ppt_schema['slides']

        # ── 1. 标题页 ──────────────────────────────────────────────────
        self._add_title_slide(prs, presentation_title, ppt_schema.get('metadata', {}))

        # ── 2. 选出主章节 ──────────────────────────────────────────────
        main_headers = self.section_handler.select_main_headers(slides_data)
        print(f'[Dark] {len(slides_data)} content slides, '
              f'{len(main_headers)} main sections: {main_headers}')

        # ── 3. 批量收集图片任务 ────────────────────────────────────────
        print('🔄 [Dark Batch] Starting image collection mode...')
        self.start_collecting()

        slide_number = 2
        failed_slides: list[dict] = []

        for slide_index, slide_data in enumerate(slides_data):
            try:
                # ── 3a. 需要时插入章节分隔页 ──────────────────────────
                if slide_data['title'] in self.section_handler.main_headers_with_numbers:
                    subtitle = self.section_handler.get_section_subtitle(
                        slide_data['title'], slides_data
                    )
                    section_slide = self.section_handler.create_section_slide(
                        prs,
                        {'title': slide_data['title'], 'subtitle': subtitle},
                        self._find_layout_by_name,
                    )
                    if section_slide:
                        slide_number += 1
                        print(f'  [Dark] Section slide → {slide_data["title"]}')

                # ── 3b. 内容页 ─────────────────────────────────────────
                slide_data['slide_number'] = str(slide_number)
                self._add_content_slide(prs, slide_data, presentation_title)
                slide_number += 1

            except Exception as exc:
                title = slide_data.get('title', f'Slide {slide_index + 1}')
                print(f'❌ [Dark] Slide {slide_number} "{title}" failed: {exc}')
                failed_slides.append({
                    'slide_number': slide_number,
                    'title': title,
                    'error': str(exc),
                })

        # ── 4. 结尾页 ──────────────────────────────────────────────────
        end_slide = self.section_handler.create_end_slide(
            prs, self._find_layout_by_name,
            title_text=_END_SLIDE_TITLE,
            subtitle_text=presentation_title,
        )
        if end_slide:
            print('[Dark] Ending slide added.')

        # ── 5. 批量处理 AI 图片 ────────────────────────────────────────
        self.stop_collecting()
        print('⚡ [Dark Batch] Processing collected image placeholders...')
        try:
            asyncio.run(self.process_all_collected_tasks())
        except Exception as exc:
            print(f'⚠️ [Dark Batch] Image batch processing failed: {exc}')

        if failed_slides:
            print(f'⚠️ [Dark] {len(failed_slides)} slides had errors: '
                  f'{[s["title"] for s in failed_slides]}')

        prs.save(output_path)
        return output_path

    # ------------------------------------------------------------------
    # 标题页
    # ------------------------------------------------------------------

    def _add_title_slide(self, prs, presentation_title: str, metadata: dict):
        """添加 Dark 标题页（布局：'Title'，type 3 + type 4）"""
        layout = self._find_layout_by_name(prs, 'Title')
        if layout is None and prs.slide_layouts:
            layout = prs.slide_layouts[0]
        if layout is None:
            return

        title_slide = prs.slides.add_slide(layout)
        title_data = {'title': presentation_title, 'metadata': metadata, 'content': []}
        _, title_font_size = self.content_processor.get_font_sizes([], presentation_title)
        self.placeholder_processor.process_title_placeholders(title_slide, title_data, title_font_size)

    # ------------------------------------------------------------------
    # 内容页
    # ------------------------------------------------------------------

    def _add_content_slide(self, prs, slide_data: dict, presentation_title: str):
        """添加单张内容幻灯片"""
        layout_raw = slide_data.get('layout')
        if isinstance(layout_raw, dict):
            layout_name = (layout_raw.get('name') or '').strip()
        elif layout_raw is None:
            layout_name = ''
        else:
            layout_name = str(layout_raw).strip()

        content_list = clean_bullets(slide_data.get('content', []))
        slide_data = {**slide_data, 'content': content_list}

        # ── 布局选择 ──────────────────────────────────────────────────
        layout = self._find_layout_by_name(prs, layout_name) if layout_name else None

        # 如果 AI 指定布局找不到，或布局没有正文占位符，换默认单图布局
        if layout is None or (content_list and not self._layout_has_body(layout)):
            fallback = self._find_content_layout(prs)
            if fallback:
                if layout is not None:
                    print(f'⚠️ [Dark] Layout "{layout_name}" has no body — '
                          f'falling back to "{fallback.name}"')
                layout = fallback

        if layout is None:
            print(f'⚠️ [Dark] No suitable layout for "{slide_data.get("title", "")}", skipping')
            return

        # ── 创建幻灯片 ────────────────────────────────────────────────
        slide = prs.slides.add_slide(layout)

        content_font_size, title_font_size = self.content_processor.get_font_sizes(
            content_list, slide_data.get('title', '')
        )

        # 填标题
        self.placeholder_processor.process_title_placeholders(slide, slide_data, title_font_size)

        # 填内容文字
        self.placeholder_processor.process_content(slide, slide_data, content_font_size)

        # 处理 AI 图片/图表占位符（进入收集队列）
        # Bug1/2 fix: 使用 _collect_visual_tasks，不触碰已写好的文字占位符
        self._collect_visual_tasks(slide, slide_data)

        # 讲者备注
        self._apply_speaker_notes(slide, slide_data)
