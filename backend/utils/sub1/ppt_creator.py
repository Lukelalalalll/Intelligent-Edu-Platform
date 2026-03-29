import os
import asyncio
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from PIL import Image
from .img_chart_processor import ImageChartProcessor
from .business.table_handler import BusinessTableHandler
from .latex_generator import process_slide_latex


class PPTCreator:
    def __init__(self, template_base_path=None):
        self.template_base_path = template_base_path or "static/ppt_templates"
        self.image_processor = ImageChartProcessor()
        
        # 收集队列相关属性
        self.collected_tasks = []  # 收集的任务队列
        self.is_collecting = False  # 是否正在收集模式
        self.batch_results = []  # 批量处理结果
        
    def start_collecting(self):
        """开始收集模式"""
        self.collected_tasks = []
        self.is_collecting = True
        self.batch_results = []
        print("🔄 [Batch Processing] Started collecting image placeholder tasks...")
        
    def stop_collecting(self):
        """停止收集模式"""
        self.is_collecting = False
        print(f"⏹️ [Batch Processing] Stopped collecting. Total collected tasks: {len(self.collected_tasks)}")
        
    async def process_all_collected_tasks(self):
        """批量处理所有收集的任务"""
        if not self.collected_tasks:
            print("ℹ️ [Batch Processing] No collected tasks to process")
            return
            
        print(f"🚀 [Batch Processing] Starting batch processing of {len(self.collected_tasks)} collected tasks...")
        
        # 准备所有图片数据
        all_image_data = []
        task_mappings = []  # 用于映射结果回对应的任务和占位符
        
        for task_idx, task_info in enumerate(self.collected_tasks):
            slide_data = task_info['slide_data']
            placeholder_infos = task_info['placeholder_infos']
            slide_title = task_info['slide_title']
            
            # 为每个任务的每个占位符创建图片数据
            for placeholder_idx, placeholder_info in enumerate(placeholder_infos):
                image_data = self._prepare_image_data(slide_data, placeholder_info, placeholder_idx)
                all_image_data.append(image_data)
                
                # 记录映射关系
                task_mappings.append({
                    'task_idx': task_idx,
                    'placeholder_idx': placeholder_idx,
                    'slide_title': slide_title
                })
        
        # 批量并行处理所有图片数据
        total_count = len(all_image_data)
        print(f"⚡ [Batch Processing] Processing {total_count} image placeholders in parallel...")
        
        self.batch_results = await self.image_processor.process_multiple_images_async(all_image_data)
        
        print(f"✅ [Batch Processing] Batch processing completed! Generated {len(self.batch_results)} images")
        
        # 应用结果到对应的占位符
        await self._apply_batch_results(task_mappings)
        
    async def _apply_batch_results(self, task_mappings):
        """应用批量处理结果到对应的占位符"""
        print(f"📌 [Batch Processing] Applying {len(self.batch_results)} results to placeholders...")
        
        applied_count = 0
        error_count = 0
        
        for result_idx, (image_path, mapping) in enumerate(zip(self.batch_results, task_mappings)):
            try:
                task_info = self.collected_tasks[mapping['task_idx']]
                slide = task_info['slide']
                placeholder_info = task_info['placeholder_infos'][mapping['placeholder_idx']]
                slide_title = mapping['slide_title']
                
                if image_path:
                    self._insert_picture_into_placeholder(slide, placeholder_info['shape'], image_path)
                    applied_count += 1
                    print(f"✅ [Batch Processing] Applied result {result_idx+1}/{len(self.batch_results)} to slide: {slide_title}")
                else:
                    print(f"⚠️ [Batch Processing] No image generated for result {result_idx+1} in slide: {slide_title}")
                    
            except Exception as e:
                error_count += 1
                print(f"❌ [Batch Processing] Error applying result {result_idx+1}: {e}")
        
        print(f"🎉 [Batch Processing] Batch application completed! Applied: {applied_count}, Errors: {error_count}")
        
    def _prepare_image_data(self, slide_data, placeholder_info, placeholder_index):
        """准备单个占位符的图片数据
        
        Args:
            slide_data (dict): 幻灯片数据
            placeholder_info (dict): 占位符信息
            placeholder_index (int): 占位符索引（用于结果映射）
            
        Returns:
            dict: 图片数据字典
        """
        title = slide_data.get('title', '')
        content_list = slide_data.get('content', [])
        chart_type = slide_data.get('chart_type', '')
        chart_reasoning = slide_data.get('chart_reasoning', '')
        original_text = slide_data.get('original_text', '')
        
        image_data = {
            'title': title,
            'content_list': content_list,
            'ratio': placeholder_info['ratio'],
            'type': placeholder_info['image_type'],
            'chart_type': chart_type,
            'chart_reasoning': chart_reasoning,
            'original_text': original_text,
            'placeholder_type': placeholder_info['placeholder_type'],  # 占位符类型信息
            'placeholder_index': placeholder_index,                    # 用于结果映射
            'aspect_ratio': placeholder_info['aspect_ratio']           # 宽高比信息
        }
        
        return image_data

    def _get_template_path(self, theme):
        """获取主题模板路径"""
        return os.path.join(self.template_base_path, f"{theme}.pptx")
    
    def _find_layout_by_name(self, prs, layout_name):
        """根据布局名称查找对应的布局"""
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout
        return None
    
    def _get_template_creator_mapping(self):
        """获取模板创建器映射
        
        Returns:
            dict: 模板名称到创建器类的映射
        """
        return {
            'business': 'BusinessPPTCreator',
            # 未来可以在这里添加更多模板
            # 'academic': 'AcademicPPTCreator',
            # 'creative': 'CreativePPTCreator',
            # 'minimal': 'MinimalPPTCreator',
        }
    
    def _should_use_specialized_creator(self, theme):
        """判断是否应该使用专门的模板创建器
        
        Args:
            theme (str): 主题名称
            
        Returns:
            bool: 是否使用专门的创建器
        """
        creator_mapping = self._get_template_creator_mapping()
        return theme.lower() in creator_mapping
    
    def _get_specialized_creator(self, theme):
        """获取专门的模板创建器实例
        
        Args:
            theme (str): 主题名称
            
        Returns:
            PPTCreator: 专门的模板创建器实例，如果不存在则返回None
        """
        creator_mapping = self._get_template_creator_mapping()
        creator_class_name = creator_mapping.get(theme.lower())
        
        if not creator_class_name:
            return None
            
        try:
            # 动态导入创建器类
            module_name = f"utils.{theme.lower()}_ppt_creator"
            module = __import__(module_name, fromlist=[creator_class_name])
            creator_class = getattr(module, creator_class_name)
            
            # 创建实例
            creator_instance = creator_class(self.template_base_path)
            print(f"Using {creator_class_name} for theme: {theme}")
            return creator_instance
            
        except (ImportError, AttributeError) as e:
            print(f"Warning: {creator_class_name} not found for theme '{theme}', falling back to default creator. Error: {e}")
            return None
    
    def _read_table_csv(self, table_index, presentation_title):
        """读取表格CSV文件
        
        Args:
            table_index (int): 表格索引
            presentation_title (str): 演示文稿标题
            
        Returns:
            dict: 包含表头和行数据的字典
        """
        csv_path = f"md/csv/{presentation_title}_tables_{table_index}.csv"
        if not os.path.exists(csv_path):
            return None
            
        try:
            df = pd.read_csv(csv_path)
            
            # 清理表头：将 'Unnamed: index' 替换为空字符串
            headers = []
            for col in df.columns:
                if 'Unnamed:' in str(col):
                    headers.append('')
                else:
                    headers.append(str(col))
            
            # 清理数据：将 NaN 替换为空字符串
            rows = []
            for _, row in df.iterrows():
                cleaned_row = []
                for value in row:
                    if pd.isna(value) or str(value).lower() == 'nan':
                        cleaned_row.append('')
                    else:
                        cleaned_row.append(str(value))
                rows.append(cleaned_row)
            
            return {
                "header": headers,
                "rows": rows
            }
        except Exception as e:
            print(f"Error reading table CSV: {e}")
            return None

    def _create_table(self, slide, table_data, left, top, width, height):
        """在幻灯片上创建表格
        
        Args:
            slide: 幻灯片对象
            table_data (dict): 表格数据
            left: 左边距
            top: 上边距
            width: 表格宽度
            height: 表格高度
        """
        rows = len(table_data["rows"]) + 1  # +1 for header
        cols = len(table_data["header"])
        
        # 创建表格
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Seaborn风格的配色方案
        # 表头背景色 - seaborn的深蓝色
        header_bg_color = RGBColor(31, 119, 180)  # #1f77b4
        # 交替行背景色 - seaborn的浅灰色
        even_row_bg_color = RGBColor(248, 248, 248)  # #f8f8f8
        odd_row_bg_color = RGBColor(255, 255, 255)   # #ffffff
        # 表头文字颜色 - 白色
        header_text_color = RGBColor(255, 255, 255)
        # 数据文字颜色 - 深灰色
        data_text_color = RGBColor(51, 51, 51)  # #333333
        
        # 设置表头
        for col_idx, header in enumerate(table_data["header"]):
            cell = table.cell(0, col_idx)
            self._set_cell_content_with_linebreaks(cell, str(header))
            
            # 设置表头背景色
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg_color
            
            # 设置表头样式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = header_text_color
        
        # 填充数据
        for row_idx, row_data in enumerate(table_data["rows"]):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                self._set_cell_content_with_linebreaks(cell, str(cell_data))
                
                # 设置交替行背景色
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = even_row_bg_color
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = odd_row_bg_color
                
                # 设置单元格样式
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = data_text_color

    def _set_cell_content_with_linebreaks(self, cell, content):
        """设置单元格内容，正确处理<br>标签为换行
        
        Args:
            cell: 表格单元格对象
            content (str): 单元格内容
        """
        # 处理<br>标签
        if '<br>' in content:
            # 分割内容为多行
            lines = content.split('<br>')
            
            # 清空单元格内容
            cell.text_frame.clear()
            
            # 添加第一行
            if lines:
                first_paragraph = cell.text_frame.paragraphs[0]
                first_paragraph.text = lines[0].strip()
                
                # 添加后续行
                for line in lines[1:]:
                    new_paragraph = cell.text_frame.add_paragraph()
                    new_paragraph.text = line.strip()
        else:
            # 没有<br>标签，直接设置文本
            cell.text = content
    
    def _determine_content_font_size(self, bullet_count, avg_words_per_bullet):
        """根据bullet points数量和平均词数确定字体大小
        
        Args:
            bullet_count (int): bullet points数量
            avg_words_per_bullet (float): 每个bullet point的平均词数
            
        Returns:
            Pt: 字体大小
        """
        # 条件1: bullet point数在3个及以下，并且每个bullet point平均词数在20以下
        if bullet_count <= 3 and avg_words_per_bullet < 20:
            return Pt(18)
        
        # 条件2: bullet point数在4-5个，并且每个bullet point平均词数在15以下
        elif bullet_count == 4 and avg_words_per_bullet < 12:
            return Pt(16)
        
        # 其他情况
        else:
            return Pt(14)
    
    def _insert_picture_into_placeholder(self, slide, placeholder, image_path):
        """将图片插入到占位符中
        
        Args:
            slide: 幻灯片对象
            placeholder: 占位符对象
            image_path (str): 图片文件路径
        """
        if not os.path.exists(image_path):
            print(f"Warning: Image file not found: {image_path}")
            return
            
        try:
            # 获取占位符的位置和尺寸
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height
            
            # 添加图片
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            print(f"Warning: Failed to insert picture: {e}")
    
    def _process_placeholders(self, slide, slide_data, presentation_title, prs=None, is_title_slide=False):
        """处理幻灯片占位符"""
        
        # 计算标题词数
        title = slide_data.get('title', '')
        title_word_count = len(title.split()) if title else 0
        
        # 计算bullet points数量和平均词数
        content_list = slide_data.get('content', [])
        bullet_count = len(content_list)
        total_words = sum(len(content.split()) for content in content_list)
        avg_words_per_bullet = total_words / bullet_count if bullet_count > 0 else 0
        
        # 确定字体大小
        title_font_size = Pt(24) if title_word_count > 4 else None
        content_font_size = self._determine_content_font_size(bullet_count, avg_words_per_bullet)
        
        # 调试信息
        print(f"Slide '{title}':")
        print(f"  - Title words: {title_word_count} (font size: {title_font_size.pt if title_font_size else 'default'})")
        print(f"  - Bullet count: {bullet_count}, avg words: {avg_words_per_bullet:.1f} (font size: {content_font_size.pt})")
        
        # 收集图片和对象占位符信息
        placeholder_infos = []
        image_data_list = []
        
        # 准备图片处理数据
        chart_type = slide_data.get('chart_type', '')
        chart_reasoning = slide_data.get('chart_reasoning', [])
        original_text = slide_data.get('original_text', '')
        
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
                
            placeholder_type = shape.placeholder_format.type
            
            # 处理大标题 (type=3)
            if placeholder_type == 3:
                shape.text = slide_data.get('title', '')
                
            # 处理副标题 (type=4)
            elif placeholder_type == 4:
                if isinstance(slide_data.get('metadata'), dict):
                    metadata_text = []
                    if slide_data['metadata'].get('author'):
                        metadata_text.append(f"Author: {slide_data['metadata']['author']}")
                    if slide_data['metadata'].get('date'):
                        metadata_text.append(f"Date: {slide_data['metadata']['date']}")
                    if slide_data['metadata'].get('description'):
                        metadata_text.append(f"Description: {slide_data['metadata']['description']}")
                    shape.text = "\n".join(metadata_text)
                else:
                    shape.text = slide_data.get('title', '')

            # 处理正文标题 (type=1)
            elif placeholder_type == 1:
                shape.text = slide_data.get('title', '')
                # 设置标题字体大小
                if title_font_size:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = title_font_size

            elif placeholder_type == 2:
                text_frame = shape.text_frame
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                text_frame.word_wrap = True

                if content_list:
                    default_font_name = None
                    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                        default_font_name = text_frame.paragraphs[0].runs[0].font.name

                    text_frame.clear()

                    for i, content in enumerate(content_list):
                        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                        p.text = content
                        p.level = 0
                        p.alignment = PP_ALIGN.LEFT

                        # 恢复预设字体，但不锁定字号（让 AutoFit 发挥作用）
                        if default_font_name and p.runs:
                            for run in p.runs:
                                run.font.name = default_font_name

                # text_frame.clear()
                #
                # if content_list:
                #     # 设置第一个段落
                #     p = text_frame.paragraphs[0]
                #     p.text = content_list[0]
                #     p.alignment = PP_ALIGN.LEFT
                #     p.level = 1
                #     p.space_after = Pt(12)  # 设置段落间距
                #     p.font.size = content_font_size
                #
                #     # 添加后续段落
                #     for content in content_list[1:]:
                #         p = text_frame.add_paragraph()
                #         p.text = content
                #         p.alignment = PP_ALIGN.LEFT
                #         p.level = 1
                #         p.space_after = Pt(12)  # 设置段落间距
                #         p.font.size = content_font_size
            
            # 收集图片占位符 (type=18)
            elif placeholder_type == 18:
                aspect_ratio = shape.width / shape.height
                if abs(aspect_ratio - 1.778) < 0.1:  # 16:9
                    ratio = 1
                elif abs(aspect_ratio - 1.333) < 0.1:  # 4:3
                    ratio = 0
                else:
                    continue  # 不支持的比例
                
                placeholder_info = {
                    'shape': shape,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'placeholder_type': placeholder_type,
                    'aspect_ratio': aspect_ratio,
                    'ratio': ratio,
                    'image_type': 'image'
                }
                placeholder_infos.append(placeholder_info)
                
                # 准备图片数据
                image_data = {
                    'title': title,
                    'content_list': content_list,
                    'ratio': ratio,
                    'type': 'image',
                    'chart_type': chart_type,
                    'chart_reasoning': chart_reasoning,
                    'original_text': original_text,
                    'placeholder_index': len(image_data_list)
                }
                image_data_list.append(image_data)
            
            # 收集对象占位符 (type=7)
            elif placeholder_type == 7:
                aspect_ratio = shape.width / shape.height
                if abs(aspect_ratio - 1.778) < 0.1:  # 16:9
                    ratio = 1
                elif abs(aspect_ratio - 1.333) < 0.1:  # 4:3
                    ratio = 0
                else:
                    continue  # 不支持的比例
                
                placeholder_info = {
                    'shape': shape,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'placeholder_type': placeholder_type,
                    'aspect_ratio': aspect_ratio,
                    'ratio': ratio,
                    'image_type': 'diagram'
                }
                placeholder_infos.append(placeholder_info)
                
                # 准备图片数据
                image_data = {
                    'title': title,
                    'content_list': content_list,
                    'ratio': ratio,
                    'type': 'diagram',
                    'chart_type': chart_type,
                    'chart_reasoning': chart_reasoning,
                    'original_text': original_text,
                    'placeholder_index': len(image_data_list)
                }
                image_data_list.append(image_data)
        
        # 处理图片占位符
        if placeholder_infos:
            slide_title = slide_data.get('title', 'Unknown')
            print(f"🔍 [Image Processing] Found {len(placeholder_infos)} placeholders in slide: {slide_title}")
            
            # 根据是否在收集模式决定处理方式
            if self.is_collecting:
                # 收集模式：将任务加入队列
                task_info = {
                    'slide': slide,
                    'slide_data': slide_data,
                    'placeholder_infos': placeholder_infos,
                    'slide_title': slide_title
                }
                self.collected_tasks.append(task_info)
                print(f"📦 [Batch Processing] Collected task for slide: {slide_title} ({len(placeholder_infos)} placeholders)")
            else:
                # 正常模式：立即处理
                print(f"🚀 [Image Processing] Processing {len(placeholder_infos)} placeholders in parallel for slide: {slide_title}")
                image_paths = asyncio.run(self.image_processor.process_multiple_images_async(image_data_list))
                
                # 应用结果到对应占位符
                for i, (placeholder_info, image_path) in enumerate(zip(placeholder_infos, image_paths)):
                    try:
                        if image_path:
                            self._insert_picture_into_placeholder(slide, placeholder_info['shape'], image_path)
                            print(f"✅ [Image Processing] Applied image {i+1}/{len(image_paths)} to placeholder successfully")
                        else:
                            print(f"⚠️ [Image Processing] No image generated for placeholder {i+1}")
                    except Exception as e:
                        print(f"❌ [Image Processing] Error applying image {i+1} to placeholder: {e}")
                
                print(f"🎉 [Image Processing] Completed processing {len(placeholder_infos)} placeholders for slide: {slide_title}")
        
        # 如果不是标题页，且有slide_number数据，则在右下角添加幻灯片编号
        if not is_title_slide and slide_data.get('slide_number'):      
            # 在右下角添加页码文本框
            left = Inches(9.40)
            top = Inches(6.95)
            width = Inches(3)
            height = Inches(0.5)
            
            # 添加幻灯片编号文本框
            slide_number_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = slide_number_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = str(slide_data.get('slide_number', ''))
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.RIGHT  # 右对齐
        
        # 处理表格数据（使用通用表格处理方法）
        BusinessTableHandler.process_tables_generic(
            slide, slide_data, presentation_title, 
            self._read_table_csv, self._create_table
        )
        
        # 处理LaTeX公式（如果存在）
        self._process_latex_formulas(slide, slide_data)
    
    def _process_latex_formulas(self, slide, slide_data):
        """处理LaTeX公式
        
        Args:
            slide: 幻灯片对象
            slide_data (dict): 幻灯片数据
        """
        # 检查是否有latex字段且包含公式
        latex_formulas = slide_data.get('latex', [])
        
        if not latex_formulas or not any(formula.strip() for formula in latex_formulas):
            return  # 没有LaTeX公式，直接返回
        
        # 生成幻灯片ID
        slide_id = f"slide_{slide_data.get('slide_number', 'unknown')}"
        
        # 处理LaTeX公式，生成图像
        formula_images = process_slide_latex(slide_data, slide_id)
        
        if not formula_images:
            print(f"⚠️ Slide {slide_id} LaTeX formulas processing failed")
            return
            
        print(f"✅ Slide {slide_id} successfully processed {len(formula_images)} formulas")
        
        # 寻找合适的占位符来插入公式图片
        self._insert_latex_images(slide, formula_images)
    
    def _insert_latex_images(self, slide, formula_images):
        """将LaTeX公式图像插入到幻灯片中，参考business/latex_processor逻辑
        
        Args:
            slide: 幻灯片对象
            formula_images (dict): 公式到图像路径的映射
        """
        # 收集其他类型的占位符信息（参考business模板处理逻辑）
        other_placeholders = self._collect_other_placeholders(slide)
        
        # 为每个公式图像分配一个占位符
        for i, (formula, image_path) in enumerate(formula_images.items()):
            if i < len(other_placeholders):
                placeholder_info = other_placeholders[i]
                try:
                    # 使用改进的图片插入方法，自动计算高度
                    self._insert_picture_with_aspect_ratio(
                        slide, 
                        placeholder_info['shape'], 
                        image_path,
                        placeholder_info['left'],
                        placeholder_info['top'],
                        placeholder_info['width'],
                        placeholder_info['height']
                    )
                    print(f"✅ Formula image inserted into the slide: {formula[:50]}...")
                except Exception as e:
                    print(f"❌ Formula image failed to insert into the slide: {e}")
            else:
                print(f"⚠️ No enough placeholders to insert formula: {formula[:50]}...")
    
    def _collect_other_placeholders(self, slide):
        """收集其他类型的占位符信息，排除常见占位符类型
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            list: 其他占位符的位置信息列表
        """
        other_placeholders = []
        excluded_types = {1, 2, 3, 4, 13}  # 排除标题、内容、页码等常见占位符
        
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
                
            placeholder_type = shape.placeholder_format.type
            
            if placeholder_type not in excluded_types:
                other_placeholders.append({
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'type': placeholder_type,
                    'shape': shape
                })
        
        # 按位置排序（从上到下，从左到右）
        other_placeholders.sort(key=lambda x: (x['top'], x['left']))
        
        return other_placeholders
    
    def _insert_picture_with_aspect_ratio(self, slide, placeholder_shape, image_path, left, top, width, height):
        """将图片插入到占位符中，根据原始宽高比计算高度
        
        Args:
            slide: 幻灯片对象
            placeholder_shape: 占位符形状对象
            image_path (str): 图片文件路径
            left: 左边距
            top: 上边距
            width: 目标宽度
            height: 占位符高度（将被忽略，根据宽高比计算）
        """
        try:
            # 获取图片的原始尺寸
            with Image.open(image_path) as img:
                original_width, original_height = img.size
                
            # 计算宽高比
            aspect_ratio = original_height / original_width
            
            # 根据目标宽度和原始宽高比计算新的高度
            calculated_height = int(width * aspect_ratio)
            
            print(f"📏 Image dimensions - Original: {original_width}x{original_height}, "
                  f"Target: {width}x{calculated_height} (aspect ratio: {aspect_ratio:.3f})")
            
            # 添加图片到指定位置，使用计算出的高度
            slide.shapes.add_picture(image_path, left, top, width, calculated_height)
            
        except Exception as e:
            print(f"❌ Failed to get image dimensions, using original height: {e}")
            # 如果获取图片尺寸失败，回退到使用原始height参数
            slide.shapes.add_picture(image_path, left, top, width, height)
    
    def create_presentation(self, ppt_schema, output_path):
        """创建演示文稿
        
        Args:
            ppt_schema (dict): PPT结构数据
            output_path (str): 输出文件路径
        """
        if not ppt_schema or 'slides' not in ppt_schema:
            raise ValueError("Invalid PPT schema")
            
        # 获取主题信息
        theme = ppt_schema.get('theme', 'Dark')
        
        # 检查是否应该使用专门的模板创建器
        if self._should_use_specialized_creator(theme):
            specialized_creator = self._get_specialized_creator(theme)
            if specialized_creator:
                return specialized_creator.create_presentation(ppt_schema, output_path)
            else:
                print(f"Specialized creator not available, falling back to default creator for theme: {theme}")
        
        # 使用默认的创建逻辑
        template_path = self._get_template_path(theme)
        
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
            
        # 创建演示文稿
        prs = Presentation(template_path)
        
        # 获取演示文稿标题
        presentation_title = ppt_schema.get('presentation_title', '')
        
        # 创建标题幻灯片
        title_layout = self._find_layout_by_name(prs, 'Title')
        if title_layout:
            title_slide = prs.slides.add_slide(title_layout)
            title_data = {
                'title': presentation_title,
                'metadata': ppt_schema.get('metadata', {})
            }
            self._process_placeholders(title_slide, title_data, presentation_title, prs, True)
        
        # 处理每个幻灯片
        print("🔄 [Batch Processing] Starting batch collection mode for image placeholders...")
        self.start_collecting()
        
        failed_slides = []
        for slide_index, slide_data in enumerate(ppt_schema['slides']):
            try:
                if 'layout' not in slide_data:
                    continue
                    
                layout_name = slide_data['layout'].get('name')
                if not layout_name:
                    continue
                    
                # 查找对应的布局
                layout = self._find_layout_by_name(prs, layout_name)
                if not layout:
                    print(f"Warning: Layout '{layout_name}' not found in template")
                    continue
                    
                # 创建幻灯片
                slide = prs.slides.add_slide(layout)
                
                # 添加页码信息到slide_data
                slide_data['slide_number'] = str(slide_index + 1)  # 从第1页开始
                
                # 处理占位符
                self._process_placeholders(slide, slide_data, presentation_title, prs)
            except Exception as e:
                slide_title = slide_data.get('title', f'Slide {slide_index + 1}')
                print(f"❌ Slide {slide_index + 1} '{slide_title}' failed: {e}")
                failed_slides.append({
                    "slide_index": slide_index + 1,
                    "title": slide_title,
                    "error": str(e)
                })
                # Continue with remaining slides — don't abort the whole presentation
        
        # 停止收集模式并执行批量处理
        self.stop_collecting()
        print("⚡ [Batch Processing] Executing batch processing for all collected image placeholders...")
        try:
            asyncio.run(self.process_all_collected_tasks())
        except Exception as e:
            print(f"⚠️ Batch image processing failed: {e} — PPT saved without some images")
        
        if failed_slides:
            print(f"⚠️ {len(failed_slides)}/{len(ppt_schema['slides'])} slides had errors: {[s['title'] for s in failed_slides]}")
        
        # 保存演示文稿
        prs.save(output_path)
        return output_path


