from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

class SimpleGroupManager:
    """
    简化的组合形状管理器
    """
    
    def __init__(self, ppt_path):
        self.prs = Presentation(ppt_path)
        self.group_templates = []
        
    def analyze_layout_groups(self, layout_name="1"):
        """
        分析版式中所有group的位置和大小
        """
        # 查找版式
        layout = None
        for l in self.prs.slide_layouts:
            if l.name == layout_name:
                layout = l
                break
        
        if not layout:
            print(f"未找到版式: {layout_name}")
            return
        
        print(f"分析版式 '{layout_name}' 中的组合形状...")
        
        # 分析版式中的所有形状
        for i, shape in enumerate(layout.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_template = self._extract_group_template(shape, i+1)
                self.group_templates.append(group_template)
        
        print(f"找到 {len(self.group_templates)} 个组合形状模板")
        return self.group_templates
    
    def _extract_group_template(self, group_shape, group_index):
        """
        提取组合形状的模板信息（位置、大小）
        """
        template = {
            'index': group_index,
            'left': group_shape.left,
            'top': group_shape.top,
            'width': group_shape.width,
            'height': group_shape.height,
            'text_boxes': []
        }
        
        print(f"\n=== 分析组合 {group_index} ===")
        print(f"组合位置: ({group_shape.left}, {group_shape.top})")
        print(f"组合大小: ({group_shape.width}, {group_shape.height})")
        print(f"组合包含 {len(group_shape.shapes)} 个子形状")
        
        # 分析组合中的每个子形状
        for j, sub_shape in enumerate(group_shape.shapes):
            print(f"\n  子形状 {j+1}:")
            print(f"    类型: {sub_shape.shape_type}")
            print(f"    名称: {sub_shape.name if hasattr(sub_shape, 'name') else 'N/A'}")
            print(f"    位置: ({sub_shape.left}, {sub_shape.top})")
            print(f"    大小: ({sub_shape.width}, {sub_shape.height})")
            
            if sub_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text_box_info = {
                    'index': j + 1,
                    'left': sub_shape.left,
                    'top': sub_shape.top,
                    'width': sub_shape.width,
                    'height': sub_shape.height,
                    'shape_type': 'TEXT_BOX'
                }
                template['text_boxes'].append(text_box_info)
                print(f"    *** 识别为文本框 ***")
            else:
                print(f"    不是文本框，跳过")
        
        print(f"组合 {group_index} 包含 {len(template['text_boxes'])} 个文本框")
        return template
    
    def create_slide_with_content(self, content_data, layout_name="1"):
        """
        使用版式创建幻灯片并插入内容
        """
        # 查找版式
        layout = None
        for l in self.prs.slide_layouts:
            if l.name == layout_name:
                layout = l
                break
        
        if not layout:
            print(f"未找到版式: {layout_name}")
            return None
        
        # 根据内容数量决定插入策略
        content_count = len(content_data)
        template_count = len(self.group_templates)
        
        print(f"内容数量: {content_count}, 模板数量: {template_count}")
        
        # 确定要插入的组合索引
        target_indices = self._get_target_indices(content_count, template_count)
        print(f"目标组合索引: {target_indices}")
        
        # 先修改layout，隐藏不需要的组合
        self._modify_layout_hide_groups(layout, target_indices)
        
        # 使用修改后的layout创建幻灯片
        slide = self.prs.slides.add_slide(layout)
        print(f"使用修改后的版式 '{layout_name}' 创建新幻灯片")
        
        # 根据策略插入文本框和内容
        for i, content in enumerate(content_data):
            if i < len(target_indices):
                template_index = target_indices[i]
                if template_index < len(self.group_templates):
                    template = self.group_templates[template_index]
                    content_item = content_data[i]
                    
                    print(f"\n处理内容 {i+1} -> 组合 {template_index+1}:")
                    
                    # 插入副标题文本框
                    if len(template['text_boxes']) > 0:
                        text_box = template['text_boxes'][0]
                        subtitle_shape = slide.shapes.add_textbox(
                            text_box['left'], text_box['top'], 
                            text_box['width'], text_box['height']
                        )
                        subtitle_shape.text = content_item.get('subtitle', f'副标题{i+1}')
                        
                        # 设置副标题样式
                        if subtitle_shape.text_frame.paragraphs:
                            paragraph = subtitle_shape.text_frame.paragraphs[0]
                            paragraph.font.size = Pt(18)
                            paragraph.font.bold = True
                        
                        print(f"  插入副标题: {subtitle_shape.text}")
                    
                    # 插入正文文本框
                    if len(template['text_boxes']) > 1:
                        text_box = template['text_boxes'][1]
                        body_shape = slide.shapes.add_textbox(
                            text_box['left'], text_box['top'], 
                            text_box['width'], text_box['height']
                        )
                        body_shape.text = content_item.get('body', f'正文内容{i+1}')
                        
                        # 设置正文样式
                        if body_shape.text_frame.paragraphs:
                            paragraph = body_shape.text_frame.paragraphs[0]
                            paragraph.font.size = Pt(14)
                            paragraph.font.bold = False
                        
                        print(f"  插入正文: {body_shape.text}")
        
        return slide
    
    def _modify_layout_hide_groups(self, layout, target_indices):
        """
        修改layout，隐藏不需要的组合形状
        """
        group_count = 0
        
        for i, shape in enumerate(layout.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_count += 1
                group_index = group_count - 1  # 转换为0基索引
                
                if group_index not in target_indices:
                    print(f"在layout中隐藏组合 {group_index + 1}")
                    self._hide_group_shapes(shape)
    
    def _hide_group_shapes(self, group_shape):
        """
        隐藏组合形状中的所有子形状
        """
        for sub_shape in group_shape.shapes:
            try:
                # 隐藏填充
                sub_shape.fill.background()
                print(f"    隐藏子形状填充成功")
            except Exception as e:
                print(f"    隐藏子形状填充失败: {e}")
            
            try:
                # 隐藏轮廓
                sub_shape.line.fill.background()
                print(f"    隐藏子形状轮廓成功")
            except Exception as e:
                print(f"    隐藏子形状轮廓失败: {e}")
            
            # 如果是组合形状，递归隐藏
            if sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._hide_group_shapes(sub_shape)
    
    def _get_target_indices(self, content_count, template_count):
        """
        根据内容数量确定目标组合索引
        """
        if content_count == 1:
            # 1个元素：插入到第1个组合
            return [0]
        elif content_count == 2:
            # 2个元素：插入到第1,3组合
            return [0, 2] if template_count >= 3 else [0, 1]
        elif content_count == 3:
            # 3个元素：插入到第1,3,5组合
            if template_count >= 5:
                return [0, 2, 4]
            elif template_count >= 3:
                return [0, 2, min(3, template_count-1)]
            else:
                return [0, 1, 2] if template_count >= 3 else [0, 1]
        elif content_count == 4:
            # 4个元素：按顺序插入到前4个组合
            return list(range(min(4, template_count)))
        elif content_count == 5:
            # 5个元素：按顺序插入到前5个组合
            return list(range(min(5, template_count)))
        else:
            # 其他情况：按顺序插入
            return list(range(min(content_count, template_count)))
    
    def save(self, output_path):
        """
        保存修改后的文件
        """
        self.prs.save(output_path)
        print(f"文件已保存: {output_path}")

def main():
    # 创建组合管理器
    manager = SimpleGroupManager("test.pptx")
    
    # 分析版式中的组合形状
    templates = manager.analyze_layout_groups("1")
    
    # 准备内容数据
    content_data = [
        {
            'subtitle': 'First Group',
            'body': 'This is the first group of content\nSupport multi-line text\nSecond line content'
        },
        {
            'subtitle': 'Second Group',
            'body': 'This is the second group of content\nContains important information\nThird line content'
        },
        {
            'subtitle': 'Third Group',
            'body': 'This is the third group of content\nDetailed explanation\nAdditional information'
        },
        {
            'subtitle': 'Fourth Group',
            'body': 'This is the fourth group of content\nSummary content\nLast line'
        },
        {
            'subtitle': 'Fifth Group',
            'body': 'This is the fifth group of content\nSummary content\nLast line'
        },
        
    ]
    
    # 创建幻灯片并插入内容
    slide = manager.create_slide_with_content(content_data, "1")
    
    # 保存文件
    manager.save("test_simple_groups.pptx")

if __name__ == "__main__":
    main() 