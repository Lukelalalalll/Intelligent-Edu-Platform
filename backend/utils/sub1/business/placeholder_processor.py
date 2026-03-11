"""Business模板占位符处理器

专门处理Business模板中各种类型占位符的处理逻辑，包括：
- 标题占位符 (type=1,3,4)
- 内容占位符 (type=2)
- 智能分配逻辑
- 占位符收集和排序
"""

from pptx.util import Pt


class BusinessPlaceholderProcessor:
    """Business模板占位符处理器"""
    
    # 占位符类型常量
    PLACEHOLDER_TITLE = 1           # 正文标题
    PLACEHOLDER_CONTENT = 2         # 内容
    PLACEHOLDER_MAIN_TITLE = 3      # 大标题
    PLACEHOLDER_SUBTITLE = 4        # 副标题
    PLACEHOLDER_OBJECT = 7          # 对象
    PLACEHOLDER_SLIDE_NUMBER = 13   # 幻灯片编号
    PLACEHOLDER_IMAGE = 18          # 图片
    
    def __init__(self):
        """初始化占位符处理器"""
        pass
    
    def process_title_placeholders(self, slide, slide_data, title_font_size):
        """处理标题类占位符
        
        Args:
            slide: 幻灯片对象
            slide_data (dict): 幻灯片数据
            title_font_size: 标题字体大小
        """
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
                
            placeholder_type = shape.placeholder_format.type
            
            # 处理大标题 (type=3)
            if placeholder_type == self.PLACEHOLDER_MAIN_TITLE:
                shape.text = slide_data.get('title', '')
                if title_font_size:
                    self._set_font_size(shape, title_font_size)
            
            # 处理副标题 (type=4)
            elif placeholder_type == self.PLACEHOLDER_SUBTITLE:
                subtitle_text = self._generate_subtitle_text(slide_data)
                shape.text = subtitle_text
                if title_font_size:
                    self._set_font_size(shape, title_font_size)

            # 处理正文标题 (type=1)
            elif placeholder_type == self.PLACEHOLDER_TITLE:
                shape.text = slide_data.get('title', '')
                if title_font_size:
                    self._set_font_size(shape, title_font_size)
    
    def collect_content_placeholders(self, slide):
        """收集内容占位符 (type=2)
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            list: 按位置排序的内容占位符列表
        """
        type2_placeholders = []
        
        for shape in slide.shapes:
            if (shape.is_placeholder and 
                shape.placeholder_format.type == self.PLACEHOLDER_CONTENT):
                type2_placeholders.append(shape)
        
        # 按位置排序：先按top位置，再按left位置
        type2_placeholders.sort(key=lambda shape: (shape.top, shape.left))
        
        return type2_placeholders
    
    def collect_other_placeholders(self, slide):
        """收集其他类型的占位符信息
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            list: 其他占位符的位置信息列表
        """
        other_placeholders = []
        excluded_types = {
            self.PLACEHOLDER_TITLE,
            self.PLACEHOLDER_CONTENT,
            self.PLACEHOLDER_MAIN_TITLE,
            self.PLACEHOLDER_SUBTITLE,
            self.PLACEHOLDER_SLIDE_NUMBER
        }
        
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
                
            placeholder_type = shape.placeholder_format.type
            
            if placeholder_type not in excluded_types:
                other_placeholders.append({
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'type': placeholder_type,
                    'shape': shape
                })
        
        # 按位置排序（从上到下，从左到右）
        other_placeholders.sort(key=lambda x: (x['top'], x['left']))
        
        return other_placeholders
    
    def apply_smart_content_distribution(self, placeholders, content_list, content_font_size):
        """智能分配内容到占位符
        
        Args:
            placeholders (list): 内容占位符列表
            content_list (list): 内容列表
            content_font_size: 内容字体大小
        """
        bullet_count = len(content_list)
        
        if bullet_count in [1, 4, 5]:
            # 顺序填入
            target_indices = list(range(bullet_count))
        elif bullet_count == 2:
            # 在索引0和2的占位符填入
            target_indices = [0, 2]
        elif bullet_count == 3:
            # 在索引0, 2, 4的占位符填入
            target_indices = [0, 2, 4]
        else:
            # 默认顺序填入
            target_indices = list(range(bullet_count))
        
        for i, content in enumerate(content_list):
            if i < len(target_indices) and target_indices[i] < len(placeholders):
                shape = placeholders[target_indices[i]]
                shape.text = content
                self._set_font_size(shape, content_font_size)
    
    def _generate_subtitle_text(self, slide_data):
        """生成副标题文本
        
        Args:
            slide_data (dict): 幻灯片数据
            
        Returns:
            str: 副标题文本
        """
        if isinstance(slide_data.get('metadata'), dict):
            metadata_text = []
            metadata = slide_data['metadata']
            
            if metadata.get('author'):
                metadata_text.append(f"Author: {metadata['author']}")
            if metadata.get('date'):
                metadata_text.append(f"Date: {metadata['date']}")
            if metadata.get('description'):
                metadata_text.append(f"Description: {metadata['description']}")
            
            return "\n".join(metadata_text)
        else:
            return slide_data.get('title', '')
    
    def _set_font_size(self, shape, font_size):
        """设置形状的字体大小
        
        Args:
            shape: 形状对象
            font_size: 字体大小
        """
        if font_size:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = font_size 