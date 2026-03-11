from pptx.enum.shapes import MSO_SHAPE_TYPE

class BusinessLayoutManager:
    """Business模板的布局管理器"""
    
    def __init__(self):
        pass
        
    def analyze_layout_groups(self, layout):
        """分析版式中所有group的位置和大小
        
        Args:
            layout: 布局对象
            
        Returns:
            list: 组合形状模板列表
        """
        group_templates = []
        
        # 收集所有group形状及其shape_id
        group_shapes_with_id = []
        for i, shape in enumerate(layout.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_shapes_with_id.append((shape, shape.shape_id, i+1))
        
        # 按shape_id排序
        group_shapes_with_id.sort(key=lambda x: x[1])
        
        # 按排序后的顺序提取模板信息
        for shape, shape_id, original_index in group_shapes_with_id:
            group_template = self._extract_group_template(shape, original_index)
            group_templates.append(group_template)
        
        return group_templates
    
    def _extract_group_template(self, group_shape, group_index):
        """提取组合形状的模板信息（位置、大小）
        
        Args:
            group_shape: 组合形状对象
            group_index (int): 组合索引
            
        Returns:
            dict: 组合形状模板信息
        """
        template = {
            'index': group_index,
            'left': group_shape.left,
            'top': group_shape.top,
            'width': group_shape.width,
            'height': group_shape.height,
            'text_boxes': []
        }
        
        # 分析组合中的每个子形状
        for j, sub_shape in enumerate(group_shape.shapes):
            if sub_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text_box_info = {
                    'index': j + 1,
                    'left': sub_shape.left,
                    'top': sub_shape.top,
                    'width': sub_shape.width,
                    'height': sub_shape.height,
                    'shape_type': 'TEXT_BOX'
                }
                template['text_boxes'].append(text_box_info)
        
        return template
    
    def modify_layout_hide_groups(self, layout, target_indices):
        """修改layout，隐藏不需要的组合形状
        
        Args:
            layout: 布局对象
            target_indices (list): 目标组合索引列表
        """
        # 收集所有group形状及其shape_id
        group_shapes_with_id = []
        for i, shape in enumerate(layout.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_shapes_with_id.append((shape, shape.shape_id, i))
        
        # 按shape_id排序
        group_shapes_with_id.sort(key=lambda x: x[1])
        
        # 按排序后的顺序进行隐藏操作
        for group_index, (shape, shape_id, original_index) in enumerate(group_shapes_with_id):
            if group_index not in target_indices:
                self._hide_group_shapes(shape)
    
    def _hide_group_shapes(self, group_shape):
        """隐藏组合形状中的所有子形状
        
        Args:
            group_shape: 组合形状对象
        """
        for sub_shape in group_shape.shapes:
            try:
                # 隐藏填充
                sub_shape.fill.background()
            except Exception:
                pass
            
            try:
                # 隐藏轮廓
                sub_shape.line.fill.background()
            except Exception:
                pass
            
            # 如果是组合形状，递归隐藏
            if sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._hide_group_shapes(sub_shape)
    
    def get_target_indices(self, content_count, template_count):
        """根据内容数量确定目标组合索引
        
        Args:
            content_count (int): 内容数量
            template_count (int): 模板数量
            
        Returns:
            list: 目标组合索引列表
        """
        if content_count == 1:
            # 1个元素：插入到第1个组合
            return [0]
        elif content_count == 2:
            # 2个元素：插入到第1,3组合
            return [0, 2] if template_count >= 3 else [0, 1]
        elif content_count == 3:
            # 3个元素：插入到第1,3,5组合
            if template_count >= 5:
                return [0, 2, 4]
            elif template_count >= 3:
                return [0, 2, min(3, template_count-1)]
            else:
                return [0, 1, 2] if template_count >= 3 else [0, 1]
        elif content_count == 4:
            # 4个元素：按顺序插入到前4个组合
            return list(range(min(4, template_count)))
        elif content_count == 5:
            # 5个元素：按顺序插入到前5个组合
            return list(range(min(5, template_count)))
        else:
            # 其他情况：按顺序插入
            return list(range(min(content_count, template_count))) 