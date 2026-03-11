import re
from pptx.util import Pt

class BusinessSectionHandler:
    """Business模板的章节处理器"""
    
    def __init__(self):
        self.main_headers_with_numbers = []
    
    def select_main_headers(self, slides_data):
        """从所有幻灯片中选择5个主要章节标题
        
        Args:
            slides_data (list): 所有幻灯片数据
            
        Returns:
            list: 选定的5个主要章节标题（去掉编号）
        """
        if len(slides_data) <= 5:
            # 去掉标题号
            cleaned_titles = []
            for slide in slides_data:
                title = slide['title']
                # 使用正则表达式去掉标题号：数字+点+空格、数字+点、数字+空格等格式
                cleaned_title = re.sub(r'^\d+(?:\.\d+)*\.?\s*', '', title).strip()
                cleaned_titles.append(cleaned_title)
            return cleaned_titles
        
        # 分析标题，选择符合一般文章顺序的主要章节
        main_sections_with_numbers = []  # 保留编号的标题
        section_keywords = [
            'introduction', 'abstract', 'method', 'methodology', 'approach',
            'experiment', 'evaluation', 'result', 'analysis', 'discussion',
            'conclusion', 'related work', 'background', 'implementation'
        ]
        
        # 第一优先级：包含关键词的标题
        for slide in slides_data:
            title = slide['title']
            title_lower = title.lower()
            
            for keyword in section_keywords:
                if keyword in title_lower:
                    if title not in main_sections_with_numbers:
                        main_sections_with_numbers.append(title)
                    break
        
        # 如果还不够5个，按顺序选择
        if len(main_sections_with_numbers) < 5:
            for slide in slides_data:
                if slide['title'] not in main_sections_with_numbers:
                    main_sections_with_numbers.append(slide['title'])
                    if len(main_sections_with_numbers) >= 5:
                        break
        
        # 在最终返回时去掉标题号
        cleaned_main_sections = []
        for title in main_sections_with_numbers:
            cleaned_title = re.sub(r'^\d+(?:\.\d+)*\.?\s*', '', title).strip()
            cleaned_main_sections.append(cleaned_title)
        
        # 保存带编号的标题到实例变量，供其他方法使用
        self.main_headers_with_numbers = main_sections_with_numbers[:5]
        
        # 确保只返回前5个（去掉编号）
        return cleaned_main_sections[:5]
    
    def get_section_content(self, current_title, next_title, all_slides):
        """获取Section页面的内容
        
        Args:
            current_title (str): 当前章节标题
            next_title (str): 下一个章节标题（可能为None，表示最后一个main_section）
            all_slides (list): 所有幻灯片数据
            
        Returns:
            dict: Section页面的内容数据
        """
        # 找到当前标题的位置
        current_index = None
        next_index = None
        
        for i, slide in enumerate(all_slides):
            if slide['title'] == current_title:
                current_index = i
                break
        
        if current_index is None:
            return None
        
        # 获取中间的所有章节标题
        intermediate_sections = []
        
        if next_title is not None:
            # 如果有下一个main_section，找到它的位置
            for i, slide in enumerate(all_slides):
                if slide['title'] == next_title:
                    next_index = i
                    break
            
            # 获取当前main_section和下一个main_section之间的所有章节
            if next_index is not None:
                for i in range(current_index + 1, next_index):
                    intermediate_sections.append(all_slides[i]['title'])
        else:
            # 如果是最后一个main_section，包含后续所有章节
            for i in range(current_index + 1, len(all_slides)):
                intermediate_sections.append(all_slides[i]['title'])
        
        # 计算在主要章节中的索引（1~5范围）
        # 使用原始标题进行匹配，因为all_slides中的标题包含编号
        main_headers_with_numbers = self.main_headers_with_numbers
        section_index = None
        for i, header in enumerate(main_headers_with_numbers, 1):
            if header == current_title:
                section_index = i
                break
        
        if section_index is None:
            # 如果找不到，使用默认值
            section_index = 1
        
        return {
            'current_title': current_title,
            'intermediate_sections': intermediate_sections,
            'section_index': section_index  # 1~5范围内的索引
        }
    
    def create_catalogue_slide(self, prs, main_headers, layout_finder):
        """创建目录页
        
        Args:
            prs: 演示文稿对象
            main_headers (list): 主要章节标题列表
            layout_finder: 布局查找器对象
            
        Returns:
            slide: 创建的目录幻灯片对象，如果创建失败则返回None
        """
        catalogue_layout = layout_finder(prs, 'Catalogue')
        if catalogue_layout:
            catalogue_slide = prs.slides.add_slide(catalogue_layout)
            
            # 收集所有type=2占位符，按索引排序
            type2_placeholders = []
            
            # 处理目录页的占位符
            for shape in catalogue_slide.shapes:
                if not shape.is_placeholder:
                    continue
                    
                placeholder_type = shape.placeholder_format.type
                
                # 处理type=3占位符 - 注入"Catalogue"字样
                if placeholder_type == 3:
                    shape.text = "Catalogue"
                
                # 收集type=2占位符
                elif placeholder_type == 2:
                    type2_placeholders.append(shape)
            
            # 为每个章节标题分配一个type=2占位符（按索引顺序）
            for i, header in enumerate(main_headers):
                if i < len(type2_placeholders):
                    shape = type2_placeholders[i]
                    text_frame = shape.text_frame
                    text_frame.text = header
            
            return catalogue_slide
        return None
    
    def create_section_slide(self, prs, section_data, layout_finder):
        """创建分节页
        
        Args:
            prs: 演示文稿对象
            section_data (dict): 分节数据
            layout_finder: 布局查找器对象
            
        Returns:
            slide: 创建的分节幻灯片对象，如果创建失败则返回None
        """
        # 使用预定义的Section-{index}布局
        section_index = section_data['section_index']
        section_layout_name = f"Section-{section_index}"
        section_layout = layout_finder(prs, section_layout_name)
        
        if section_layout:
            section_slide = prs.slides.add_slide(section_layout)
            print(f"Created section slide with layout: {section_layout_name}")
            
            # 处理分节页的占位符
            for shape in section_slide.shapes:
                if not shape.is_placeholder:
                    continue
                    
                placeholder_type = shape.placeholder_format.type
                
                # 处理type=3 - 设置为当前章节标题（去掉编号）
                if placeholder_type == 3:
                    current_title = section_data['current_title']
                    cleaned_title = re.sub(r'^\d+(?:\.\d+)*\.?\s*', '', current_title).strip()
                    shape.text = cleaned_title
                
                # 处理type=4 - 设置为中间章节标题（去掉编号）
                elif placeholder_type == 4:
                    if section_data['intermediate_sections']:
                        cleaned_sections = []
                        for section in section_data['intermediate_sections']:
                            cleaned_section = re.sub(r'^\d+(?:\.\d+)*\.?\s*', '', section).strip()
                            cleaned_sections.append(cleaned_section)
                        shape.text = "\n".join(cleaned_sections)
                        
                        # 如果intermediate章节数量超过4个，设置字体大小为14
                        if len(section_data['intermediate_sections']) > 4:
                            for paragraph in shape.text_frame.paragraphs:
                                paragraph.font.size = Pt(14)
                    else:
                        shape.text = ""
            
            return section_slide
        else:
            print(f"Warning: Layout '{section_layout_name}' not found, using default Section layout")
            # 如果找不到对应的布局，回退到默认的Section布局
            section_layout = layout_finder(prs, 'Section')
            if section_layout:
                section_slide = prs.slides.add_slide(section_layout)
                return section_slide
            else:
                print(f"Error: Default Section layout not found")
                return None 