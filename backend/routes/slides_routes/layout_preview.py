"""Authenticated layout preview generation with template/layout whitelisting."""
from __future__ import annotations

import io
import logging
import os
import shutil
import subprocess
import tempfile
import time
import zipfile
from pathlib import Path
from urllib.parse import unquote

from cachetools import TTLCache
from fastapi import APIRouter, Depends, HTTPException, Query, Request
from fastapi.responses import FileResponse

from backend.config import Config
from backend.core.security import get_current_user

logger = logging.getLogger(__name__)
router = APIRouter()
public_router = APIRouter()

STATIC_IMG_ROOT = Path(os.path.dirname(Config.PPT_TEMPLATES_FOLDER)).resolve() / "img"
SOFFICE_BIN = shutil.which("soffice") or shutil.which("libreoffice")
_LAYOUT_PREVIEW_RATE_LIMIT = 20
_LAYOUT_PREVIEW_WINDOW_SECONDS = 60
_layout_preview_requests: TTLCache = TTLCache(maxsize=1024, ttl=_LAYOUT_PREVIEW_WINDOW_SECONDS)


def _assert_within(base_dir: Path, candidate: Path) -> Path:
    base = base_dir.resolve()
    resolved = candidate.resolve()
    try:
        resolved.relative_to(base)
    except ValueError as exc:
        raise ValueError("Resolved path escaped the allowed directory") from exc
    return resolved


def _client_ip(request: Request) -> str:
    forwarded = (request.headers.get("x-forwarded-for") or "").split(",")[0].strip()
    if forwarded:
        return forwarded
    return str(getattr(request.client, "host", "") or "unknown")


def _enforce_preview_rate_limit(request: Request) -> None:
    key = _client_ip(request)
    count = int(_layout_preview_requests.get(key, 0) or 0) + 1
    _layout_preview_requests[key] = count
    if count > _LAYOUT_PREVIEW_RATE_LIMIT:
        raise HTTPException(status_code=429, detail="Too many preview requests")


def _available_template_paths() -> dict[str, Path]:
    templates_root = Path(Config.PPT_TEMPLATES_FOLDER).resolve()
    if not templates_root.is_dir():
        return {}
    paths: dict[str, Path] = {}
    for file in templates_root.glob("*.pptx"):
        paths[file.stem.casefold()] = _assert_within(templates_root, file)
    return paths


def _resolve_theme_template_path(theme: str) -> tuple[str, Path]:
    requested = unquote(str(theme or "").strip())
    if not requested:
        raise ValueError("Theme is required")
    template_paths = _available_template_paths()
    match = template_paths.get(requested.casefold())
    if match is None:
        raise ValueError(f"Theme '{requested}' is not available.")
    return match.stem, match


def _get_layout_names(template_path: Path) -> list[str]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx not installed") from exc

    prs = Presentation(str(template_path))
    return [str(layout.name or "").strip() for layout in prs.slide_layouts if str(layout.name or "").strip()]


def _resolve_layout_name(template_path: Path, layout_name: str) -> str:
    requested = unquote(str(layout_name or "").strip())
    if not requested:
        raise ValueError("Layout is required")
    for candidate in _get_layout_names(template_path):
        if candidate == requested:
            return candidate
    raise ValueError(f"Layout '{requested}' is not available for theme '{template_path.stem}'.")


def _cache_path(theme: str, layout_name: str) -> Path:
    return _assert_within(STATIC_IMG_ROOT, STATIC_IMG_ROOT / theme / f"{layout_name}.png")


def _dedup_zip(src_bytes: bytes) -> bytes:
    seen: set[str] = set()
    out = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(src_bytes), "r") as zin:
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename in seen:
                    continue
                seen.add(item.filename)
                zout.writestr(item, zin.read(item.filename))
    return out.getvalue()


def _render_layout_png(theme: str, layout_name: str) -> str:
    resolved_theme, template_path = _resolve_theme_template_path(theme)
    resolved_layout_name = _resolve_layout_name(template_path, layout_name)
    cache = _cache_path(resolved_theme, resolved_layout_name)
    if cache.exists():
        return str(cache)

    if not SOFFICE_BIN:
        raise RuntimeError("LibreOffice (soffice) not found on PATH.")

    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx not installed") from exc

    prs = Presentation(str(template_path))
    target_layout = next(
        (layout for layout in prs.slide_layouts if str(layout.name or "").strip() == resolved_layout_name),
        None,
    )
    if target_layout is None:
        raise ValueError(f"Layout '{resolved_layout_name}' not found in theme '{resolved_theme}'.")

    slide = prs.slides.add_slide(target_layout)
    for placeholder in slide.placeholders:
        try:
            if placeholder.placeholder_format.type == 1:
                placeholder.text = resolved_layout_name
        except Exception:
            pass

    buf = io.BytesIO()
    prs.save(buf)
    fixed_bytes = _dedup_zip(buf.getvalue())

    with tempfile.TemporaryDirectory() as tmp:
        pptx_file = os.path.join(tmp, "slide.pptx")
        with open(pptx_file, "wb") as handle:
            handle.write(fixed_bytes)

        t0 = time.time()
        result = subprocess.run(
            [
                SOFFICE_BIN,
                "--headless",
                "--convert-to", "png",
                "--outdir", tmp,
                pptx_file,
            ],
            capture_output=True,
            text=True,
            timeout=60,
        )
        elapsed = time.time() - t0
        logger.info(
            "[layout-preview] soffice finished in %.1fs for %s/%s (rc=%d)",
            elapsed,
            resolved_theme,
            resolved_layout_name,
            result.returncode,
        )

        import glob

        pngs = sorted(glob.glob(os.path.join(tmp, "*.png")))
        if not pngs:
            raise RuntimeError(
                f"LibreOffice produced no PNG for {resolved_theme}/{resolved_layout_name}. "
                f"stderr: {result.stderr[:400]}"
            )

        src_png = pngs[-1]

        try:
            from PIL import Image

            img = Image.open(src_png).convert("RGB")
            img = img.resize((480, 270), Image.LANCZOS)
            cache.parent.mkdir(parents=True, exist_ok=True)
            img.save(str(cache), "PNG", optimize=True)
        except ImportError:
            cache.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(src_png, str(cache))

    return str(cache)


@router.get("/layout-preview")
@public_router.get("/layout-preview", include_in_schema=False)
def get_layout_preview(
    request: Request,
    theme: str = Query(..., description="Theme name, e.g. 'Business'"),
    layout: str = Query(..., description="Layout name, e.g. 'Section-1'"),
    user: dict = Depends(get_current_user),
):
    """Return a real screenshot of the requested PPTX layout slide as PNG."""
    try:
        _enforce_preview_rate_limit(request)
        cache_file = _render_layout_png(theme, layout)
        return FileResponse(
            cache_file,
            media_type="image/png",
            headers={
                "Cache-Control": "private, max-age=3600",
                "X-Layout-Preview-Theme": theme,
                "X-Layout-Preview-Layout": layout,
            },
        )
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc))
    except ValueError as exc:
        raise HTTPException(status_code=404, detail=str(exc))
    except RuntimeError as exc:
        logger.error("[layout-preview] RuntimeError: %s", exc)
        raise HTTPException(status_code=500, detail=str(exc))
    except Exception:
        logger.exception("[layout-preview] Unexpected error for %s/%s", theme, layout)
        raise HTTPException(status_code=500, detail="Layout preview generation failed.")
