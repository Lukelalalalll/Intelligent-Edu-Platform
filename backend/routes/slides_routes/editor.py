"""Editor routes — LLM layout assignment, session creation, export, image upload."""
from __future__ import annotations

import json
import logging
import os
import re
import uuid
import hashlib
from typing import Literal, Optional

from fastapi import HTTPException, Query, UploadFile, File
from fastapi.responses import FileResponse
from pydantic import BaseModel
from pptx import Presentation

from backend.config import Config
from backend.core.ai_provider import resolve_provider
from backend.services.ai_gateway_service import AIGatewayService
from backend.services.slides.output.editor_session import (
    create_session,
    get_session_meta,
    get_slide_png_path,
    export_pptx,
    re_render_session,
    SESSION_DIR,
)
from backend.services.slides.output.list_placeholders import PPTTemplateManager
from .router import slides_router, public_slides_router

logger = logging.getLogger(__name__)

# ── Schemas ───────────────────────────────────────────────────────────────────


class AutoAssignLayoutsRequest(BaseModel):
    provider: Optional[Literal["coze", "local_ollama"]] = "local_ollama"
    theme: str
    ppt_schema: dict


class RenderEditorSessionRequest(BaseModel):
    theme: str
    ppt_schema: dict


class ExportPptxRequest(BaseModel):
    session_id: str
    theme: str
    ppt_schema: dict
    edits: list[dict] = []
    slide_images: list[dict] = []


class ReRenderSessionRequest(BaseModel):
    session_id: str
    edits: list[dict] = []
    slide_images: list[dict] = []


# ── LLM Layout Assignment ────────────────────────────────────────────────────

_LAYOUT_SYSTEM_PROMPT = """\
You are a PPT design expert. Given the information for each slide and the available theme layouts, select the most suitable layout for each slide.

Rules:
- The first slide (and only the first) should typically use a "Title Slide" or the theme's cover/Title layout
- Section separator slides should use a Section-type layout
- **Important**: When bullet_count >= 1, you MUST choose a content layout that includes a body area (e.g. "Title and Content", "B1-P1-H", "Single Picture", etc.). Never use "Title Slide" or a title-only layout — content will be lost
- Use standard content layouts for 3-5 bullet points
- Use an Ending layout for the last slide (if available)
- Return a strict JSON array, no extra text"""


def _build_assignment_prompt(theme: str, layout_names: list[str], slides: list[dict]) -> str:
    slides_summary = []
    for s in slides:
        slides_summary.append({
            "index": s.get("index", 0),
            "title": s.get("title", ""),
            "bullet_count": len(s.get("content", [])),
        })
    return (
        f"Theme: {theme}\n"
        f"Available layouts: {json.dumps(layout_names, ensure_ascii=False)}\n\n"
        f"Slide information:\n{json.dumps(slides_summary, ensure_ascii=False, indent=2)}\n\n"
        '[Return format (JSON array)]:\n[{ "index": 0, "layout": "layout name" }, ...]'
    )


def _parse_llm_assignments(raw: str) -> list[dict]:
    """Extract JSON array from LLM response (may be wrapped in markdown fences)."""
    # Strip markdown code fences if present
    cleaned = re.sub(r"```(?:json)?\s*", "", raw).strip().rstrip("`")
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        # Try to find the first JSON array in the response
        match = re.search(r"\[.*\]", cleaned, re.DOTALL)
        if match:
            return json.loads(match.group(0))
        raise ValueError(f"Could not parse LLM response as JSON: {raw[:300]}")


@slides_router.post("/auto-assign-layouts")
async def auto_assign_layouts(req: AutoAssignLayoutsRequest):
    """Use LLM to auto-assign the best layout for each slide."""
    provider = resolve_provider(req.provider, feature="layout-assignment")

    # Get available layouts for the theme
    manager = PPTTemplateManager(Config.PPT_TEMPLATES_FOLDER)
    try:
        placeholders = manager.get_placeholders(req.theme)
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))

    layout_names = [p["name"] for p in placeholders]

    slides = req.ppt_schema.get("slides", [])
    # Add index if missing
    for i, s in enumerate(slides):
        s.setdefault("index", i)

    user_prompt = _build_assignment_prompt(req.theme, layout_names, slides)

    ai_service = AIGatewayService()
    try:
        raw_response = await ai_service.chat_with_provider(
            message=user_prompt,
            context={"system_override": _LAYOUT_SYSTEM_PROMPT},
            provider=provider,
        )
        assignments = _parse_llm_assignments(raw_response)
    except Exception as e:
        logger.error("[auto-assign] LLM failed: %s", e)
        raise HTTPException(status_code=502, detail=f"LLM layout assignment failed: {e}")

    # Apply assignments to schema
    assign_map = {a["index"]: a.get("layout", "") for a in assignments}
    for slide in slides:
        idx = slide.get("index", 0)
        if idx in assign_map:
            assigned = assign_map[idx]
            # Validate it's a real layout name
            if assigned in layout_names:
                slide["layout"] = assigned
            else:
                # Fallback: pick first content-style layout
                slide["layout"] = layout_names[0] if layout_names else ""

    updated_schema = {**req.ppt_schema, "slides": slides}

    # Ensure presentation has a title — generate one from AI if missing
    pres_title = updated_schema.get("presentation_title", "").strip()
    if not pres_title and slides:
        # Build a quick title from slide titles
        slide_titles = [s.get("title", "") for s in slides[:5] if s.get("title")]
        if slide_titles:
            pres_title = slide_titles[0]
        else:
            # Ask AI for a title based on content
            try:
                content_summary = "; ".join(
                    s.get("title", "") or (s.get("content", [""])[0] if s.get("content") else "")
                    for s in slides[:5]
                )
                title_response = await ai_service.chat_with_provider(
                    message=f"Based on the following presentation content, generate a concise title (output only the title, no extra text):\n{content_summary}",
                    context={"system_override": "You are a title generation assistant. Output only one concise title, no extra explanation."},
                    provider=provider,
                )
                pres_title = title_response.strip().strip('"\'')
            except Exception:
                pres_title = "Presentation"
        updated_schema["presentation_title"] = pres_title

    # Also ensure first slide has title text
    if slides and not slides[0].get("title"):
        slides[0]["title"] = pres_title

    return {"ppt_schema": updated_schema}


# ── Editor Session ────────────────────────────────────────────────────────────


@slides_router.post("/render-editor-session")
def render_editor_session(req: RenderEditorSessionRequest):
    """Generate PPTX, render every slide to PNG, and return element tree."""
    try:
        session = create_session(req.theme, req.ppt_schema)
        return session
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except RuntimeError as e:
        logger.error("[editor-session] %s", e)
        raise HTTPException(status_code=500, detail=str(e))


@slides_router.get("/editor-img/{session_id}/{page}.png")
def get_editor_slide_image(session_id: str, page: int):
    """Serve a rendered slide PNG for the editor canvas."""
    try:
        png_path = get_slide_png_path(session_id, page)
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="Slide image not found")
    return FileResponse(png_path, media_type="image/png", headers={"Cache-Control": "public, max-age=3600"})


@slides_router.post("/re-render-session")
def re_render_session_endpoint(req: ReRenderSessionRequest):
    """Apply text edits to a session's PPTX, re-render PNGs, return fresh meta."""
    try:
        meta = re_render_session(req.session_id, req.edits, req.slide_images)
        return meta
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except RuntimeError as e:
        logger.error("[re-render-session] %s", e)
        raise HTTPException(status_code=500, detail=str(e))


# ── Export ────────────────────────────────────────────────────────────────────


@slides_router.post("/export-pptx")
def export_pptx_endpoint(req: ExportPptxRequest):
    """Apply edits and return final .pptx for download."""
    try:
        out_path = export_pptx(req.session_id, req.theme, req.ppt_schema, req.edits, req.slide_images)
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))

    expected_slide_count = -1
    try:
        meta = get_session_meta(req.session_id)
        expected_slide_count = len(meta.get("slides", []) or [])
    except Exception:
        expected_slide_count = -1

    file_size = os.path.getsize(out_path)
    file_md5 = hashlib.md5(open(out_path, "rb").read()).hexdigest()
    try:
        slide_count = len(Presentation(out_path).slides)
    except Exception:
        slide_count = -1

    if expected_slide_count >= 0 and slide_count >= 0 and slide_count != expected_slide_count:
        logger.error(
            "[export-pptx] slide-count mismatch: session_id=%s expected=%d actual=%d file=%s req_theme=%s",
            req.session_id,
            expected_slide_count,
            slide_count,
            os.path.basename(out_path),
            req.theme,
        )
        raise HTTPException(
            status_code=500,
            detail=(
                "Export slide count mismatch "
                f"(expected={expected_slide_count}, actual={slide_count}). "
                "Please regenerate this session and retry."
            ),
        )

    logger.info(
        "[export-pptx] session_id=%s req_theme=%s file=%s size=%d md5=%s slides=%d edits=%d images=%d",
        req.session_id,
        req.theme,
        os.path.basename(out_path),
        file_size,
        file_md5,
        slide_count,
        len(req.edits or []),
        len(req.slide_images or []),
    )

    filename = f"{req.ppt_schema.get('presentation_title', 'presentation')}.pptx"
    return FileResponse(
        out_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
        headers={
            "Cache-Control": "no-store, no-cache, must-revalidate, max-age=0",
            "Pragma": "no-cache",
            "Expires": "0",
        },
    )


# ── Image upload ──────────────────────────────────────────────────────────────


@slides_router.post("/upload-image")
async def upload_slide_image(file: UploadFile = File(...)):
    """Upload an image for use in the slide editor."""
    assets_dir = os.path.join(SESSION_DIR, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    asset_id = uuid.uuid4().hex[:12]
    ext = os.path.splitext(file.filename or "img.png")[1] or ".png"
    save_path = os.path.join(assets_dir, f"{asset_id}{ext}")

    content = await file.read()
    if len(content) > 10 * 1024 * 1024:  # 10 MB limit
        raise HTTPException(status_code=413, detail="File too large (max 10 MB)")

    with open(save_path, "wb") as f:
        f.write(content)

    return {
        "asset_id": asset_id,
        "url": f"/api/slides/asset/{asset_id}{ext}",
    }


@slides_router.get("/asset/{filename}")
def serve_asset(filename: str):
    """Serve an uploaded image asset."""
    assets_dir = os.path.join(SESSION_DIR, "assets")
    file_path = os.path.join(assets_dir, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="Asset not found")
    return FileResponse(file_path, headers={"Cache-Control": "public, max-age=86400"})
