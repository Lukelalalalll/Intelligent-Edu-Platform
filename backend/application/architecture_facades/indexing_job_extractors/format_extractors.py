from __future__ import annotations

import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)


def _extract_pdf_markdown(source_path: Path, *, use_fast: bool) -> str:
    from backend.utils.pdf_extractor import extract_text_from_pdf

    return str(extract_text_from_pdf(source_path, use_fast=use_fast) or "").strip()


def _extract_fast_text(source_path: Path) -> str:
    suffix = source_path.suffix.lower()
    if suffix == ".docx":
        return _extract_text_from_docx(source_path)
    if suffix == ".pptx":
        return _extract_text_from_pptx(source_path)
    if suffix == ".xlsx":
        return _extract_text_from_xlsx(source_path)
    if suffix in {".md", ".markdown"}:
        return _strip_markdown(source_path.read_text(encoding="utf-8", errors="replace"))
    if suffix in {".html", ".htm"}:
        return _strip_html(source_path.read_text(encoding="utf-8", errors="replace"))
    if suffix in {".txt", ".text", ".csv"}:
        return source_path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        return _extract_pdf_markdown(source_path, use_fast=True)
    return source_path.read_text(encoding="utf-8", errors="replace")


def _extract_text_from_docx(source_path: Path) -> str:
    try:
        from docx import Document as DocxDocument
    except ImportError:
        logger.warning("python-docx not installed; cannot process %s", source_path.name)
        return source_path.read_text(encoding="utf-8", errors="replace")

    try:
        doc = DocxDocument(str(source_path))
    except Exception:
        logger.exception("Failed to open DOCX: %s", source_path.name)
        return ""

    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = (para.style.name or "").lower()
        if style_name.startswith("heading"):
            try:
                level = int(style_name.replace("heading", "").strip())
            except ValueError:
                level = 1
            parts.append(f"{'#' * max(1, min(level, 6))} {text}")
        else:
            parts.append(text)

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            header_sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
            rows.insert(1, header_sep)
            parts.append("\n".join(rows))

    return "\n\n".join(parts)


def _extract_text_from_pptx(source_path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed; cannot process %s", source_path.name)
        return ""

    try:
        prs = Presentation(str(source_path))
    except Exception:
        logger.exception("Failed to open PPTX: %s", source_path.name)
        return ""

    slides: list[str] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        lines = [f"# Slide {slide_index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text or "").strip()
                if text:
                    lines.append(text)
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def _extract_text_from_xlsx(source_path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed; cannot process %s", source_path.name)
        return ""

    try:
        wb = load_workbook(str(source_path), read_only=True, data_only=True)
    except Exception:
        logger.exception("Failed to open XLSX: %s", source_path.name)
        return ""

    sheets: list[str] = []
    for ws in wb.worksheets:
        rows: list[str] = [f"# Sheet {ws.title}"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell not in (None, "")]
            if cells:
                rows.append("| " + " | ".join(cell.replace("|", "\\|") for cell in cells) + " |")
        sheets.append("\n".join(rows))
    return "\n\n".join(sheets)


def _strip_markdown(text: str) -> str:
    text = re.sub(r"^(#{1,6})\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"__(.+?)__", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"_(.+?)_", r"\1", text)
    text = re.sub(r"~~(.+?)~~", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"^[-*_]{3,}\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"^(\s*)[*\-+]\s+", r"\1", text, flags=re.MULTILINE)
    text = re.sub(r"^(\s*)\d+\.\s+", r"\1", text, flags=re.MULTILINE)
    text = re.sub(r"^>\s?", "", text, flags=re.MULTILINE)
    return text


def _strip_html(text: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return re.sub(r"<[^>]+>", " ", text)

    soup = BeautifulSoup(text, "html.parser")
    return soup.get_text("\n")
