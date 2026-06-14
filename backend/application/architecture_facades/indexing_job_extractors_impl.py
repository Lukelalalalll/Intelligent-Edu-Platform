from __future__ import annotations

import json
import logging
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


@dataclass(slots=True)
class ParsedDocumentArtifact:
    kind: str
    filename: str
    content: str
    content_type: str


@dataclass(slots=True)
class ParsedDocumentResult:
    text: str
    normalized_markdown: str
    structure: dict[str, Any]
    quality_report: dict[str, Any]
    parser_used: str
    parser_strategy: str
    fallback_chain: list[str]
    artifacts: list[ParsedDocumentArtifact]


def extract_document_payload(
    source_path: Path,
    *,
    parser_strategy: str = "auto",
    index_profile: str = "quality",
    use_fast: bool = False,
) -> ParsedDocumentResult:
    """Extract a document into normalized markdown + structure + quality metadata."""
    strategy = str(parser_strategy or "auto").strip().lower() or "auto"
    profile = str(index_profile or "quality").strip().lower() or "quality"
    suffix = source_path.suffix.lower()
    fallback_chain: list[str] = []

    preferred_use_fast = use_fast or profile == "fast" or strategy == "fast"

    if strategy == "docling":
        result = _extract_with_docling(source_path)
        if result is None:
            raise RuntimeError(f"Docling extraction unavailable for {source_path.name}")
        return _finalize_parse_result(
            source_path,
            result["markdown"],
            parser_used="docling",
            parser_strategy=strategy,
            fallback_chain=fallback_chain,
            structure_hint=result.get("structure"),
        )

    if strategy == "marker":
        markdown = _extract_pdf_markdown(source_path, use_fast=False)
        if not markdown.strip():
            raise RuntimeError(f"marker extraction unavailable for {source_path.name}")
        return _finalize_parse_result(
            source_path,
            markdown,
            parser_used="marker",
            parser_strategy=strategy,
            fallback_chain=fallback_chain,
        )

    if strategy == "fast":
        markdown = _extract_fast_text(source_path)
        if not markdown.strip():
            raise RuntimeError(f"fast extraction unavailable for {source_path.name}")
        return _finalize_parse_result(
            source_path,
            markdown,
            parser_used="fast",
            parser_strategy=strategy,
            fallback_chain=fallback_chain,
        )

    if strategy == "auto":
        if not preferred_use_fast and _docling_enabled() and _supports_docling_suffix(suffix):
            docling_result = _extract_with_docling(source_path)
            if docling_result is not None:
                finalized = _finalize_parse_result(
                    source_path,
                    docling_result["markdown"],
                    parser_used="docling",
                    parser_strategy=strategy,
                    fallback_chain=fallback_chain,
                    structure_hint=docling_result.get("structure"),
                )
                if _passes_quality_gate(finalized.quality_report):
                    return finalized
                fallback_chain.append("docling")
                logger.info(
                    "Docling quality gate failed for %s; falling back",
                    source_path.name,
                )

        if suffix == ".pdf":
            markdown = _extract_pdf_markdown(source_path, use_fast=preferred_use_fast)
            parser_used = "fast" if preferred_use_fast else "marker"
            if markdown.strip():
                return _finalize_parse_result(
                    source_path,
                    markdown,
                    parser_used=parser_used,
                    parser_strategy=strategy,
                    fallback_chain=fallback_chain,
                )
            fallback_chain.append(parser_used)

        markdown = _extract_fast_text(source_path)
        if markdown.strip():
            return _finalize_parse_result(
                source_path,
                markdown,
                parser_used="fast",
                parser_strategy=strategy,
                fallback_chain=fallback_chain,
            )

    raise RuntimeError(f"Could not extract document payload from {source_path.name}")


def _docling_enabled() -> bool:
    try:
        from backend.config import Config

        return bool(getattr(Config, "RAG_ENABLE_DOCLING", True))
    except Exception:
        return True


def extract_text_from_path(source_path: Path, use_fast: bool = False) -> str:
    """Backward-compatible wrapper used by older callers."""
    result = extract_document_payload(
        source_path,
        parser_strategy="fast" if use_fast else "auto",
        index_profile="fast" if use_fast else "quality",
        use_fast=use_fast,
    )
    return result.text


def _extract_with_docling(source_path: Path) -> dict[str, Any] | None:
    try:
        from docling.document_converter import DocumentConverter
    except ImportError:
        logger.info("Docling not installed; skipping for %s", source_path.name)
        return None

    try:
        converter = DocumentConverter()
        conversion = converter.convert(str(source_path))
        document = conversion.document
        markdown = str(document.export_to_markdown() or "").strip()
        if not markdown:
            return None
        structure = _build_structure_from_markdown(markdown, source_path.name)
        structure["docling_json"] = _safe_docling_json(document)
        return {"markdown": markdown, "structure": structure}
    except Exception:
        logger.exception("Docling extraction failed for %s", source_path.name)
        return None


def _safe_docling_json(document: Any) -> dict[str, Any]:
    try:
        if hasattr(document, "export_to_dict"):
            payload = document.export_to_dict()
            if isinstance(payload, dict):
                return payload
        if hasattr(document, "model_dump"):
            payload = document.model_dump()
            if isinstance(payload, dict):
                return payload
    except Exception:
        logger.debug("Could not serialize Docling document", exc_info=True)
    return {}


def _extract_pdf_markdown(source_path: Path, *, use_fast: bool) -> str:
    from backend.utils.pdf_extractor import extract_text_from_pdf

    return str(extract_text_from_pdf(source_path, use_fast=use_fast) or "").strip()


def _extract_fast_text(source_path: Path) -> str:
    suffix = source_path.suffix.lower()
    if suffix == ".docx":
        return _extract_text_from_docx(source_path)
    if suffix == ".pptx":
        return _extract_text_from_pptx(source_path)
    if suffix == ".xlsx":
        return _extract_text_from_xlsx(source_path)
    if suffix in {".md", ".markdown"}:
        return _strip_markdown(source_path.read_text(encoding="utf-8", errors="replace"))
    if suffix in {".html", ".htm"}:
        return _strip_html(source_path.read_text(encoding="utf-8", errors="replace"))
    if suffix in {".txt", ".text", ".csv"}:
        return source_path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        return _extract_pdf_markdown(source_path, use_fast=True)
    return source_path.read_text(encoding="utf-8", errors="replace")


def _extract_text_from_docx(source_path: Path) -> str:
    try:
        from docx import Document as DocxDocument
    except ImportError:
        logger.warning("python-docx not installed; cannot process %s", source_path.name)
        return source_path.read_text(encoding="utf-8", errors="replace")

    try:
        doc = DocxDocument(str(source_path))
    except Exception:
        logger.exception("Failed to open DOCX: %s", source_path.name)
        return ""

    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = (para.style.name or "").lower()
        if style_name.startswith("heading"):
            try:
                level = int(style_name.replace("heading", "").strip())
            except ValueError:
                level = 1
            parts.append(f"{'#' * max(1, min(level, 6))} {text}")
        else:
            parts.append(text)

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            header_sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
            rows.insert(1, header_sep)
            parts.append("\n".join(rows))

    return "\n\n".join(parts)


def _extract_text_from_pptx(source_path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed; cannot process %s", source_path.name)
        return ""

    try:
        prs = Presentation(str(source_path))
    except Exception:
        logger.exception("Failed to open PPTX: %s", source_path.name)
        return ""

    slides: list[str] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        lines = [f"# Slide {slide_index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text or "").strip()
                if text:
                    lines.append(text)
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def _extract_text_from_xlsx(source_path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed; cannot process %s", source_path.name)
        return ""

    try:
        wb = load_workbook(str(source_path), read_only=True, data_only=True)
    except Exception:
        logger.exception("Failed to open XLSX: %s", source_path.name)
        return ""

    sheets: list[str] = []
    for ws in wb.worksheets:
        rows: list[str] = [f"# Sheet {ws.title}"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell not in (None, "")]
            if cells:
                rows.append("| " + " | ".join(cell.replace("|", "\\|") for cell in cells) + " |")
        sheets.append("\n".join(rows))
    return "\n\n".join(sheets)


def _strip_markdown(text: str) -> str:
    text = re.sub(r"^(#{1,6})\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"__(.+?)__", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"_(.+?)_", r"\1", text)
    text = re.sub(r"~~(.+?)~~", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"^[-*_]{3,}\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"^(\s*)[*\-+]\s+", r"\1", text, flags=re.MULTILINE)
    text = re.sub(r"^(\s*)\d+\.\s+", r"\1", text, flags=re.MULTILINE)
    text = re.sub(r"^>\s?", "", text, flags=re.MULTILINE)
    return text


def _strip_html(text: str) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return re.sub(r"<[^>]+>", " ", text)

    soup = BeautifulSoup(text, "html.parser")
    return soup.get_text("\n")


def _finalize_parse_result(
    source_path: Path,
    markdown: str,
    *,
    parser_used: str,
    parser_strategy: str,
    fallback_chain: list[str],
    structure_hint: dict[str, Any] | None = None,
) -> ParsedDocumentResult:
    normalized = _normalize_markdown(markdown)
    structure = structure_hint or _build_structure_from_markdown(normalized, source_path.name)
    quality_report = _build_quality_report(normalized, structure)
    artifacts = [
        ParsedDocumentArtifact(
            kind="normalized_markdown",
            filename=f"{source_path.stem}.normalized.md",
            content=normalized,
            content_type="text/markdown",
        ),
        ParsedDocumentArtifact(
            kind="structure_json",
            filename=f"{source_path.stem}.structure.json",
            content=json.dumps(structure, ensure_ascii=False, indent=2),
            content_type="application/json",
        ),
        ParsedDocumentArtifact(
            kind="quality_report_json",
            filename=f"{source_path.stem}.quality_report.json",
            content=json.dumps(quality_report, ensure_ascii=False, indent=2),
            content_type="application/json",
        ),
    ]
    return ParsedDocumentResult(
        text=normalized,
        normalized_markdown=normalized,
        structure=structure,
        quality_report=quality_report,
        parser_used=parser_used,
        parser_strategy=parser_strategy,
        fallback_chain=list(fallback_chain),
        artifacts=artifacts,
    )


def _normalize_markdown(markdown: str) -> str:
    text = str(markdown or "").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _build_structure_from_markdown(markdown: str, doc_name: str) -> dict[str, Any]:
    lines = markdown.splitlines()
    sections: list[dict[str, Any]] = []
    blocks: list[dict[str, Any]] = []
    stack: list[tuple[int, str]] = []
    current_section = {
        "title": "Document",
        "heading_level": 0,
        "heading_path": "Document",
        "page_start": 1,
        "page_end": 1,
        "text": "",
    }
    current_lines: list[str] = []
    current_block_lines: list[str] = []
    block_type = "paragraph"
    page_num = 1

    def flush_block() -> None:
        nonlocal current_block_lines, block_type
        content = "\n".join(current_block_lines).strip()
        if not content:
            current_block_lines = []
            block_type = "paragraph"
            return
        blocks.append(
            {
                "element_type": block_type,
                "heading_path": current_section["heading_path"],
                "page_start": page_num,
                "page_end": page_num,
                "text": content,
            }
        )
        current_block_lines = []
        block_type = "paragraph"

    def flush_section() -> None:
        nonlocal current_section, current_lines
        content = "\n".join(current_lines).strip()
        if content or current_section["title"] != "Document":
            section = dict(current_section)
            section["text"] = content
            sections.append(section)
        current_lines = []

    heading_re = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
    for raw_line in lines:
        line = raw_line.rstrip()
        if "\f" in line:
            page_num += line.count("\f")
            line = line.replace("\f", "").rstrip()
        heading_match = heading_re.match(line)
        if heading_match:
            flush_block()
            flush_section()
            level = len(heading_match.group(1))
            title = heading_match.group(2).strip()
            while stack and stack[-1][0] >= level:
                stack.pop()
            stack.append((level, title))
            current_section = {
                "title": title,
                "heading_level": level,
                "heading_path": " > ".join(item[1] for item in stack),
                "page_start": page_num,
                "page_end": page_num,
                "text": "",
            }
            continue

        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            if block_type != "table":
                flush_block()
                block_type = "table"
            current_block_lines.append(line)
        elif re.match(r"^[-*+]\s+", stripped) or re.match(r"^\d+\.\s+", stripped):
            if block_type != "list":
                flush_block()
                block_type = "list"
            current_block_lines.append(line)
        elif stripped:
            if block_type not in {"paragraph", "formula"}:
                flush_block()
            if _looks_like_formula(stripped):
                if block_type != "formula":
                    flush_block()
                    block_type = "formula"
            else:
                if block_type != "paragraph":
                    flush_block()
                    block_type = "paragraph"
            current_block_lines.append(line)
        else:
            flush_block()

        current_lines.append(line)
        current_section["page_end"] = page_num

    flush_block()
    flush_section()
    return {
        "doc_name": doc_name,
        "sections": sections,
        "blocks": blocks,
    }


def _build_quality_report(markdown: str, structure: dict[str, Any]) -> dict[str, Any]:
    blocks = list(structure.get("blocks") or [])
    sections = list(structure.get("sections") or [])
    pages = [
        int(block.get("page_start") or 1)
        for block in blocks
        if int(block.get("page_start") or 1) > 0
    ]
    page_count = max(pages, default=1)
    page_blocks: dict[int, int] = {page: 0 for page in range(1, page_count + 1)}
    for block in blocks:
        page_blocks[int(block.get("page_start") or 1)] = page_blocks.get(int(block.get("page_start") or 1), 0) + 1

    empty_pages = sum(1 for value in page_blocks.values() if value == 0)
    heading_count = sum(1 for section in sections if int(section.get("heading_level") or 0) > 0)
    table_count = sum(1 for block in blocks if str(block.get("element_type")) == "table")
    formula_count = sum(1 for block in blocks if str(block.get("element_type")) == "formula")
    total_chars = len(markdown)
    formula_density = round(formula_count / max(len(blocks), 1), 4)
    avg_block_len = round(sum(len(str(block.get("text") or "")) for block in blocks) / max(len(blocks), 1), 2)

    replacement_chars = markdown.count("\ufffd")
    mojibake_hits = len(re.findall(r"[ÃÂÐÑØÞ]{2,}|鈭|锛|銆|绗", markdown))
    garbled_text_ratio = round((replacement_chars + mojibake_hits) / max(total_chars, 1), 4)

    ocr_like_hits = len(re.findall(r"\b(?:ocr|scan|scanned)\b", markdown, flags=re.IGNORECASE))
    short_lines = sum(1 for line in markdown.splitlines() if 0 < len(line.strip()) <= 2)
    ocr_ratio = round(min(1.0, (ocr_like_hits + short_lines) / max(len(markdown.splitlines()), 1)), 4)
    quality_status = "ok"
    if total_chars < 120 or avg_block_len < 20 or garbled_text_ratio > 0.08:
        quality_status = "needs_review"
    if total_chars < 40:
        quality_status = "poor"

    return {
        "page_count": page_count,
        "empty_page_ratio": round(empty_pages / max(page_count, 1), 4),
        "ocr_ratio": ocr_ratio,
        "heading_count": heading_count,
        "table_count": table_count,
        "formula_density": formula_density,
        "garbled_text_ratio": garbled_text_ratio,
        "avg_block_len": avg_block_len,
        "quality_status": quality_status,
        "block_count": len(blocks),
        "section_count": len(sections),
    }


def _passes_quality_gate(quality_report: dict[str, Any]) -> bool:
    return (
        str(quality_report.get("quality_status") or "ok") == "ok"
        and float(quality_report.get("garbled_text_ratio") or 0.0) <= 0.08
        and float(quality_report.get("empty_page_ratio") or 0.0) <= 0.6
        and int(quality_report.get("block_count") or 0) > 0
    )


def _looks_like_formula(text: str) -> bool:
    return bool(re.search(r"[=<>]|\\[a-zA-Z]+|\$\$?|∑|∫|√|±", text))


def _supports_docling_suffix(suffix: str) -> bool:
    return suffix in {
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".html",
        ".htm",
        ".md",
        ".markdown",
        ".txt",
        ".png",
        ".jpg",
        ".jpeg",
    }
