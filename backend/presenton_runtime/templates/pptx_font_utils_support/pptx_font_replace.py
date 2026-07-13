from __future__ import annotations

import os
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from typing import Dict, Optional, Sequence

from pptx import Presentation

from .common import FONT_TAGS, PPT_NS, TEXT_STYLE_TAGS
from .font_matching import _font_style_variant, normalize_font_family_name


def create_font_alias_config(
    raw_fonts: list[str],
    extra_includes: Optional[list[str]] = None,
    temp_dir: Optional[str] = None,
    explicit_aliases: Optional[Dict[str, str]] = None,
    protected_font_names: Optional[Sequence[str]] = None,
) -> str:
    mappings = {}
    explicit_aliases = {src: dst for src, dst in (explicit_aliases or {}).items() if src and dst and src != dst}
    protected_names = {name for name in (protected_font_names or []) if name}
    skip_normalization = protected_names.union(explicit_aliases.keys()).union(explicit_aliases.values())
    for font_name in raw_fonts:
        if font_name in skip_normalization:
            continue
        normalized = normalize_font_family_name(font_name)
        if normalized and normalized != font_name:
            mappings[font_name] = normalized
    fd, fonts_conf_path = tempfile.mkstemp(prefix="fonts_alias_", suffix=".conf", dir=temp_dir)
    os.close(fd)
    with open(fonts_conf_path, "w", encoding="utf-8") as cfg:
        cfg.write("<?xml version='1.0'?>\n<!DOCTYPE fontconfig SYSTEM \"urn:fontconfig:fonts.dtd\">\n<fontconfig>\n  <include>/etc/fonts/fonts.conf</include>\n")
        for include_path in extra_includes or []:
            if include_path:
                cfg.write(f"  <include>{include_path}</include>\n")
        for source, target in {**mappings, **explicit_aliases}.items():
            cfg.write(f"\n  <match target=\"pattern\">\n    <test name=\"family\" compare=\"eq\">\n      <string>{source}</string>\n    </test>\n    <edit name=\"family\" mode=\"assign\" binding=\"strong\">\n      <string>{target}</string>\n    </edit>\n  </match>\n")
        cfg.write("\n</fontconfig>\n")
    return fonts_conf_path


def _replace_fonts_in_xml_root(
    root: ET.Element,
    font_mapping: Dict[str, str],
    font_variant_mapping: Optional[Dict[str, Dict[str, str]]] = None,
) -> bool:
    def first_typeface(style_elem: Optional[ET.Element]) -> Optional[str]:
        if style_elem is None:
            return None
        for font_tag in FONT_TAGS:
            font_elem = style_elem.find(font_tag, PPT_NS)
            if font_elem is not None and font_elem.get("typeface"):
                return font_elem.get("typeface")
        return None

    changed = False
    for style_tag in TEXT_STYLE_TAGS:
        for style_elem in root.findall(f".//{style_tag}", PPT_NS):
            for font_tag in FONT_TAGS:
                font_elem = style_elem.find(font_tag, PPT_NS)
                typeface = font_elem.get("typeface") if font_elem is not None else None
                if not typeface:
                    continue
                replacement = ((font_variant_mapping or {}).get(typeface) or {}).get(_font_style_variant(typeface, style_elem, []))
                replacement = replacement or font_mapping.get(typeface)
                if replacement and replacement != typeface:
                    font_elem.set("typeface", replacement)
                    changed = True

    run_tags = {f"{{{PPT_NS['a']}}}r", f"{{{PPT_NS['a']}}}fld"}
    for paragraph in root.findall(".//a:p", PPT_NS):
        p_pr = paragraph.find("a:pPr", PPT_NS)
        paragraph_default = p_pr.find("a:defRPr", PPT_NS) if p_pr is not None else None
        inherited_typeface = first_typeface(paragraph_default)
        if not inherited_typeface:
            continue
        variant_mapping = (font_variant_mapping or {}).get(inherited_typeface)
        if not variant_mapping:
            original_typeface = next((source for source, replacement in font_mapping.items() if replacement == inherited_typeface), None)
            variant_mapping = (font_variant_mapping or {}).get(original_typeface or "")
            inherited_typeface = original_typeface or inherited_typeface
        if not variant_mapping:
            continue
        for child in paragraph:
            if child.tag not in run_tags:
                continue
            r_pr = child.find("a:rPr", PPT_NS)
            if r_pr is None:
                r_pr = ET.Element(f"{{{PPT_NS['a']}}}rPr")
                child.insert(0, r_pr)
            if first_typeface(r_pr):
                continue
            replacement = variant_mapping.get(_font_style_variant(inherited_typeface, r_pr, [paragraph_default]))
            if replacement:
                latin = ET.SubElement(r_pr, f"{{{PPT_NS['a']}}}latin")
                latin.set("typeface", replacement)
                changed = True
    return changed


def _replace_fonts_in_pptx_xml(
    pptx_path: str,
    font_mapping: Dict[str, str],
    output_path: str,
    font_variant_mapping: Optional[Dict[str, Dict[str, str]]] = None,
) -> None:
    xml_prefixes = ("ppt/slides/", "ppt/slideLayouts/", "ppt/slideMasters/", "ppt/charts/")
    with zipfile.ZipFile(pptx_path, "r") as src, zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as dst:
        for info in src.infolist():
            data = src.read(info.filename)
            if info.filename.endswith(".xml") and info.filename.startswith(xml_prefixes):
                try:
                    root = ET.fromstring(data)
                    if _replace_fonts_in_xml_root(root, font_mapping, font_variant_mapping):
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                except Exception:
                    pass
            dst.writestr(info, data)


def replace_fonts_in_pptx(
    pptx_path: str,
    font_mapping: Dict[str, str],
    output_path: str,
    font_variant_mapping: Optional[Dict[str, Dict[str, str]]] = None,
) -> None:
    if font_variant_mapping or font_mapping:
        _replace_fonts_in_pptx_xml(pptx_path, font_mapping, output_path, font_variant_mapping)
        return
    prs = Presentation(pptx_path)
    for slides in (prs.slides, prs.slide_layouts, prs.slide_masters):
        for slide in slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name and run.font.name in font_mapping:
                                run.font.name = font_mapping[run.font.name]
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.font.name and run.font.name in font_mapping:
                                        run.font.name = font_mapping[run.font.name]
    prs.save(output_path)
