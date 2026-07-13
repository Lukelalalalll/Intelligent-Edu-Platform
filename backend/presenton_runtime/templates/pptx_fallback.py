from __future__ import annotations

import html
from typing import Iterable, List

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ModuleNotFoundError:  # pragma: no cover - dependency is required in runtime/tests
    Presentation = None  # type: ignore[assignment]
    MSO_SHAPE_TYPE = None  # type: ignore[assignment]

try:
    from services.slides.output.editor_session.rendering import render_slides_via_pillow
except ModuleNotFoundError:  # pragma: no cover - backend test import path
    from backend.services.slides.output.editor_session.rendering import (
        render_slides_via_pillow,
    )


FALLBACK_PREVIEW_WIDTH = 1280
FALLBACK_PREVIEW_HEIGHT = 720
EMU_PER_INCH = 914400


def _require_presentation():
    if Presentation is None:
        raise RuntimeError("python-pptx is required for PPTX fallback rendering")
    return Presentation


def _load_presentation(pptx_path: str):
    presentation_cls = _require_presentation()
    try:
        return presentation_cls(pptx_path)
    except Exception as exc:  # pragma: no cover - library exception shape is not stable
        raise RuntimeError(f"Unable to parse PPTX for fallback rendering: {exc}") from exc


def render_fallback_slide_pngs_from_pptx(
    pptx_path: str,
    *,
    max_slides: int | None = None,
) -> List[bytes]:
    try:
        with open(pptx_path, "rb") as pptx_file:
            pptx_bytes = pptx_file.read()
    except OSError as exc:
        raise RuntimeError(f"Unable to read PPTX for fallback previews: {exc}") from exc

    try:
        slide_pngs = render_slides_via_pillow(pptx_bytes)
    except Exception as exc:  # pragma: no cover - library exception shape is not stable
        raise RuntimeError(f"Unable to render PPTX fallback previews: {exc}") from exc

    return slide_pngs[:max_slides] if max_slides else slide_pngs


def extract_slide_htmls_from_pptx(
    pptx_path: str,
    *,
    max_slides: int | None = None,
    width: int = FALLBACK_PREVIEW_WIDTH,
    height: int = FALLBACK_PREVIEW_HEIGHT,
) -> List[str]:
    presentation = _load_presentation(pptx_path)
    slides = list(presentation.slides)
    if max_slides:
        slides = slides[:max_slides]

    slide_width = int(getattr(presentation, "slide_width", 0) or 0)
    slide_height = int(getattr(presentation, "slide_height", 0) or 0)
    if slide_width <= 0:
        slide_width = 10 * EMU_PER_INCH
    if slide_height <= 0:
        slide_height = 7_500_000

    return [
        _build_fallback_slide_html(
            slide,
            slide_width_emu=slide_width,
            slide_height_emu=slide_height,
            width=width,
            height=height,
        )
        for slide in slides
    ]


def _build_fallback_slide_html(
    slide,
    *,
    slide_width_emu: int,
    slide_height_emu: int,
    width: int,
    height: int,
) -> str:
    blocks = sorted(
        _iter_slide_blocks(
            slide,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            width=width,
            height=height,
        ),
        key=lambda block: (block["top"], block["left"]),
    )

    html_parts = [
        (
            '<div class="slide-container">'
            f'<div class="slide-content fallback-slide" style="position:relative;width:{width}px;'
            f'height:{height}px;background:#ffffff;overflow:hidden;">'
        )
    ]

    if not blocks:
        html_parts.append(
            '<div style="position:absolute;inset:0;display:flex;align-items:center;'
            'justify-content:center;color:#9ca3af;font-size:24px;">(Empty slide)</div>'
        )

    for block in blocks:
        style = (
            f"position:absolute;left:{block['left']}px;top:{block['top']}px;"
            f"width:{block['width']}px;height:{block['height']}px;"
        )
        if block["kind"] == "text":
            html_parts.append(
                f'<div class="fallback-text" style="{style}'
                "padding:4px 6px;color:#1f2937;font-size:18px;line-height:1.35;"
                'white-space:pre-wrap;overflow:hidden;">'
                f"{_escape_with_breaks(block['text'])}</div>"
            )
        elif block["kind"] == "table":
            html_parts.append(
                f'<div class="fallback-table" style="{style}'
                'padding:8px;border:1px solid #d1d5db;background:#f9fafb;overflow:hidden;">'
                '<div style="font-size:12px;color:#6b7280;margin-bottom:6px;">Table</div>'
                f"{block['table_html']}</div>"
            )
        else:
            label = html.escape(block["label"])
            html_parts.append(
                f'<div class="fallback-placeholder" style="{style}'
                'display:flex;align-items:center;justify-content:center;'
                'border:1px dashed #cbd5e1;background:#f8fafc;color:#64748b;'
                'font-size:16px;text-transform:uppercase;">'
                f"{label}</div>"
            )

    html_parts.append("</div></div>")
    return "".join(html_parts)


def _iter_slide_blocks(
    slide,
    *,
    slide_width_emu: int,
    slide_height_emu: int,
    width: int,
    height: int,
) -> Iterable[dict]:
    for shape in slide.shapes:
        bounds = _shape_bounds(
            shape,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            width=width,
            height=height,
        )
        if not bounds:
            continue

        left, top, block_width, block_height = bounds
        if getattr(shape, "has_text_frame", False):
            text = _shape_text(shape)
            if text:
                yield {
                    "kind": "text",
                    "left": left,
                    "top": top,
                    "width": block_width,
                    "height": block_height,
                    "text": text,
                }
                continue

        if getattr(shape, "has_table", False):
            table_html = _shape_table_html(shape)
            if table_html:
                yield {
                    "kind": "table",
                    "left": left,
                    "top": top,
                    "width": block_width,
                    "height": block_height,
                    "table_html": table_html,
                }
                continue

        label = _shape_placeholder_label(shape)
        if label:
            yield {
                "kind": "placeholder",
                "left": left,
                "top": top,
                "width": block_width,
                "height": block_height,
                "label": label,
            }


def _shape_bounds(
    shape,
    *,
    slide_width_emu: int,
    slide_height_emu: int,
    width: int,
    height: int,
) -> tuple[int, int, int, int] | None:
    try:
        left = int((int(getattr(shape, "left", 0) or 0) / slide_width_emu) * width)
        top = int((int(getattr(shape, "top", 0) or 0) / slide_height_emu) * height)
        block_width = int((int(getattr(shape, "width", 0) or 0) / slide_width_emu) * width)
        block_height = int((int(getattr(shape, "height", 0) or 0) / slide_height_emu) * height)
    except Exception:
        return None

    block_width = max(48, min(block_width or 180, width))
    block_height = max(28, min(block_height or 48, height))
    left = max(0, min(left, max(0, width - block_width)))
    top = max(0, min(top, max(0, height - block_height)))
    return left, top, block_width, block_height


def _shape_text(shape) -> str:
    lines: List[str] = []
    for paragraph in getattr(shape.text_frame, "paragraphs", []):
        text = (getattr(paragraph, "text", "") or "").strip()
        if text:
            lines.append(text)
    return "\n".join(lines)


def _shape_table_html(shape) -> str:
    rows_html: List[str] = []
    for row in getattr(shape.table, "rows", []):
        cells_html: List[str] = []
        for cell in getattr(row, "cells", []):
            cell_text = _table_cell_text(cell)
            cells_html.append(
                '<td style="border:1px solid #d1d5db;padding:4px 6px;vertical-align:top;">'
                f"{_escape_with_breaks(cell_text)}</td>"
            )
        rows_html.append(f"<tr>{''.join(cells_html)}</tr>")
    if not rows_html:
        return ""
    return (
        '<table style="width:100%;border-collapse:collapse;font-size:12px;color:#334155;">'
        f"{''.join(rows_html)}</table>"
    )


def _table_cell_text(cell) -> str:
    lines: List[str] = []
    for paragraph in getattr(cell.text_frame, "paragraphs", []):
        text = (getattr(paragraph, "text", "") or "").strip()
        if text:
            lines.append(text)
    return "\n".join(lines)


def _shape_placeholder_label(shape) -> str | None:
    shape_type = getattr(shape, "shape_type", None)
    if MSO_SHAPE_TYPE is not None:
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "Image"
        if shape_type == MSO_SHAPE_TYPE.CHART:
            return "Chart"
        if shape_type == MSO_SHAPE_TYPE.MEDIA:
            return "Media"
    name = (getattr(shape, "name", "") or "").strip()
    if not name:
        return None
    lowered = name.lower()
    if "picture" in lowered or "image" in lowered:
        return "Image"
    if "chart" in lowered:
        return "Chart"
    return None


def _escape_with_breaks(value: str) -> str:
    return html.escape(value).replace("\n", "<br/>")
